"""
Host dashboard routes for Family Feud.

Owns: Host dashboard, round management, code management, settings,
broadcast, reset, training data, and all related host-facing endpoints.
"""

import os
import json
import base64
import time
import html
import urllib.request
import urllib.error
from datetime import datetime
from flask import Blueprint, request, render_template, redirect, url_for, jsonify, session, flash

from config import (
    logger, reset_state,
    BASE_DIR,
    AI_SCORING_ENABLED, AI_MODEL_CHOICES,
    GITHUB_TOKEN, GITHUB_REPO,
    QR_DEFAULT_URL, CORRECTIONS_FILE,
    time_ago, format_timestamp,
)
from auth import host_required
from database import (
    db_connect,
    load_fixed_codes,
    ensure_fixed_codes,
    get_setting,
    set_setting,
)
from ai import (
    load_corrections_history,
    get_current_ai_model,
)

host_bp = Blueprint('host', __name__)

# Game configuration - 8 rounds
ROUNDS_CONFIG = [
    {"round": 1, "answers": 4},
    {"round": 2, "answers": 5},
    {"round": 3, "answers": 6},
    {"round": 4, "answers": 4},
    {"round": 5, "answers": 5},
    {"round": 6, "answers": 3},
    {"round": 7, "answers": 5},
    {"round": 8, "answers": 4}
]

# Pre-built surveys for quick round creation via dropdown
PREBUILT_SURVEYS = {
    "survey1": {
        "name": "Survey 1",
        "rounds": [
            {"question": "Name Something Parents Warn Their Children Not To Get Their Fingers Caught In", "answers": ["Door", "Fan", "Outlet", "Cookie Jar"], "answer1_count": 45},
            {"question": "A Young Person \u201cFights For Their Right To Party.\u201d What Might An Old Person Fight For The Right To Do?", "answers": ["Sleep", "Vote", "Retire", "Keep License", "Get Social Security"], "answer1_count": 32},
            {"question": "Name Something That Goes Well With Pizza", "answers": ["Beer", "Soda", "Salad", "Breadstick/Knots", "Chicken Wings", "Chips"], "answer1_count": 36},
            {"question": "Name Something From The Laundry That\u2019s Impossible To Fold Neatly", "answers": ["Fitted Sheets", "Socks", "Underwear", "Blouse"], "answer1_count": 35},
            {"question": "Name A Place Where You\u2019d Be Mortified If Your Cell Phone Went Off", "answers": ["Church", "Funeral", "Movie Theater", "Job Interview", "Wedding"], "answer1_count": 39},
            {"question": "Name Something You Should Switch Off Before Going To Bed", "answers": ["Lights", "Phone", "TV"], "answer1_count": 67},
            {"question": "Name A Common Sickness That Kids Seem To Get More Often Than Adults", "answers": ["Cold", "Flu", "Chicken Pox", "Ear Infection", "Strep Throat"], "answer1_count": 32},
            {"question": "Name A Reason Why A Man Would Wax Hair Off Part Of His Body", "answers": ["Too Hairy", "For Spouse/Date", "Body Builder", "Swimmer"], "answer1_count": 34},
        ]
    },
    "survey2": {
        "name": "Survey 2",
        "rounds": [
            {"question": "Name Something Permanent On a Criminal\u2019s Skin That Police Use To Be Sure They\u2019ve Got Their Man", "answers": ["Tattoo", "Fingerprint", "Scar", "Birthmark"], "answer1_count": 41},
            {"question": "What Might Someone Use While Cutting Their Own Hair?", "answers": ["Scissors", "Mirror", "Clippers", "Comb", "Bowl"], "answer1_count": 48},
            {"question": "Name Something Babies And Puppies Have In Common", "answers": ["Cute", "Drooling", "Need Attention", "Playful", "Sleep A lot", "Cry"], "answer1_count": 34},
            {"question": "What Diaper Bag Item Would A Parent Hate To Be Without?", "answers": ["Diapers", "Wipes", "Bottle", "Pacifier"], "answer1_count": 49},
            {"question": "Name Something Twins Might Always Share", "answers": ["Looks", "Parents", "Genes", "Birthday", "Last Name"], "answer1_count": 40},
            {"question": "Something Specific People Do In Front Of Mirror", "answers": ["Apply Makeup", "Check Outfit", "Pose Naked"], "answer1_count": 63},
            {"question": "Name A Type Of Sauce That You\u2019d Never Put On Pasta", "answers": ["Apple Sauce", "Hot Sauce", "Ketchup", "BBQ", "Chocolate"], "answer1_count": 30},
            {"question": "Name Something A Child Does To Prove They\u2019re Too Sick For School", "answers": ["Cough", "Vomit", "Cry", "Take Temperature"], "answer1_count": 50},
        ]
    },
    "survey3": {
        "name": "Survey 3",
        "rounds": [
            {"question": "Name Something You Might Adjust When You Get Into A Rental Car", "answers": ["Seat", "Mirrors", "Seat Belt", "Steering Wheel"], "answer1_count": 58},
            {"question": "Name Something A Woman Should Know A Man Before Marrying Him", "answers": ["Income", "Age", "Does He Have Kids", "His Name", "Past Relationships"], "answer1_count": 39},
            {"question": "Name Something You Need In Order To Make A Garden", "answers": ["Seeds", "Soil", "Water", "Hoe", "Shovel", "Plot of Land"], "answer1_count": 35},
            {"question": "Name A Place Where You Hear People Being Paged Over A Loudspeaker", "answers": ["Hospital", "Airport", "School", "Store"], "answer1_count": 31},
            {"question": "Tell Me Something You Do When You Stay Up Late At Night", "answers": ["Watch TV/Movie", "Read", "Snack", "Drink", "Play Phone/Video Games"], "answer1_count": 58},
            {"question": "Name A Crime That Some People Probably Commit Every Day", "answers": ["Speeding", "Jaywalking", "Littering"], "answer1_count": 62},
            {"question": "Name A Reason Why A Person Might Prefer To Own A Dog Over A Cat", "answers": ["Protection", "Loyalty", "Cat Allergies", "Friendlier", "More fun to play with"], "answer1_count": 54},
            {"question": "Name A Phrase You\u2019d Say To Your Partner That Starts With \u201cYou Drive Me __.\u201d", "answers": ["Crazy/Nuts", "Wild", "Up a Wall", "To Drink"], "answer1_count": 58},
        ]
    },
    "survey4": {
        "name": "Survey 4",
        "rounds": [
            {"question": "We Asked 100 Women: Name A Gift That You\u2019d Always Be Happy To Get From Your Partner", "answers": ["Flowers", "Jewelry", "Money", "Chocolate"], "answer1_count": 43},
            {"question": "Name A Slow-Moving Vehicle That You Hate To Get Stuck Behind", "answers": ["Bus", "Semi-Truck", "Tractor", "Garbage Truck", "Dump Truck"], "answer1_count": 34},
            {"question": "Name A Last Minute Problem That Could Make You Late For Work", "answers": ["Traffic", "Car Trouble", "Lost Keys", "Child is Sick", "No Gas", "Bad Hair"], "answer1_count": 35},
            {"question": "Name Something Parents Warn Their Children Not To Get Their Fingers Caught In", "answers": ["Door", "Fan", "Outlet", "Cookie Jar"], "answer1_count": 45},
            {"question": "Name Something You Spray On Yourself That Would Sting If It Got In Your Eyes", "answers": ["Perfume", "Insect Repellent", "Hairspray", "Sunscreen/Tan", "Deodorant"], "answer1_count": 34},
            {"question": "Name Something You Dunk", "answers": ["Basketball", "Donuts", "Cookies"], "answer1_count": 59},
            {"question": "Name Something You Wear That Covers Your Ears", "answers": ["Earmuffs", "Hat", "Headphones", "Scarf", "Hood"], "answer1_count": 49},
            {"question": "Name Something A Politician Does When Scandalous News Breaks Out About Them", "answers": ["Lie/Deny It", "Go Into Hiding", "Apology/Press Conference", "Resign"], "answer1_count": 49},
        ]
    },
    "survey5": {
        "name": "Survey 5",
        "rounds": [
            {"question": "Name a chore kids try to avoid", "answers": ["Doing the Dishes", "Cleaning Their Room", "Taking Out The Trash", "Mowing The Lawn"], "answer1_count": 43},
            {"question": "What Would You Hear On The Radio That Would Make You Turn The Station?", "answers": ["Commercial", "News", "Bad Song", "Static", "Cursing"], "answer1_count": 34},
            {"question": "Name Something People Do With Both Hands", "answers": ["Drive", "Dishes", "Type on Keyboard", "Clap", "Cook", "Eat"], "answer1_count": 40},
            {"question": "Name A Day Of The Year That Some People Don\u2019t Want To Spend Alone", "answers": ["Christmas", "Valentines Day", "Birthday", "New Years Eve"], "answer1_count": 40},
            {"question": "Name Something You Might Pay Someone To Do While You\u2019re Away On Vacation", "answers": ["Care for Pets", "House Sit", "Water Plants", "Babysit", "Collect Mail"], "answer1_count": 28},
            {"question": "What Do You Find Out About A Town By Reading Signs On The Side Of The Road?", "answers": ["Population", "Town Name", "Speed Limit"], "answer1_count": 60},
            {"question": "Other Than Academics Why Might A Teen Choose A Certain College?", "answers": ["Sports Team", "Location", "Friends are Going", "Party School", "Cost of Tuition"], "answer1_count": 47},
            {"question": "Name Something That\u2019s On Your Dinner Table Every Night That The Dog Won\u2019t Beg For", "answers": ["Veggies/Salad", "Salt", "Silverware/Plates", "Napkins"], "answer1_count": 45},
        ]
    },
    "survey6": {
        "name": "Survey 6",
        "rounds": [
            {"question": "If You\u2019re Driving In The Middle Of No Where, What Animal Might You See Crossing The Street?", "answers": ["Deer", "Cow", "Moose", "Coyote"], "answer1_count": 50},
            {"question": "how Many Hours Of Sleep Does The average person Need In Oder To Wake Up Refreshed?", "answers": ["8", "7", "6", "10", "9"], "answer1_count": 47},
            {"question": "Name Something There Are Seven Of", "answers": ["Dwarfs", "Deadly Sins", "Wonders of the World", "Days Per Week", "Sins", "Continents"], "answer1_count": 28},
            {"question": "Name An Activity That\u2019d Be Hard To Do By Candlelight", "answers": ["Read", "Cook", "Write", "Sewing/Knitting"], "answer1_count": 62},
            {"question": "Name Something That Happens To An Old Person\u2019s Body, That You\u2019d Be Surprised To Hear A teen Complaining About", "answers": ["Wrinkles", "Arthritis", "Gray Hair", "Sagging", "Back Ache"], "answer1_count": 50},
            {"question": "Name something you might find in a kitchen drawer", "answers": ["Beach", "Spa", "Park"], "answer1_count": 51},
            {"question": "Name A Good Place To Put Your Hands While Kissing Someone", "answers": ["Their Face", "Around Their Neck", "Their Hips", "Their Back", "Their Shoulders"], "answer1_count": 27},
            {"question": "Instead Of Their First Name, What Might A Parent Shout When Calling For Their Child?", "answers": ["Whole Name", "Nickname", "Hey!", "Siblings Name"], "answer1_count": 38},
        ]
    },
    "survey7": {
        "name": "Survey 7",
        "rounds": [
            {"question": "Name A Place An Animal Might Take A Bath, But You Never Would", "answers": ["Lake/Pond", "Puddle", "River", "Bird Bath"], "answer1_count": 51},
            {"question": "Name a job title someone might have in a big company", "answers": ["CEO", "President", "Vice President", "Supervisor", "Manager"], "answer1_count": 39},
            {"question": "Name A Job Where It Would Be Okay To Yell At Work", "answers": ["Construction", "Sports", "Teacher", "Police", "Stock Brocker", "Auctioneer"], "answer1_count": 43},
            {"question": "What Are 2 Brothers Most Likely To Fight Over?", "answers": ["Girls", "Toys", "TV Remote", "Attention"], "answer1_count": 45},
            {"question": "Name A Way You Can Tell A Storm Is Coming", "answers": ["Dark Clouds", "Lightning", "Wind Changes", "Smell", "Drizzling"], "answer1_count": 61},
            {"question": "Name Something A Plane Can't Fly Without", "answers": ["Wings", "Fuel", "A Pilot"], "answer1_count": 42},
            {"question": "Tell Me A Reason You Might Be Low On Sleep", "answers": ["Overworked", "Kids/New Baby", "Can\u2019t Sleep", "Sick", "Studying"], "answer1_count": 44},
            {"question": "Name A Color Baby Clothes Comes in", "answers": ["Pink", "Blue", "Yellow", "Green"], "answer1_count": 47},
        ]
    },
    "survey8": {
        "name": "Survey 8",
        "rounds": [
            {"question": "Last thing you\u2019d want to happen at the airport", "answers": ["Miss Flight", "Lose Luggage", "Stopped By Security", "Delayed"], "answer1_count": 0},
            {"question": "Something you do when approached by a salesperson", "answers": ["Avoid Them", "Ask For Help", "Smile", "Say Hi", "Just Looking"], "answer1_count": 0},
            {"question": "Someone you hope never writes a tell-all book", "answers": ["Parent", "Significant Other", "Ex", "Best Friend", "Sibling", "Son/Daughter"], "answer1_count": 0},
            {"question": "Something people check on their smartwatch", "answers": ["Steps", "Notifications", "Heart Rate", "Time"], "answer1_count": 0},
            {"question": "Famous phrase from The Wizard of Oz", "answers": ["Off To See The Wizard", "No Place Like Home", "Follow The Yellow Brick Road", "I\u2019ll Get You My Pretty", "Lions, Tigers & Bears!"], "answer1_count": 0},
            {"question": "Place people stash a spare charging cable", "answers": ["Car", "Work Desk", "Backpack/Purse"], "answer1_count": 0},
            {"question": "Chore that takes less than 10 minutes", "answers": ["Take Out Trash", "Wipe Counter", "Make The Bed", "Load Dishwasher", "Water Plants"], "answer1_count": 0},
            {"question": "Feature people look for in a new phone", "answers": ["Battery Life", "Camera", "Price", "Storage"], "answer1_count": 0},
        ]
    },
    "survey9": {
        "name": "Survey 9",
        "rounds": [
            {"question": "How Many Dates Should You Go On Before Kissing Someone?", "answers": ["2", "3", "1", "5"], "answer1_count": 43},
            {"question": "What Would You Hear On The Radio That Would Make You Turn The Station?", "answers": ["Commercial", "News", "Bad Song", "Static", "Cursing"], "answer1_count": 34},
            {"question": "Name Something People Do With Both Hands", "answers": ["Drive", "Dishes", "Tie Shoelaces", "Clap", "Cook", "Put on a Coat"], "answer1_count": 40},
            {"question": "Name A Day Of The Year That Some People Don\u2019t Want To Spend Alone", "answers": ["Christmas", "Valentines Day", "Birthday", "New Years Eve"], "answer1_count": 40},
            {"question": "Name Something You Might Pay Someone To Do While You\u2019re Away On Vacation", "answers": ["Care for Pets", "House Sit", "Water Plants", "Babysit", "Collect Mail"], "answer1_count": 28},
            {"question": "What Do You Find Out About A Town By Reading Signs On The Side Of The Road?", "answers": ["Population", "Town Name", "Speed Limit"], "answer1_count": 60},
            {"question": "Other Than Academics Why Might A Teen Choose A Certain College?", "answers": ["Sports Team", "Location", "Friends are Going", "Party School", "Cost of Tuition"], "answer1_count": 47},
            {"question": "Name Something That\u2019s On Your Dinner Table Every Night That The Dog Won\u2019t Beg For", "answers": ["Veggies/Salad", "Salt", "Silverware/Plates", "Napkins"], "answer1_count": 45},
        ]
    },
}

def get_qr_base_url():
    """Get QR code base URL from settings or config defaults."""
    return get_setting('qr_base_url', QR_DEFAULT_URL)


def parse_pptx(filepath):
    """Parse PowerPoint file and extract questions/answers

    IMPROVED VERSION - Handles text boxes with answer/count pairs
    Correctly distinguishes rank indicators (1,2,3) from answer counts (10,20,43)
    """
    from pptx import Presentation

    prs = Presentation(filepath)
    slides = list(prs.slides)

    rounds_data = []

    # Strategy: Find question slides, then parse the next slide as answers
    i = 0
    while i < len(slides):
        slide = slides[i]

        # Extract all text from current slide
        all_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    all_text.append(text)

        # Check if this is a question slide
        is_question_slide = False
        question_text = ""

        for text in all_text:
            # Look for question markers
            if 'Survey Has' in text and 'Responses' in text:
                is_question_slide = True
            # Extract actual question (not the metadata)
            elif len(text) > 10 and 'Round #' not in text and 'Survey Has' not in text:
                if not question_text:  # Take first substantial text
                    question_text = text

        if is_question_slide and i + 1 < len(slides):
            # Next slide should have answers
            answer_slide = slides[i + 1]
            answers = []

            # Extract all text from answer slide
            answer_text_elements = []
            for shape in answer_slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        answer_text_elements.append(text)

            # Parse answer/count pairs from text
            j = 0
            while j < len(answer_text_elements):
                text = answer_text_elements[j]

                # Skip UI elements
                skip_keywords = ['Round:', 'ROUND', 'Score Multiplier:', 'BACK TO SCORES',
                               'NEXT ROUND', 'And The Survey Says', 'X', '\u00ab', '\u00bb',
                               'type only', '(type', 'Click', 'Press']

                if any(keyword in text for keyword in skip_keywords):
                    j += 1
                    continue

                # Skip rank indicators (single-digit numbers 1-8)
                if text.isdigit() and len(text) <= 2 and int(text) <= 8:
                    j += 1
                    continue

                # If it's text (potential answer), look for count
                if not text.isdigit() and len(text) > 1:
                    answer_text = text
                    count = 0

                    # Look ahead for count
                    if j + 1 < len(answer_text_elements):
                        next_text = answer_text_elements[j + 1]
                        if next_text.isdigit():
                            try:
                                count_value = int(next_text)
                                # Count numbers are typically 5+ (answer counts, not ranks)
                                # But also accept low counts (some answers might have count of 1-4)
                                # The key is: if we just read answer text, next number IS the count
                                if count_value > 0:
                                    answers.append({'answer': answer_text, 'count': count_value})
                                    j += 1  # Skip the count
                            except ValueError:
                                pass

                j += 1

            # Add round if we found answers
            if answers:
                rounds_data.append({
                    'question': question_text,
                    'answers': answers
                })

            # Skip the answer slide
            i += 2
        else:
            i += 1

    return rounds_data


# ============= HOST ROUTES =============

@host_bp.route('/')
def index():
    return redirect(url_for('auth.host_login'))

@host_bp.route('/host')
@host_required
def host_dashboard():
    """Main host dashboard"""
    logger.debug("[HOST] host_dashboard() - loading dashboard")
    with db_connect() as conn:
        codes_raw = conn.execute("""
            SELECT code, used, team_name, reconnected, last_heartbeat
            FROM team_codes
            ORDER BY id ASC
        """).fetchall()

        # Process codes to add active status
        codes = []
        for code in codes_raw:
            code_dict = dict(code)
            # Calculate if team is active (heartbeat within last 30 seconds)
            if code['last_heartbeat']:
                try:
                    # Parse timestamp
                    last_hb = datetime.fromisoformat(code['last_heartbeat'])
                    now = datetime.now()
                    time_diff = (now - last_hb).total_seconds()
                    code_dict['is_active'] = time_diff < 30
                except:
                    code_dict['is_active'] = False
            else:
                code_dict['is_active'] = False
            codes.append(code_dict)

        rounds = conn.execute("SELECT * FROM rounds ORDER BY round_number ASC").fetchall()
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()

        # Count unscored submissions for active round
        unscored_count = 0
        submission_count = 0
        if active_round:
            unscored_count = conn.execute("""
                SELECT COUNT(*) as cnt FROM submissions
                WHERE round_id = ? AND scored = 0
            """, (active_round['id'],)).fetchone()['cnt']

            # Total submissions for active round
            submission_count = conn.execute("""
                SELECT COUNT(*) as cnt FROM submissions
                WHERE round_id = ?
            """, (active_round['id'],)).fetchone()['cnt']

    logger.debug(f"[HOST] host_dashboard() - {len(codes)} codes, {len(rounds)} rounds, "
                 f"active_round={'R'+str(active_round['round_number']) if active_round else 'None'}, "
                 f"submissions={submission_count}, unscored={unscored_count}")
    return render_template('host.html',
                         codes=codes,
                         rounds=[dict(r) for r in rounds],
                         active_round=dict(active_round) if active_round else None,
                         unscored_count=unscored_count,
                         submission_count=submission_count,
                         rounds_config=ROUNDS_CONFIG,
                         ai_scoring_available=AI_SCORING_ENABLED)

@host_bp.route('/host/codes-status')
@host_required
def codes_status():
    """API endpoint - returns code statuses as JSON for auto-refresh"""
    logger.debug("[CODES] codes_status() called")
    with db_connect() as conn:
        codes = conn.execute("""
            SELECT code, used, team_name
            FROM team_codes
            ORDER BY id ASC
        """).fetchall()

        codes_data = []
        for code in codes:
            codes_data.append({
                'code': code['code'],
                'used': bool(code['used']),
                'team_name': code['team_name'] if code['team_name'] else None
            })

        used_count = sum(1 for c in codes_data if c['used'])

        logger.debug(f"[CODES] codes_status() returning {len(codes_data)} total, {used_count} used")

        return jsonify({
            'codes': codes_data,
            'total': len(codes_data),
            'used': used_count
        })

@host_bp.route('/host/generate-codes', methods=['POST'])
@host_required
def generate_codes():
    """Reload fixed team codes from codes.json"""
    logger.info("[CODES] generate_codes() - reloading fixed codes from codes.json")
    ensure_fixed_codes()
    codes = load_fixed_codes()
    logger.info(f"[CODES] generate_codes() - {len(codes)} fixed codes loaded")
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
        <title>Codes Generated</title>
        <style>
            body {{
                font-family: Arial;
                background: linear-gradient(135deg, #1e3c72 0%, #2a5298 100%);
                color: #ffd700;
                text-align: center;
                padding: 50px;
            }}
            .box {{
                background: #000;
                border: 3px solid #ffd700;
                padding: 40px;
                border-radius: 15px;
                max-width: 500px;
                margin: 0 auto;
            }}
            h1 {{ font-size: 3em; margin-bottom: 20px; }}
            button {{
                background: #ffd700;
                color: #000;
                border: none;
                padding: 20px 40px;
                font-size: 1.5em;
                font-weight: bold;
                border-radius: 10px;
                cursor: pointer;
                margin-top: 20px;
            }}
            button:hover {{ background: #ffed4e; }}
        </style>
    </head>
    <body>
        <div class="box">
            <h1>\u2705 Success!</h1>
            <p style="font-size: 1.5em;">{len(codes)} fixed team codes loaded!</p>
            <button onclick="window.location.href='/host'">Back to Dashboard</button>
        </div>
    </body>
    </html>
    """

@host_bp.route('/host/reclaim-code/<code>', methods=['POST'])
@host_required
def reclaim_code(code):
    """Reclaim a used code - removes team and frees code for reuse"""
    code = code.upper()
    logger.debug(f"[CODES] reclaim_code() - attempting to reclaim code={code}")

    with db_connect() as conn:
        code_row = conn.execute("SELECT * FROM team_codes WHERE code = ?", (code,)).fetchone()

        if not code_row:
            logger.warning(f"[CODES] reclaim_code() - code={code} not found")
            return jsonify({"success": False, "message": "Code not found"}), 404

        if not code_row['used']:
            logger.warning(f"[CODES] reclaim_code() - code={code} is not in use")
            return jsonify({"success": False, "message": "Code is not in use"}), 400

        team_name = code_row['team_name']

        # Delete all submissions for this team
        conn.execute("DELETE FROM submissions WHERE code = ?", (code,))

        # Reset the code (mark as unused, clear team name and reconnect flag)
        conn.execute("""
            UPDATE team_codes
            SET used = 0, team_name = NULL, reconnected = 0, last_heartbeat = NULL
            WHERE code = ?
        """, (code,))

        conn.commit()

        logger.info(f"[CODES] reclaim_code() - code={code} reclaimed from team='{team_name}', submissions deleted")

        return jsonify({
            "success": True,
            "message": f"Code {code} reclaimed. Team '{team_name}' removed."
        })

@host_bp.route('/host/print-codes')
@host_required
def print_codes():
    """Generate landscape HTML page with QR codes for mobile play (replaces paper)"""
    logger.debug("[CODES] print_codes() - generating mobile play QR code page")
    codes = load_fixed_codes()
    if not codes:
        return "No codes available. Generate codes first!", 400
    qr_base_url = get_qr_base_url()
    return render_template('print_qr_codes.html', codes=codes, qr_base_url=qr_base_url,
                           mode='play', title='Mobile Play \u2014 Scan to Play on Your Phone')

@host_bp.route('/host/print-codes-landscape')
@host_required
def print_codes_landscape():
    """Generate landscape HTML page with QR codes for view-only status (companion to paper)"""
    logger.debug("[CODES] print_codes_landscape() - generating view-only QR code page")
    codes = load_fixed_codes()
    if not codes:
        return "No codes available. Generate codes first!", 400
    qr_base_url = get_qr_base_url()
    return render_template('print_qr_codes.html', codes=codes, qr_base_url=qr_base_url,
                           mode='view', title='View Only \u2014 See Your Submitted Answers')

@host_bp.route('/host/print-answer-sheets')
@host_required
def print_answer_sheets():
    """Generate printable answer sheets with pre-printed codes.
    Accepts ?group=1 (codes 1-30) or ?group=2 (codes 31-60).
    """
    all_codes = load_fixed_codes()
    group = request.args.get('group', '1')
    if group == '2':
        codes = all_codes[30:60]
        group_label = 'Group 2 (31-60)'
    else:
        codes = all_codes[0:30]
        group_label = 'Group 1 (1-30)'
    logger.info(f"[CODES] print_answer_sheets() - generating {group_label} ({len(codes)} codes)")
    qr_base_url = get_qr_base_url()

    return render_template('print_answer_sheets.html', codes=codes, group_label=group_label, rounds_config=ROUNDS_CONFIG, qr_base_url=qr_base_url)

@host_bp.route('/host/upload-answers', methods=['POST'])
@host_required
def upload_answers():
    """Upload DOCX or PPTX answer sheet and auto-create all rounds"""
    logger.info("[UPLOAD] upload_answers() - file upload started")
    try:
        if 'file' not in request.files:
            logger.warning("[UPLOAD] No file in request")
            flash('No file uploaded!', 'error')
            return redirect(url_for('.host_dashboard'))

        file = request.files['file']
        if file.filename == '':
            logger.warning("[UPLOAD] No file selected")
            flash('No file selected!', 'error')
            return redirect(url_for('.host_dashboard'))

        # Accept .docx, .pptx, and .pptm files
        file_ext = os.path.splitext(file.filename)[1].lower()
        logger.info(f"[UPLOAD] File received: '{file.filename}', type={file_ext}")
        if file_ext not in ['.docx', '.pptx', '.pptm']:
            logger.warning(f"[UPLOAD] Invalid file type: {file_ext}")
            flash('Please upload a .docx, .pptx, or .pptm file!', 'error')
            return redirect(url_for('.host_dashboard'))

        # Save temp file
        temp_path = os.path.join(BASE_DIR, f'temp_answers_{int(time.time())}{file_ext}')
        file.save(temp_path)

        # Parse based on file type
        rounds_data = []

        if file_ext == '.docx':
            # DOCX parsing (existing code)
            from docx import Document
            doc = Document(temp_path)

            # ROBUST: Extract questions with flexible matching (handles both - and \u2013 dashes)
            questions = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text and len(text) > 0 and text[0].isdigit():
                    # Match both regular dash (-) and em-dash (\u2013)
                    if '-' in text or '\u2013' in text:
                        # Split on either dash type
                        separator = '\u2013' if '\u2013' in text else '-'
                        parts = text.split(separator, 1)
                        if len(parts) > 1:
                            question = parts[1].strip()
                            questions.append(question)

            # ROBUST: Parse ALL 8 tables regardless of question count
            for table_idx, table in enumerate(doc.tables):
                if table_idx >= 8:
                    break

                answers = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if len(cells) >= 3:
                        # Skip header rows - only process if first cell is a number (rank)
                        if not cells[0] or not cells[0].strip().isdigit():
                            continue

                        answer = cells[1]
                        points_count = cells[2]

                        # ROBUST: Flexible count parsing (handles various spacing)
                        count = 0
                        if points_count:
                            # Try both dash types
                            for separator in ['-', '\u2013']:
                                if separator in points_count:
                                    parts = points_count.split(separator)
                                    if len(parts) > 1:
                                        try:
                                            # Extract just the digits from the second part
                                            count_str = ''.join(filter(str.isdigit, parts[1]))
                                            if count_str:
                                                count = int(count_str)
                                            break
                                        except ValueError:
                                            count = 0

                        answers.append({'answer': answer, 'count': count})

                # Use question by index, or empty string if not found
                question = questions[table_idx] if table_idx < len(questions) else ''

                rounds_data.append({
                    'question': question,
                    'answers': answers
                })

        elif file_ext in ['.pptx', '.pptm']:
            # PowerPoint parsing (new code)
            rounds_data = parse_pptx(temp_path)

        # Always create all rounds found (should be 8)
        with db_connect() as conn:
            conn.execute("DELETE FROM rounds")
            conn.execute("DELETE FROM submissions")

            for idx, round_data in enumerate(rounds_data):
                round_num = idx + 1
                config = ROUNDS_CONFIG[idx]
                num_answers = config['answers']

                fields = ['round_number', 'question', 'num_answers', 'is_active']
                values = [round_num, round_data['question'], num_answers, 0]

                for i in range(1, num_answers + 1):
                    if i <= len(round_data['answers']):
                        fields.append(f'answer{i}')
                        fields.append(f'answer{i}_count')
                        values.append(round_data['answers'][i-1]['answer'])
                        values.append(round_data['answers'][i-1]['count'])

                placeholders = ','.join(['?'] * len(values))
                conn.execute(f"INSERT INTO rounds ({','.join(fields)}) VALUES ({placeholders})", values)

            conn.commit()

        if os.path.exists(temp_path):
            os.remove(temp_path)

        rounds_created = len(rounds_data)
        logger.info(f"[UPLOAD] Complete: {rounds_created} rounds created from '{file.filename}'")
        for idx, rd in enumerate(rounds_data):
            logger.debug(f"[UPLOAD]   Round {idx+1}: Q='{rd['question'][:60]}', {len(rd['answers'])} answers")
        flash(f'\u2705 Success! {rounds_created} rounds created!', 'success')
        return redirect(url_for('.host_dashboard'))

    except FileNotFoundError as e:
        logger.error(f"[UPLOAD] FileNotFoundError: {e}")
        try:
            if 'temp_path' in locals() and os.path.exists(temp_path):
                os.remove(temp_path)
        except:
            pass
        flash(f'\u274c File error: Could not read the uploaded file. Please try again.', 'error')
        return redirect(url_for('.host_dashboard'))
    except ImportError as e:
        logger.error(f"[UPLOAD] ImportError (missing library): {e}")
        try:
            if 'temp_path' in locals() and os.path.exists(temp_path):
                os.remove(temp_path)
        except:
            pass
        flash(f'\u274c Missing library: {str(e)}. Please install required dependencies.', 'error')
        return redirect(url_for('.host_dashboard'))
    except Exception as e:
        logger.error(f"[UPLOAD] Unexpected error: {type(e).__name__}: {e}")
        import traceback
        logger.error(f"[UPLOAD] Traceback:\n{traceback.format_exc()}")
        try:
            if 'temp_path' in locals() and os.path.exists(temp_path):
                os.remove(temp_path)
        except:
            pass

        # Provide helpful error messages based on the error type
        error_msg = str(e)
        if 'pptx' in error_msg.lower() or 'presentation' in error_msg.lower():
            flash(f'\u274c PowerPoint parsing error: The file format may be corrupted or unsupported. Details: {error_msg}', 'error')
        elif 'docx' in error_msg.lower() or 'document' in error_msg.lower():
            flash(f'\u274c Word document parsing error: The file format may be corrupted. Details: {error_msg}', 'error')
        elif 'table' in error_msg.lower():
            flash(f'\u274c Table parsing error: Could not read answer tables. Make sure your file has the correct format. Details: {error_msg}', 'error')
        else:
            flash(f'\u274c Upload failed: {error_msg}', 'error')

        return redirect(url_for('.host_dashboard'))

@host_bp.route('/host/round/create', methods=['POST'])
@host_required
def create_round():
    """Create a round manually"""
    round_num = int(request.form.get('round_number'))
    question = request.form.get('question', '').strip()
    logger.info(f"[ROUND] create_round() - round_num={round_num}, question='{question[:50]}'")

    config = next((r for r in ROUNDS_CONFIG if r['round'] == round_num), None)
    if not config:
        logger.warning(f"[ROUND] create_round() - invalid round number: {round_num}")
        return "Invalid round number", 400

    with db_connect() as conn:
        conn.execute("UPDATE rounds SET is_active = 0")
        conn.execute("""
            INSERT INTO rounds (round_number, question, num_answers, is_active)
            VALUES (?, ?, ?, 1)
        """, (round_num, question, config['answers']))
        conn.commit()
    logger.info(f"[ROUND] create_round() - round {round_num} created and activated")
    return redirect(url_for('.host_dashboard'))

@host_bp.route('/host/round/<int:round_id>/activate', methods=['POST'])
@host_required
def activate_round(round_id):
    """Activate a specific round"""
    logger.info(f"[ROUND] activate_round() - requesting activation of round_id={round_id}")
    with db_connect() as conn:
        # CRITICAL FIX: Validate that round has answers before activating
        round_data = conn.execute(
            "SELECT answer1, question FROM rounds WHERE id = ?",
            (round_id,)
        ).fetchone()

        if not round_data:
            logger.warning(f"[ROUND] activate_round() - round_id={round_id} not found")
            flash('\u274c Round not found!', 'error')
            return redirect(url_for('.host_dashboard'))

        if not round_data['answer1']:
            logger.warning(f"[ROUND] activate_round() - round_id={round_id} has no answers, blocking activation")
            flash('\u274c Cannot activate round without answers! Please set answers first.', 'error')
            return redirect(url_for('.host_dashboard'))

        # CRITICAL FIX: Use transaction to prevent race conditions
        # Deactivate ALL rounds, then activate the selected one atomically
        conn.execute("BEGIN IMMEDIATE")  # Lock database to prevent race conditions
        try:
            conn.execute("UPDATE rounds SET is_active = 0")
            conn.execute("UPDATE rounds SET is_active = 1 WHERE id = ?", (round_id,))
            conn.commit()
            logger.info(f"[ROUND] activate_round() - round_id={round_id} now active (deactivated all others)")
            flash(f'\u2705 Round activated: {round_data["question"]}', 'success')
        except Exception as e:
            conn.rollback()
            logger.error(f"[ROUND] activate_round() - error: {e}")
            flash(f'\u274c Error activating round: {str(e)}', 'error')

    return redirect(url_for('.host_dashboard'))

@host_bp.route('/host/round/<int:round_id>/answers', methods=['POST'])
@host_required
def set_answers(round_id):
    """Set correct answers for a round"""
    logger.debug(f"[ROUND] set_answers() - round_id={round_id}")
    with db_connect() as conn:
        round_info = conn.execute("SELECT * FROM rounds WHERE id = ?", (round_id,)).fetchone()
        num_answers = round_info['num_answers']
        logger.debug(f"[ROUND] Setting {num_answers} answers for round {round_info['round_number']}")

        fields = []
        values = []
        for i in range(1, 7):
            if i <= num_answers:
                fields.append(f'answer{i} = ?')
                if i == 1:
                    fields.append(f'answer{i}_count = ?')
                    values.append(request.form.get(f'answer{i}', '').strip())
                    values.append(int(request.form.get(f'answer{i}_count', 0) or 0))
                else:
                    values.append(request.form.get(f'answer{i}', '').strip())

        values.append(round_id)
        conn.execute(f"UPDATE rounds SET {', '.join(fields)} WHERE id = ?", values)
        conn.commit()
    logger.info(f"[ROUND] set_answers() - answers saved for round_id={round_id}")
    return redirect(url_for('.host_dashboard'))

@host_bp.route('/host/check-active-round')
@host_required
def check_active_round():
    """API endpoint to check if there's an active round (for AJAX polling)"""
    with db_connect() as conn:
        active_round = conn.execute("SELECT id FROM rounds WHERE is_active = 1").fetchone()
        has_active = active_round is not None
        logger.debug(f"[API] check_active_round() = {has_active}")
        return jsonify({'has_active_round': has_active})

@host_bp.route('/host/start-next-round', methods=['POST'])
@host_required
def start_next_round():
    """Move to next round"""
    logger.debug("[ROUND] start_next_round() called")
    with db_connect() as conn:
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()

        if active_round:
            current_num = active_round['round_number']
            logger.info(f"[ROUND] Current active round: {current_num}, advancing to {current_num + 1}")
            # Deactivate current
            conn.execute("UPDATE rounds SET is_active = 0 WHERE id = ?", (active_round['id'],))

            # Activate next round
            next_round = conn.execute("""
                SELECT * FROM rounds WHERE round_number = ?
            """, (current_num + 1,)).fetchone()

            if next_round:
                conn.execute("UPDATE rounds SET is_active = 1 WHERE id = ?", (next_round['id'],))
                conn.commit()
                logger.info(f"[ROUND] Activated round {current_num + 1} (id={next_round['id']})")
            else:
                # No more rounds - game over
                conn.commit()
                logger.info(f"[ROUND] No round {current_num + 1} found - game complete!")
                flash('All rounds complete!', 'info'); return redirect(url_for('.host_dashboard'))

        conn.commit()

    return redirect(url_for('.host_dashboard'))

@host_bp.route('/host/round/<int:round_id>/edit-answer/<int:answer_num>')
@host_required
def edit_single_answer(round_id, answer_num):
    """Edit a single answer"""
    logger.info(f"[ROUND] edit_single_answer() - round_id={round_id}, answer_num={answer_num}")
    with db_connect() as conn:
        round_info = conn.execute("SELECT * FROM rounds WHERE id = ?", (round_id,)).fetchone()

        if not round_info:
            logger.warning(f"[ROUND] edit_single_answer() - round_id={round_id} not found")
            flash('Round not found!', 'error'); return redirect(url_for('.host_dashboard'))

        current_answer = round_info[f'answer{answer_num}']
        current_count = round_info['answer1_count'] if answer_num == 1 else None

    return render_template('edit_answer.html',
                         round=dict(round_info),
                         answer_num=answer_num,
                         current_answer=current_answer,
                         current_count=current_count)

@host_bp.route('/host/round/<int:round_id>/update-answer/<int:answer_num>', methods=['POST'])
@host_required
def update_single_answer(round_id, answer_num):
    """Update a single answer"""
    new_answer = request.form.get('answer', '').strip()
    logger.info(f"[ROUND] update_single_answer() - round_id={round_id}, answer_num={answer_num}, new_answer='{new_answer}'")

    with db_connect() as conn:
        if answer_num == 1:
            new_count = int(request.form.get('count', 0) or 0)
            conn.execute("""
                UPDATE rounds
                SET answer1 = ?, answer1_count = ?
                WHERE id = ?
            """, (new_answer, new_count, round_id))
        else:
            conn.execute(f"""
                UPDATE rounds
                SET answer{answer_num} = ?
                WHERE id = ?
            """, (new_answer, round_id))

        conn.commit()

    flash('\u2705 Answer updated!', 'success'); return redirect(url_for('.host_dashboard'))

@host_bp.route('/host/create-round-manual')
@host_required
def create_round_manual_form():
    """Show manual round creation form"""
    return render_template('create_round_manual.html',
                         rounds_config=ROUNDS_CONFIG,
                         prebuilt_surveys=PREBUILT_SURVEYS)

@host_bp.route('/host/create-round-manual/submit', methods=['POST'])
@host_required
def create_round_manual_submit():
    """Process manual round creation for ALL 8 rounds"""
    logger.info("[ROUND] create_round_manual_submit() - creating all 8 rounds manually")
    try:
        with db_connect() as conn:
            # Delete any existing rounds and submissions
            conn.execute("DELETE FROM rounds")
            conn.execute("DELETE FROM submissions")

            # Create all 8 rounds
            for config in ROUNDS_CONFIG:
                round_num = config['round']
                num_answers = config['answers']

                # Get question for this round
                question = request.form.get(f'question{round_num}', '').strip()

                # Build insert for this round
                fields = ['round_number', 'question', 'num_answers', 'is_active']
                is_active = 0
                values = [round_num, question, num_answers, is_active]

                # Get answers for this round
                for i in range(1, num_answers + 1):
                    answer = request.form.get(f'round{round_num}_answer{i}', '').strip()
                    fields.append(f'answer{i}')
                    values.append(answer)

                    # Get count only for answer #1
                    if i == 1:
                        count = int(request.form.get(f'round{round_num}_answer1_count', 0) or 0)
                        fields.append(f'answer{i}_count')
                        values.append(count)

                # Insert this round
                placeholders = ','.join(['?'] * len(values))
                conn.execute(f"INSERT INTO rounds ({','.join(fields)}) VALUES ({placeholders})", values)

            conn.commit()

        logger.info("[ROUND] create_round_manual_submit() - all 8 rounds created successfully")
        flash('\u2705 All 8 rounds created!', 'success'); return redirect(url_for('.host_dashboard'))

    except Exception as e:
        logger.error(f"[ROUND] create_round_manual_submit() error: {e}")
        flash(f'Error creating rounds: {str(e)}', 'error'); return redirect(url_for('.host_dashboard'))

@host_bp.route('/host/reset', methods=['POST'])
@host_required
def reset_game():
    """Reset game but keep codes and team names - for setup fixes"""
    logger.info("[HOST] reset_game() - resetting game (keeping teams)")
    with db_connect() as conn:
        conn.execute("DELETE FROM submissions")
        conn.execute("DELETE FROM rounds")
        # DO NOT touch team_codes table - keep teams joined!
        conn.commit()
    logger.info("[HOST] reset_game() - submissions and rounds deleted, team_codes untouched")
    flash('Game reset! Teams are still joined. Upload new questions to start fresh.', 'success')
    return redirect(url_for('.host_dashboard'))

@host_bp.route('/host/reset-all', methods=['POST'])
@host_required
def reset_all():
    """Reset everything - clear teams but keep code values"""
    with db_connect() as conn:
        conn.execute("DELETE FROM submissions")
        conn.execute("DELETE FROM rounds")
        # Reset codes to unused but keep the code values (HNCL, LZLX, etc)
        conn.execute("UPDATE team_codes SET used = 0, team_name = NULL")
        conn.commit()

    # Increment reset counter to invalidate all team sessions
    reset_state['counter'] += 1
    logger.info(f"[HOST] reset_all() - reset counter incremented to {reset_state['counter']}")
    logger.info("[HOST] All team sessions are now invalid - teams will see Game Over page")

    flash('Everything reset! All codes are now unused and ready for new teams.', 'success')
    return redirect(url_for('.host_dashboard'))

@host_bp.route('/host/settings', methods=['GET', 'POST'])
@host_required
def settings():
    """Settings page for configuring game options"""
    if request.method == 'POST':
        # Get form data
        qr_base_url = request.form.get('qr_base_url', '').strip()
        logger.info(f"[SETTINGS] settings() POST - qr_base_url='{qr_base_url}'")

        # Basic validation
        if not qr_base_url:
            flash('QR Base URL cannot be empty!', 'error')
        elif ' ' in qr_base_url:
            flash('URL cannot contain spaces!', 'error')
        else:
            # Save setting
            if set_setting('qr_base_url', qr_base_url, 'Base URL for QR codes on printed sheets'):
                flash('Settings saved successfully!', 'success')
            else:
                flash('Failed to save settings. Please try again.', 'error')

        return redirect(url_for('.settings'))

   # GET - show form with current settings
    current_qr_url = get_qr_base_url()
    allow_team_registration = get_setting('allow_team_registration', 'true') == 'true'
    system_paused = get_setting('system_paused', 'false') == 'true'
    broadcast_message = get_setting('broadcast_message', '')
    ai_scoring_enabled = get_setting('ai_scoring_enabled', 'true') == 'true'
    extended_thinking_enabled = get_setting('extended_thinking_enabled', 'false') == 'true'
    thinking_budget_tokens = int(get_setting('thinking_budget_tokens', '10000'))

    # Count corrections in current session
    corrections_count = len(load_corrections_history())

    return render_template('settings.html',
                         qr_base_url=current_qr_url,
                         allow_team_registration=allow_team_registration,
                         system_paused=system_paused,
                         broadcast_message=broadcast_message,
                         ai_scoring_available=AI_SCORING_ENABLED,
                         ai_scoring_enabled=ai_scoring_enabled,
                         corrections_count=corrections_count,
                         ai_model_choices=AI_MODEL_CHOICES,
                         current_ai_model=get_current_ai_model(),
                         extended_thinking_enabled=extended_thinking_enabled,
                         thinking_budget_tokens=thinking_budget_tokens)

@host_bp.route('/host/save-training', methods=['POST'])
@host_required
def save_training():
    """Save AI corrections to GitHub repo for long-term persistence."""
    if not GITHUB_TOKEN:
        return jsonify({'success': False, 'error': 'GITHUB_TOKEN not configured. Set it in Render environment variables.'}), 400

    corrections = load_corrections_history()
    if not corrections:
        return jsonify({'success': False, 'error': 'No corrections to save.'}), 400

    try:
        file_path = 'corrections_history.json'
        api_url = f'https://api.github.com/repos/{GITHUB_REPO}/contents/{file_path}'

        # First, get the current file SHA (needed for updates)
        get_req = urllib.request.Request(api_url, headers={
            'Authorization': f'token {GITHUB_TOKEN}',
            'Accept': 'application/vnd.github.v3+json'
        })

        existing_sha = None
        existing_data = []
        try:
            with urllib.request.urlopen(get_req) as resp:
                file_info = json.loads(resp.read().decode())
                existing_sha = file_info.get('sha')
                # Decode existing content and merge
                existing_content = base64.b64decode(file_info.get('content', '')).decode('utf-8')
                existing_data = json.loads(existing_content) if existing_content.strip() else []
        except urllib.error.HTTPError as e:
            if e.code == 404:
                existing_data = []  # File doesn't exist yet
            else:
                raise

        # Merge: add new corrections that aren't already in the file
        # Use a simple dedup by converting to comparable tuples
        existing_set = set()
        for c in existing_data:
            key = (c.get('team_answer', ''), c.get('survey_answer', ''), c.get('correction_type', ''), c.get('question', ''))
            existing_set.add(key)

        new_corrections = []
        for c in corrections:
            key = (c.get('team_answer', ''), c.get('survey_answer', ''), c.get('correction_type', ''), c.get('question', ''))
            if key not in existing_set:
                new_corrections.append(c)

        if not new_corrections:
            return jsonify({'success': True, 'message': f'All {len(corrections)} corrections already saved. No new data.'})

        merged = existing_data + new_corrections
        content_b64 = base64.b64encode(json.dumps(merged, indent=2).encode('utf-8')).decode('utf-8')

        # Commit to GitHub
        payload = json.dumps({
            'message': f'Update AI training data (+{len(new_corrections)} corrections, {len(merged)} total)',
            'content': content_b64,
            'sha': existing_sha  # None if new file
        }).encode('utf-8')

        # Remove sha key if None (new file)
        payload_dict = json.loads(payload)
        if payload_dict.get('sha') is None:
            del payload_dict['sha']
        payload = json.dumps(payload_dict).encode('utf-8')

        put_req = urllib.request.Request(api_url, data=payload, method='PUT', headers={
            'Authorization': f'token {GITHUB_TOKEN}',
            'Accept': 'application/vnd.github.v3+json',
            'Content-Type': 'application/json'
        })

        with urllib.request.urlopen(put_req) as resp:
            if resp.status in (200, 201):
                logger.info(f"[AI-CORRECTIONS] Saved {len(new_corrections)} new corrections to GitHub ({len(merged)} total)")
                return jsonify({'success': True, 'message': f'Saved {len(new_corrections)} new corrections to GitHub ({len(merged)} total)'})

        return jsonify({'success': False, 'error': 'Unexpected response from GitHub'}), 500

    except urllib.error.HTTPError as e:
        error_body = e.read().decode() if e.fp else str(e)
        logger.error(f"[AI-CORRECTIONS] GitHub API error: {e.code} - {error_body}")
        return jsonify({'success': False, 'error': f'GitHub API error ({e.code}). Check your token permissions.'}), 500
    except Exception as e:
        logger.error(f"[AI-CORRECTIONS] Failed to save to GitHub: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500


@host_bp.route('/host/clear-training', methods=['POST'])
@host_required
def clear_training():
    """Clear all AI training corrections from local file and database."""
    try:
        # Clear the local JSON file
        with open(CORRECTIONS_FILE, 'w') as f:
            json.dump([], f)
        logger.info("[AI-CORRECTIONS] Cleared corrections_history.json")

        # Clear the database table
        with db_connect() as conn:
            conn.execute("DELETE FROM ai_corrections")
            conn.commit()
        logger.info("[AI-CORRECTIONS] Cleared ai_corrections table")

        return jsonify({'success': True, 'message': 'All training data cleared.'})
    except Exception as e:
        logger.error(f"[AI-CORRECTIONS] Failed to clear training data: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500


@host_bp.route('/host/toggle-setting', methods=['POST'])
@host_required
def toggle_setting():
    """Toggle a boolean setting"""
    setting_key = request.form.get('setting_key')

    if setting_key in ['allow_team_registration', 'system_paused', 'ai_scoring_enabled', 'extended_thinking_enabled']:
        current_value = get_setting(setting_key, 'true' if setting_key == 'ai_scoring_enabled' else 'false')
        new_value = 'false' if current_value == 'true' else 'true'
        logger.info(f"[SETTINGS] toggle_setting() - {setting_key}: '{current_value}' -> '{new_value}'")

        set_setting(setting_key, new_value, '')

        # User-friendly messages
        if setting_key == 'allow_team_registration':
            if new_value == 'true':
                flash('\u2705 Team registration enabled - New teams can join!', 'success')
            else:
                flash('\ud83d\udead Team registration disabled - No new teams can join', 'success')
        elif setting_key == 'system_paused':
            if new_value == 'true':
                # Auto-disable registration when pausing
                set_setting('allow_team_registration', 'false', '')
                flash('\u23f8\ufe0f System PAUSED - Team registration also disabled', 'success')
            else:
                flash('\u25b6\ufe0f System RESUMED - Remember to re-enable registration if needed', 'success')
        elif setting_key == 'ai_scoring_enabled':
            if new_value == 'true':
                flash('\ud83e\udd16 AI Scoring enabled - AI button will appear on scoring queue', 'success')
            else:
                flash('\ud83e\udd16 AI Scoring disabled - AI button hidden from scoring queue', 'success')
        elif setting_key == 'extended_thinking_enabled':
            if new_value == 'true':
                flash('\ud83e\udde0 Extended Thinking enabled - AI will think deeper (higher cost)', 'success')
            else:
                flash('\ud83e\udde0 Extended Thinking disabled - Using standard mode', 'success')

    return redirect(url_for('.settings'))

@host_bp.route('/host/set-ai-model', methods=['POST'])
@host_required
def set_ai_model():
    """Set the AI model for scoring and photo scanning"""
    model_id = request.form.get('ai_model', '').strip()

    valid_ids = [m['id'] for m in AI_MODEL_CHOICES]
    if model_id not in valid_ids:
        flash('Invalid model selection.', 'error')
        return redirect(url_for('.settings'))

    set_setting('ai_model', model_id, 'AI model for scoring and photo scan')

    model_name = next((m['name'] for m in AI_MODEL_CHOICES if m['id'] == model_id), model_id)
    logger.info(f"[SETTINGS] AI model changed to: {model_id}")
    flash(f'AI Model set to {model_name}', 'success')

    return redirect(url_for('.settings'))

@host_bp.route('/host/set-thinking-budget', methods=['POST'])
@host_required
def set_thinking_budget():
    """Set the token budget for extended thinking"""
    budget_str = request.form.get('thinking_budget', '').strip()

    try:
        budget = int(budget_str)
    except (ValueError, TypeError):
        flash('Invalid budget value. Must be a number.', 'error')
        return redirect(url_for('.settings'))

    if budget < 1024:
        flash('Thinking budget must be at least 1,024 tokens.', 'error')
        return redirect(url_for('.settings'))

    if budget > 128000:
        flash('Thinking budget cannot exceed 128,000 tokens.', 'error')
        return redirect(url_for('.settings'))

    set_setting('thinking_budget_tokens', str(budget), 'Token budget for extended thinking')

    logger.info(f"[SETTINGS] Thinking budget changed to: {budget}")
    flash(f'Thinking budget set to {budget:,} tokens', 'success')

    return redirect(url_for('.settings'))

@host_bp.route('/host/toggle-sleep', methods=['POST'])
@host_required
def toggle_sleep():
    """Toggle server sleep mode"""
    current_value = get_setting('server_sleep', 'false')
    new_value = 'false' if current_value == 'true' else 'true'

    set_setting('server_sleep', new_value, 'Server sleep mode - stops auto-refresh')

    if new_value == 'true':
        logger.info("[SETTINGS] Server sleep mode ENABLED - team auto-refresh will stop")
        flash('\ud83d\udca4 Server sleep mode enabled - All auto-refresh stopped', 'success')
    else:
        logger.info("[SETTINGS] Server sleep mode DISABLED - team auto-refresh resumed")
        flash('\u23f0 Server awake - Auto-refresh resumed', 'success')

    return jsonify({'success': True, 'sleep_mode': new_value})

@host_bp.route('/host/get-sleep-status')
@host_required
def get_sleep_status():
    """Get current sleep mode status"""
    sleep_mode = get_setting('server_sleep', 'false')
    logger.debug(f"[API] get_sleep_status() -> {sleep_mode}")
    return jsonify({'sleep_mode': sleep_mode})

@host_bp.route('/host/send-broadcast', methods=['POST'])
@host_required
def send_broadcast():
    """Send broadcast message to all teams"""
    message = request.form.get('message', '').strip()
    logger.info(f"[HOST] send_broadcast() - message='{message[:50]}' (len={len(message)})")

    if not message:
        flash('\u26a0\ufe0f Message cannot be empty!', 'error')
        return redirect(url_for('.settings'))

    # Security: Length limit (200 chars max)
    if len(message) > 200:
        flash('\u26a0\ufe0f Message too long! Maximum 200 characters.', 'error')
        return redirect(url_for('.settings'))

    # Security: HTML escape to prevent XSS
    message = html.escape(message)

    # Store message with timestamp
    broadcast_data = {
        'message': message,
        'timestamp': time.time()
    }

    set_setting('broadcast_message', json.dumps(broadcast_data), 'Broadcast message to all teams')
    flash(f'\ud83d\udce2 Message sent to all teams!', 'success')

    return redirect(url_for('.settings'))

@host_bp.route('/host/clear-broadcast', methods=['POST'])
@host_required
def clear_broadcast():
    """Clear broadcast message"""
    logger.info("[HOST] clear_broadcast() - broadcast message cleared")
    # Set empty broadcast with current timestamp
    broadcast_data = {
        'message': '',
        'timestamp': 0
    }
    set_setting('broadcast_message', json.dumps(broadcast_data), 'Broadcast message to all teams')
    flash('\ud83d\uddd1\ufe0f Broadcast message cleared', 'success')
    return redirect(url_for('.settings'))

@host_bp.route('/host/close-round', methods=['POST'])
@host_required
def close_round():
    """Close submissions for the active round and move to scoring"""
    logger.info("[ROUND] close_round() called")
    with db_connect() as conn:
        # Get active round
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()

        if not active_round:
            logger.warning("[ROUND] close_round() - no active round to close")
            flash('\u26a0\ufe0f No active round to close', 'error')
            return redirect(url_for('.host_dashboard'))

        # Mark round as closed
        conn.execute("UPDATE rounds SET submissions_closed = 1 WHERE id = ?", (active_round['id'],))
        conn.commit()

        # Count submissions
        sub_count = conn.execute("SELECT COUNT(*) as cnt FROM submissions WHERE round_id = ?",
                                (active_round['id'],)).fetchone()['cnt']

        # Check if all submissions are already scored
        unscored_count = conn.execute(
            "SELECT COUNT(*) as cnt FROM submissions WHERE round_id = ? AND scored = 0",
            (active_round['id'],)
        ).fetchone()['cnt']

        logger.info(f"[ROUND] Round {active_round['round_number']} closed - {sub_count} submissions, {unscored_count} unscored")
        flash(f'\ud83d\udd12 Round {active_round["round_number"]} closed! {sub_count} teams submitted.', 'success')

        # If all teams are already scored, skip scoring queue and go straight to scored teams
        if unscored_count == 0 and sub_count > 0:
            logger.info(f"[ROUND] All {sub_count} teams already scored - redirecting to scored_teams")
            return redirect(url_for('scored_teams'))

        return redirect(url_for('scoring_queue'))
