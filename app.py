import os
import json
import re
import base64
import sqlite3
import secrets
import string
import time
import logging
import urllib.request
import urllib.error
from datetime import datetime
from functools import wraps
from flask import Flask, request, render_template, redirect, url_for, jsonify, session, flash
from difflib import SequenceMatcher

# AI Scoring - optional dependency
try:
    import anthropic
    ANTHROPIC_AVAILABLE = True
except ImportError:
    ANTHROPIC_AVAILABLE = False

# ===== LOGGING CONFIGURATION =====
# Set LOG_LEVEL env var to control verbosity (default: INFO)
# Use LOG_LEVEL=DEBUG for verbose output when troubleshooting
log_level_str = os.environ.get('LOG_LEVEL', 'INFO').upper()
log_level = getattr(logging, log_level_str, logging.INFO)

if os.environ.get('RENDER'):
    # Production on Render - log to stdout only
    logging.basicConfig(
        level=log_level,
        format='[%(asctime)s] [%(levelname)s] %(message)s',
        datefmt='%H:%M:%S',
        handlers=[logging.StreamHandler()]
    )
    logger = logging.getLogger(__name__)
    logger.info("FAMILY FEUD - SERVER STARTING (RENDER)")
else:
    # Local development - log to file and console
    LOG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'logs')
    os.makedirs(LOG_DIR, exist_ok=True)

    # Create log filename with timestamp
    log_filename = datetime.now().strftime('%Y-%m-%d_%H-%M-%S.log')
    log_filepath = os.path.join(LOG_DIR, log_filename)

    logging.basicConfig(
        level=log_level,
        format='[%(asctime)s] [%(levelname)s] %(message)s',
        datefmt='%H:%M:%S',
        handlers=[
            logging.FileHandler(log_filepath),
            logging.StreamHandler()
        ]
    )

    logger = logging.getLogger(__name__)
    logger.info(f"FAMILY FEUD - SERVER STARTING (log file: {log_filepath})")

# Suppress Flask/Werkzeug per-request logging noise
logging.getLogger('werkzeug').setLevel(logging.WARNING)
logger.info(f"Log level: {logging.getLevelName(log_level)} (set LOG_LEVEL=DEBUG for verbose output)")

app = Flask(__name__)
APP_VERSION = "v2.1.0 - Fusion"
# Use environment variable for secret key in production, generate random for local dev
app.secret_key = os.environ.get('SECRET_KEY', secrets.token_hex(32))

# Generate unique startup ID - changes on EVERY server restart
# This ensures old sessions are invalidated when server restarts
# SERVER RESTART = FRESH START (no data persistence)
STARTUP_ID = str(int(time.time() * 1000000))  # Unique timestamp
logger.info(f"Server startup ID: {STARTUP_ID}")
logger.info("All sessions from previous server runs are now invalid")

# Reset counter - increments when host clicks "Reset All" button
# This invalidates all team sessions without restarting server
RESET_COUNTER = 0
logger.info(f"Reset counter initialized: {RESET_COUNTER}")

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_PATH = os.path.join(BASE_DIR, "feud.db")
CORRECTIONS_FILE = os.path.join(BASE_DIR, "corrections_history.json")


def load_corrections_history():
    """Load persistent corrections from JSON file (survives deploys)."""
    try:
        if os.path.exists(CORRECTIONS_FILE):
            with open(CORRECTIONS_FILE, 'r') as f:
                data = json.load(f)
                logger.info(f"[AI-CORRECTIONS] Loaded {len(data)} corrections from history file")
                return data
    except Exception as e:
        logger.warning(f"[AI-CORRECTIONS] Failed to load corrections history: {e}")
    return []


def save_correction_to_history(correction):
    """Append a correction to the persistent JSON file."""
    try:
        history = load_corrections_history()
        history.append(correction)
        with open(CORRECTIONS_FILE, 'w') as f:
            json.dump(history, f, indent=2)
        logger.info(f"[AI-CORRECTIONS] Saved correction to history file (total: {len(history)})")
    except Exception as e:
        logger.warning(f"[AI-CORRECTIONS] Failed to save correction to history: {e}")


# Host password protection - set via environment variable or use default
HOST_PASSWORD = os.environ.get('HOST_PASSWORD', 'localdev')
if not os.environ.get('HOST_PASSWORD'):
    logger.warning("No HOST_PASSWORD env var set — using default development password")
else:
    logger.info("Host password protection enabled (custom password set)")

# AI Scoring configuration
ANTHROPIC_API_KEY = os.environ.get('ANTHROPIC_API_KEY')
ENABLE_AI_SCORING = os.environ.get('ENABLE_AI_SCORING', 'false').lower() == 'true'
AI_SCORING_ENABLED = ENABLE_AI_SCORING and ANTHROPIC_AVAILABLE and bool(ANTHROPIC_API_KEY)
if AI_SCORING_ENABLED:
    logger.info("AI Scoring: ENABLED (env var ON, SDK installed, API key configured)")
elif ENABLE_AI_SCORING and not ANTHROPIC_AVAILABLE:
    logger.warning("AI Scoring: DISABLED - anthropic SDK not installed")
elif ENABLE_AI_SCORING and not ANTHROPIC_API_KEY:
    logger.warning("AI Scoring: DISABLED - ANTHROPIC_API_KEY not set")
else:
    logger.info("AI Scoring: DISABLED (ENABLE_AI_SCORING not set)")

# AI Model selection - which Claude model to use
AI_MODEL_DEFAULT = os.environ.get('AI_MODEL', 'claude-sonnet-4-20250514')
logger.info(f"AI Model default: {AI_MODEL_DEFAULT}")

AI_MODEL_CHOICES = [
    {'id': 'claude-sonnet-4-20250514', 'name': 'Claude Sonnet 4', 'description': 'Balanced quality & cost', 'cost_note': '~$0.01/scoring'},
    {'id': 'claude-opus-4-20250514', 'name': 'Claude Opus 4', 'description': 'Highest quality, more expensive', 'cost_note': '~$0.05/scoring'},
    {'id': 'claude-haiku-4-5-20251001', 'name': 'Claude Haiku 4.5', 'description': 'Fastest & cheapest', 'cost_note': '~$0.002/scoring'},
]

# GitHub API for saving AI training data
GITHUB_TOKEN = os.environ.get('GITHUB_TOKEN')
GITHUB_REPO = os.environ.get('GITHUB_REPO', 'teksabian/family-feud')
if GITHUB_TOKEN:
    logger.info(f"GitHub API: ENABLED (repo: {GITHUB_REPO})")
else:
    logger.info("GitHub API: DISABLED (no GITHUB_TOKEN env var)")

@app.context_processor
def inject_version():
    """Make app version and cache buster available in all templates.

    {{ app_version }} - Display version string (e.g. "v2.0.0 - Fusion")
    {{ cache_bust }}  - Query param for static assets, changes every deploy
                        Usage: href="...?v={{ cache_bust }}"
    """
    return dict(app_version=APP_VERSION, cache_bust=STARTUP_ID)

@app.after_request
def add_cache_headers(response):
    """Prevent browsers from caching HTML pages after deployment.

    Static assets use ?v= query params for cache busting.
    HTML responses get no-cache so phones always get fresh pages on reload.
    """
    if 'text/html' in response.content_type:
        response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        response.headers['Pragma'] = 'no-cache'
        response.headers['Expires'] = '0'
    return response

# Polling endpoints that fire every 5s per client - never log these
QUIET_PATHS = frozenset([
    '/api/heartbeat',
    '/api/check-round-status',
    '/api/broadcast-message',
    '/host/codes-status',
    '/host/check-active-round',
    '/host/count-unscored',
    '/host/team-status',
    '/host/get-sleep-status',
])

@app.before_request
def log_request():
    """Log incoming requests - skip static files and high-frequency polling"""
    if request.path.startswith('/static'):
        return
    if request.path in QUIET_PATHS:
        return
    if request.path.startswith('/api/view-status/'):
        return
    code = session.get('code', '-')
    team = session.get('team_name', '-')
    logger.debug(f"[REQUEST] {request.method} {request.path} | code={code} team={team} ip={request.remote_addr}")

def host_required(f):
    """Decorator to protect host routes - requires password authentication"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not session.get('host_authenticated'):
            if request.path in QUIET_PATHS:
                logger.debug(f"[HOST] Auth check FAILED for {request.path} - redirecting to login")
            else:
                logger.info(f"[HOST] Auth check FAILED for {request.path} - redirecting to login")
            return redirect(url_for('host_login'))
        logger.debug(f"[HOST] Auth check passed for {request.path}")
        return f(*args, **kwargs)
    return decorated_function

def team_session_valid(f):
    """Decorator to validate team session - checks startup_id and reset_counter"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        # CRITICAL: Check reset_counter and startup_id BEFORE checking if session exists
        # This ensures Game Over page shows even if session was cleared

        # Use DEBUG for polling endpoints to avoid log noise
        log = logger.debug if request.path in QUIET_PATHS else logger.info

        # Check if startup_id in session matches current server startup
        # If server restarted, STARTUP_ID changes = all old sessions invalid
        session_startup_id = session.get('startup_id')

        if session_startup_id is not None and session_startup_id != STARTUP_ID:
            # Server was restarted - show game over page
            log(f"Team session invalid - server restarted (session startup_id: {session_startup_id}, current: {STARTUP_ID})")
            session.clear()
            return render_template('game_over.html', reason='server_restart')

        # Check if reset_counter matches (Reset All button invalidates sessions)
        session_reset_counter = session.get('reset_counter', 0)

        if session_reset_counter != RESET_COUNTER:
            # Game was reset - show game over page
            log(f"Team session invalid - game was reset (session counter: {session_reset_counter}, current: {RESET_COUNTER})")
            session.clear()
            return render_template('game_over.html', reason='game_reset')

        # NOW check if team has a session (after checking reset/restart)
        if 'code' not in session:
            log("[TEAM] No team session found - redirecting to join")
            return redirect(url_for('join'))

        logger.debug(f"[TEAM] Session valid for code={session.get('code')} team={session.get('team_name')} path={request.path}")
        return f(*args, **kwargs)
    return decorated_function

# Game configuration - 8 rounds
ROUNDS_CONFIG = [
    {"round": 1, "answers": 4},
    {"round": 2, "answers": 5},
    {"round": 3, "answers": 6},
    {"round": 4, "answers": 4},
    {"round": 5, "answers": 5},
    {"round": 6, "answers": 3},
    {"round": 7, "answers": 5},
    {"round": 8, "answers": 4}
]

# Pre-built surveys for quick round creation via dropdown
PREBUILT_SURVEYS = {
    "survey1": {
        "name": "Survey 1",
        "rounds": [
            {"question": "Name Something Parents Warn Their Children Not To Get Their Fingers Caught In", "answers": ["Door", "Fan", "Outlet", "Cookie Jar"], "answer1_count": 45},
            {"question": "A Young Person \u201cFights For Their Right To Party.\u201d What Might An Old Person Fight For The Right To Do?", "answers": ["Sleep", "Vote", "Retire", "Keep License", "Get Social Security"], "answer1_count": 32},
            {"question": "Name Something That Goes Well With Pizza", "answers": ["Beer", "Soda", "Salad", "Breadstick/Knots", "Chicken Wings", "Chips"], "answer1_count": 36},
            {"question": "Name Something From The Laundry That\u2019s Impossible To Fold Neatly", "answers": ["Fitted Sheets", "Socks", "Underwear", "Blouse"], "answer1_count": 35},
            {"question": "Name A Place Where You\u2019d Be Mortified If Your Cell Phone Went Off", "answers": ["Church", "Funeral", "Movie Theater", "Job Interview", "Wedding"], "answer1_count": 39},
            {"question": "Name Something You Should Switch Off Before Going To Bed", "answers": ["Lights", "Phone", "TV"], "answer1_count": 67},
            {"question": "Name A Common Sickness That Kids Seem To Get More Often Than Adults", "answers": ["Cold", "Flu", "Chicken Pox", "Ear Infection", "Strep Throat"], "answer1_count": 32},
            {"question": "Name A Reason Why A Man Would Wax Hair Off Part Of His Body", "answers": ["Too Hairy", "For Spouse/Date", "Body Builder", "Swimmer"], "answer1_count": 34},
        ]
    },
    "survey2": {
        "name": "Survey 2",
        "rounds": [
            {"question": "Name Something Permanent On a Criminal\u2019s Skin That Police Use To Be Sure They\u2019ve Got Their Man", "answers": ["Tattoo", "Fingerprint", "Scar", "Birthmark"], "answer1_count": 41},
            {"question": "What Might Someone Use While Cutting Their Own Hair?", "answers": ["Scissors", "Mirror", "Clippers", "Comb", "Bowl"], "answer1_count": 48},
            {"question": "Name Something Babies And Puppies Have In Common", "answers": ["Cute", "Drooling", "Need Attention", "Playful", "Sleep A lot", "Cry"], "answer1_count": 34},
            {"question": "What Diaper Bag Item Would A Parent Hate To Be Without?", "answers": ["Diapers", "Wipes", "Bottle", "Pacifier"], "answer1_count": 49},
            {"question": "Name Something Twins Might Always Share", "answers": ["Looks", "Parents", "Genes", "Birthday", "Last Name"], "answer1_count": 40},
            {"question": "Something Specific People Do In Front Of Mirror", "answers": ["Apply Makeup", "Check Outfit", "Pose Naked"], "answer1_count": 63},
            {"question": "Name A Type Of Sauce That You\u2019d Never Put On Pasta", "answers": ["Apple Sauce", "Hot Sauce", "Ketchup", "BBQ", "Chocolate"], "answer1_count": 30},
            {"question": "Name Something A Child Does To Prove They\u2019re Too Sick For School", "answers": ["Cough", "Vomit", "Cry", "Take Temperature"], "answer1_count": 50},
        ]
    },
    "survey3": {
        "name": "Survey 3",
        "rounds": [
            {"question": "Name Something You Might Adjust When You Get Into A Rental Car", "answers": ["Seat", "Mirrors", "Seat Belt", "Steering Wheel"], "answer1_count": 58},
            {"question": "Name Something A Woman Should Know A Man Before Marrying Him", "answers": ["Income", "Age", "Does He Have Kids", "His Name", "Past Relationships"], "answer1_count": 39},
            {"question": "Name Something You Need In Order To Make A Garden", "answers": ["Seeds", "Soil", "Water", "Hoe", "Shovel", "Plot of Land"], "answer1_count": 35},
            {"question": "Name A Place Where You Hear People Being Paged Over A Loudspeaker", "answers": ["Hospital", "Airport", "School", "Store"], "answer1_count": 31},
            {"question": "Tell Me Something You Do When You Stay Up Late At Night", "answers": ["Watch TV/Movie", "Read", "Snack", "Drink", "Play Phone/Video Games"], "answer1_count": 58},
            {"question": "Name A Crime That Some People Probably Commit Every Day", "answers": ["Speeding", "Jaywalking", "Littering"], "answer1_count": 62},
            {"question": "Name A Reason Why A Person Might Prefer To Own A Dog Over A Cat", "answers": ["Protection", "Loyalty", "Cat Allergies", "Friendlier", "More fun to play with"], "answer1_count": 54},
            {"question": "Name A Phrase You\u2019d Say To Your Partner That Starts With \u201cYou Drive Me __.\u201d", "answers": ["Crazy/Nuts", "Wild", "Up a Wall", "To Drink"], "answer1_count": 58},
        ]
    },
    "survey4": {
        "name": "Survey 4",
        "rounds": [
            {"question": "We Asked 100 Women: Name A Gift That You\u2019d Always Be Happy To Get From Your Partner", "answers": ["Flowers", "Jewelry", "Money", "Chocolate"], "answer1_count": 43},
            {"question": "Name A Slow-Moving Vehicle That You Hate To Get Stuck Behind", "answers": ["Bus", "Semi-Truck", "Tractor", "Garbage Truck", "Dump Truck"], "answer1_count": 34},
            {"question": "Name A Last Minute Problem That Could Make You Late For Work", "answers": ["Traffic", "Car Trouble", "Lost Keys", "Child is Sick", "No Gas", "Bad Hair"], "answer1_count": 35},
            {"question": "Name Something Parents Warn Their Children Not To Get Their Fingers Caught In", "answers": ["Door", "Fan", "Outlet", "Cookie Jar"], "answer1_count": 45},
            {"question": "Name Something You Spray On Yourself That Would Sting If It Got In Your Eyes", "answers": ["Perfume", "Insect Repellent", "Hairspray", "Sunscreen/Tan", "Deodorant"], "answer1_count": 34},
            {"question": "Name Something You Dunk", "answers": ["Basketball", "Donuts", "Cookies"], "answer1_count": 59},
            {"question": "Name Something You Wear That Covers Your Ears", "answers": ["Earmuffs", "Hat", "Headphones", "Scarf", "Hood"], "answer1_count": 49},
            {"question": "Name Something A Politician Does When Scandalous News Breaks Out About Them", "answers": ["Lie/Deny It", "Go Into Hiding", "Apology/Press Conference", "Resign"], "answer1_count": 49},
        ]
    },
    "survey5": {
        "name": "Survey 5",
        "rounds": [
            {"question": "Name a chore kids try to avoid", "answers": ["Doing the Dishes", "Cleaning Their Room", "Taking Out The Trash", "Mowing The Lawn"], "answer1_count": 43},
            {"question": "What Would You Hear On The Radio That Would Make You Turn The Station?", "answers": ["Commercial", "News", "Bad Song", "Static", "Cursing"], "answer1_count": 34},
            {"question": "Name Something People Do With Both Hands", "answers": ["Drive", "Dishes", "Type on Keyboard", "Clap", "Cook", "Eat"], "answer1_count": 40},
            {"question": "Name A Day Of The Year That Some People Don\u2019t Want To Spend Alone", "answers": ["Christmas", "Valentines Day", "Birthday", "New Years Eve"], "answer1_count": 40},
            {"question": "Name Something You Might Pay Someone To Do While You\u2019re Away On Vacation", "answers": ["Care for Pets", "House Sit", "Water Plants", "Babysit", "Collect Mail"], "answer1_count": 28},
            {"question": "What Do You Find Out About A Town By Reading Signs On The Side Of The Road?", "answers": ["Population", "Town Name", "Speed Limit"], "answer1_count": 60},
            {"question": "Other Than Academics Why Might A Teen Choose A Certain College?", "answers": ["Sports Team", "Location", "Friends are Going", "Party School", "Cost of Tuition"], "answer1_count": 47},
            {"question": "Name Something That\u2019s On Your Dinner Table Every Night That The Dog Won\u2019t Beg For", "answers": ["Veggies/Salad", "Salt", "Silverware/Plates", "Napkins"], "answer1_count": 45},
        ]
    },
    "survey6": {
        "name": "Survey 6",
        "rounds": [
            {"question": "If You\u2019re Driving In The Middle Of No Where, What Animal Might You See Crossing The Street?", "answers": ["Deer", "Cow", "Moose", "Coyote"], "answer1_count": 50},
            {"question": "how Many Hours Of Sleep Does The average person Need In Oder To Wake Up Refreshed?", "answers": ["8", "7", "6", "10", "9"], "answer1_count": 47},
            {"question": "Name Something There Are Seven Of", "answers": ["Dwarfs", "Deadly Sins", "Wonders of the World", "Days Per Week", "Sins", "Continents"], "answer1_count": 28},
            {"question": "Name An Activity That\u2019d Be Hard To Do By Candlelight", "answers": ["Read", "Cook", "Write", "Sewing/Knitting"], "answer1_count": 62},
            {"question": "Name Something That Happens To An Old Person\u2019s Body, That You\u2019d Be Surprised To Hear A teen Complaining About", "answers": ["Wrinkles", "Arthritis", "Gray Hair", "Sagging", "Back Ache"], "answer1_count": 50},
            {"question": "Name something you might find in a kitchen drawer", "answers": ["Beach", "Spa", "Park"], "answer1_count": 51},
            {"question": "Name A Good Place To Put Your Hands While Kissing Someone", "answers": ["Their Face", "Around Their Neck", "Their Hips", "Their Back", "Their Shoulders"], "answer1_count": 27},
            {"question": "Instead Of Their First Name, What Might A Parent Shout When Calling For Their Child?", "answers": ["Whole Name", "Nickname", "Hey!", "Siblings Name"], "answer1_count": 38},
        ]
    },
    "survey7": {
        "name": "Survey 7",
        "rounds": [
            {"question": "Name A Place An Animal Might Take A Bath, But You Never Would", "answers": ["Lake/Pond", "Puddle", "River", "Bird Bath"], "answer1_count": 51},
            {"question": "Name a job title someone might have in a big company", "answers": ["CEO", "President", "Vice President", "Supervisor", "Manager"], "answer1_count": 39},
            {"question": "Name A Job Where It Would Be Okay To Yell At Work", "answers": ["Construction", "Sports", "Teacher", "Police", "Stock Brocker", "Auctioneer"], "answer1_count": 43},
            {"question": "What Are 2 Brothers Most Likely To Fight Over?", "answers": ["Girls", "Toys", "TV Remote", "Attention"], "answer1_count": 45},
            {"question": "Name A Way You Can Tell A Storm Is Coming", "answers": ["Dark Clouds", "Lightning", "Wind Changes", "Smell", "Drizzling"], "answer1_count": 61},
            {"question": "Name Something A Plane Can't Fly Without", "answers": ["Wings", "Fuel", "A Pilot"], "answer1_count": 42},
            {"question": "Tell Me A Reason You Might Be Low On Sleep", "answers": ["Overworked", "Kids/New Baby", "Can\u2019t Sleep", "Sick", "Studying"], "answer1_count": 44},
            {"question": "Name A Color Baby Clothes Comes in", "answers": ["Pink", "Blue", "Yellow", "Green"], "answer1_count": 47},
        ]
    },
    "survey8": {
        "name": "Survey 8",
        "rounds": [
            {"question": "Last thing you\u2019d want to happen at the airport", "answers": ["Miss Flight", "Lose Luggage", "Stopped By Security", "Delayed"], "answer1_count": 0},
            {"question": "Something you do when approached by a salesperson", "answers": ["Avoid Them", "Ask For Help", "Smile", "Say Hi", "Just Looking"], "answer1_count": 0},
            {"question": "Someone you hope never writes a tell-all book", "answers": ["Parent", "Significant Other", "Ex", "Best Friend", "Sibling", "Son/Daughter"], "answer1_count": 0},
            {"question": "Something people check on their smartwatch", "answers": ["Steps", "Notifications", "Heart Rate", "Time"], "answer1_count": 0},
            {"question": "Famous phrase from The Wizard of Oz", "answers": ["Off To See The Wizard", "No Place Like Home", "Follow The Yellow Brick Road", "I\u2019ll Get You My Pretty", "Lions, Tigers & Bears!"], "answer1_count": 0},
            {"question": "Place people stash a spare charging cable", "answers": ["Car", "Work Desk", "Backpack/Purse"], "answer1_count": 0},
            {"question": "Chore that takes less than 10 minutes", "answers": ["Take Out Trash", "Wipe Counter", "Make The Bed", "Load Dishwasher", "Water Plants"], "answer1_count": 0},
            {"question": "Feature people look for in a new phone", "answers": ["Battery Life", "Camera", "Price", "Storage"], "answer1_count": 0},
        ]
    },
    "survey9": {
        "name": "Survey 9",
        "rounds": [
            {"question": "How Many Dates Should You Go On Before Kissing Someone?", "answers": ["2", "3", "1", "5"], "answer1_count": 43},
            {"question": "What Would You Hear On The Radio That Would Make You Turn The Station?", "answers": ["Commercial", "News", "Bad Song", "Static", "Cursing"], "answer1_count": 34},
            {"question": "Name Something People Do With Both Hands", "answers": ["Drive", "Dishes", "Tie Shoelaces", "Clap", "Cook", "Put on a Coat"], "answer1_count": 40},
            {"question": "Name A Day Of The Year That Some People Don\u2019t Want To Spend Alone", "answers": ["Christmas", "Valentines Day", "Birthday", "New Years Eve"], "answer1_count": 40},
            {"question": "Name Something You Might Pay Someone To Do While You\u2019re Away On Vacation", "answers": ["Care for Pets", "House Sit", "Water Plants", "Babysit", "Collect Mail"], "answer1_count": 28},
            {"question": "What Do You Find Out About A Town By Reading Signs On The Side Of The Road?", "answers": ["Population", "Town Name", "Speed Limit"], "answer1_count": 60},
            {"question": "Other Than Academics Why Might A Teen Choose A Certain College?", "answers": ["Sports Team", "Location", "Friends are Going", "Party School", "Cost of Tuition"], "answer1_count": 47},
            {"question": "Name Something That\u2019s On Your Dinner Table Every Night That The Dog Won\u2019t Beg For", "answers": ["Veggies/Salad", "Salt", "Silverware/Plates", "Napkins"], "answer1_count": 45},
        ]
    },
}

def db_connect():
    logger.debug("[DB] Opening database connection")
    conn = sqlite3.connect(DB_PATH, timeout=30)
    conn.row_factory = sqlite3.Row
    # Production SQLite settings for concurrent writes
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA synchronous=NORMAL")
    conn.execute("PRAGMA foreign_keys=ON")
    return conn

def generate_team_code():
    """Generate 4-letter code like BAJK (no numbers for easier mobile typing)"""
    # Only letters that are unambiguous in handwriting AND OCR
    # Removed: C/G (similar), D/O (similar), I/L (identical), O/Q (similar), U/V (similar)
    chars = 'ABEFHJKMNPRSTWXYZ'
    code = ''.join(secrets.choice(chars) for _ in range(4))
    logger.debug(f"[CODES] Generated team code: {code}")
    return code

def load_fixed_codes():
    """Load the fixed team codes from codes.json"""
    codes_file = os.path.join(os.path.dirname(__file__), 'codes.json')
    with open(codes_file, 'r') as f:
        codes = json.load(f)
    return codes

def get_qr_base_url():
    """Get QR code base URL from settings, env vars, or defaults."""
    qr_url_from_env = os.environ.get('QR_BASE_URL')
    if qr_url_from_env:
        default_url = qr_url_from_env
    elif os.environ.get('RENDER'):
        default_url = os.environ.get('RENDER_EXTERNAL_URL', 'https://pubfeud.gamenightguild.net')
    else:
        default_url = 'http://localhost:5000'
    return get_setting('qr_base_url', default_url)

def ensure_fixed_codes():
    """Insert fixed codes from codes.json into the database if they don't exist.
    Also removes any codes NOT in the fixed list (leftover from old random generation).
    """
    codes = load_fixed_codes()
    with db_connect() as conn:
        # Remove codes not in the fixed list
        placeholders = ','.join(['?'] * len(codes))
        conn.execute(f"DELETE FROM team_codes WHERE code NOT IN ({placeholders})", codes)
        # Insert fixed codes if not already present
        for code in codes:
            try:
                conn.execute("INSERT INTO team_codes (code, used) VALUES (?, 0)", (code,))
            except sqlite3.IntegrityError:
                pass  # Already exists
        conn.commit()
    logger.info(f"[CODES] {len(codes)} fixed codes loaded from codes.json")

def init_db():
    with db_connect() as conn:
        conn.execute("""
            CREATE TABLE IF NOT EXISTS team_codes (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                code TEXT UNIQUE NOT NULL,
                used INTEGER DEFAULT 0,
                team_name TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        
        conn.execute("""
            CREATE TABLE IF NOT EXISTS rounds (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                round_number INTEGER NOT NULL,
                question TEXT,
                num_answers INTEGER NOT NULL,
                answer1 TEXT,
                answer1_count INTEGER,
                answer2 TEXT,
                answer2_count INTEGER,
                answer3 TEXT,
                answer3_count INTEGER,
                answer4 TEXT,
                answer4_count INTEGER,
                answer5 TEXT,
                answer5_count INTEGER,
                answer6 TEXT,
                answer6_count INTEGER,
                is_active INTEGER DEFAULT 0,
                submissions_closed INTEGER DEFAULT 0,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        
        conn.execute("""
            CREATE TABLE IF NOT EXISTS submissions (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                code TEXT NOT NULL,
                round_id INTEGER NOT NULL,
                answer1 TEXT,
                answer2 TEXT,
                answer3 TEXT,
                answer4 TEXT,
                answer5 TEXT,
                answer6 TEXT,
                tiebreaker INTEGER,
                score INTEGER DEFAULT 0,
                scored INTEGER DEFAULT 0,
                scored_at TIMESTAMP,
                submitted_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (round_id) REFERENCES rounds (id),
                UNIQUE(code, round_id)
            )
        """)
        
        conn.execute("""
            CREATE TABLE IF NOT EXISTS settings (
                key TEXT PRIMARY KEY,
                value TEXT,
                description TEXT,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        conn.execute("""
            CREATE TABLE IF NOT EXISTS ai_corrections (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                round_id INTEGER NOT NULL,
                submission_id INTEGER NOT NULL,
                question TEXT NOT NULL,
                team_answer TEXT NOT NULL,
                survey_answer TEXT,
                survey_num INTEGER,
                correction_type TEXT NOT NULL,
                ai_reasoning TEXT,
                host_reason TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)

        conn.commit()

        logger.info("Database tables initialized")
        
        # Migration: Add checked_answers column if it doesn't exist
        try:
            conn.execute("SELECT checked_answers FROM submissions LIMIT 1")
        except:
            logger.info("Adding checked_answers column to submissions table...")
            conn.execute("ALTER TABLE submissions ADD COLUMN checked_answers TEXT")
            conn.commit()
            logger.info("Migration complete: checked_answers column added")
        
        # Migration: Add submissions_closed column to rounds table if it doesn't exist
        try:
            conn.execute("SELECT submissions_closed FROM rounds LIMIT 1")
        except:
            logger.info("Adding submissions_closed column to rounds table...")
            conn.execute("ALTER TABLE rounds ADD COLUMN submissions_closed INTEGER DEFAULT 0")
            conn.commit()
            logger.info("Migration complete: submissions_closed column added")
        
        # Migration: Add previous_score column to submissions table if it doesn't exist
        try:
            conn.execute("SELECT previous_score FROM submissions LIMIT 1")
        except:
            logger.info("Adding previous_score column to submissions table...")
            conn.execute("ALTER TABLE submissions ADD COLUMN previous_score INTEGER DEFAULT NULL")
            conn.commit()
            logger.info("Migration complete: previous_score column added")
        
        # Migration: Add last_heartbeat column to team_codes table (v1.1.0)
        try:
            conn.execute("SELECT last_heartbeat FROM team_codes LIMIT 1")
        except:
            logger.info("Adding last_heartbeat column to team_codes table...")
            conn.execute("ALTER TABLE team_codes ADD COLUMN last_heartbeat TIMESTAMP DEFAULT NULL")
            conn.commit()
            logger.info("Migration complete: last_heartbeat column added")
        
        # Migration: Add reconnected column to team_codes table (v1.1.0)
        try:
            conn.execute("SELECT reconnected FROM team_codes LIMIT 1")
        except:
            logger.info("Adding reconnected column to team_codes table...")
            conn.execute("ALTER TABLE team_codes ADD COLUMN reconnected INTEGER DEFAULT 0")
            conn.commit()
            logger.info("Migration complete: reconnected column added")
        
        # Migration: Add winner_code column to rounds table (v1.1.0)
        try:
            conn.execute("SELECT winner_code FROM rounds LIMIT 1")
        except:
            logger.info("Adding winner_code column to rounds table...")
            conn.execute("ALTER TABLE rounds ADD COLUMN winner_code TEXT DEFAULT NULL")
            conn.commit()
            logger.info("Migration complete: winner_code column added")
        
        # Migration: Add host_reason column to ai_corrections table (v2.0.4)
        try:
            conn.execute("SELECT host_reason FROM ai_corrections LIMIT 1")
        except:
            logger.info("Adding host_reason column to ai_corrections table...")
            conn.execute("ALTER TABLE ai_corrections ADD COLUMN host_reason TEXT DEFAULT NULL")
            conn.commit()
            logger.info("Migration complete: host_reason column added")

        # Initialize default settings if they don't exist
        default_settings = [
            ('allow_team_registration', 'true', 'Allow new teams to join'),
            ('system_paused', 'false', 'System pause status'),
            ('broadcast_message', '', 'Broadcast message to all teams'),
            ('server_sleep', 'false', 'Server sleep mode - stops auto-refresh'),
            ('ai_model', '', 'AI model for scoring and photo scan'),
            ('extended_thinking_enabled', 'false', 'Enable extended thinking for AI calls'),
            ('thinking_budget_tokens', '10000', 'Token budget for extended thinking'),
        ]
        
        for key, value, description in default_settings:
            existing = conn.execute("SELECT value FROM settings WHERE key = ?", (key,)).fetchone()
            if not existing:
                conn.execute(
                    "INSERT INTO settings (key, value, description) VALUES (?, ?, ?)",
                    (key, value, description)
                )
        conn.commit()
        
        # Auto-generate 30 codes if table is empty
        codes_exist = conn.execute("SELECT COUNT(*) as cnt FROM team_codes").fetchone()['cnt']
        if codes_exist == 0:
            logger.info("Auto-generating 30 team codes...")
            generated = 0
            attempts = 0
            while generated < 30 and attempts < 100:
                code = generate_team_code()
                try:
                    conn.execute("INSERT INTO team_codes (code, used) VALUES (?, 0)", (code,))
                    conn.commit()
                    generated += 1
                except sqlite3.IntegrityError:
                    attempts += 1
                    continue
            logger.info(f"Generated {generated} team codes")
        else:
            logger.info(f"Database already has {codes_exist} team codes")

# ============= NUCLEAR RESET ON EVERY SERVER START =============
def nuke_all_data():
    """NUCLEAR OPTION: Clear ALL game data on server startup
    
    SERVER RESTART = FRESH START
    - No old teams
    - No old submissions  
    - No old rounds
    - EVERYTHING IS BRAND NEW
    """
    logger.info("[RESET] Nuclear reset - clearing all game data")

    with db_connect() as conn:
        # DELETE EVERYTHING
        conn.execute("DELETE FROM submissions")
        conn.execute("DELETE FROM rounds")
        conn.execute("UPDATE team_codes SET used = 0, team_name = NULL")
        conn.commit()

    logger.info("[RESET] All data cleared - server is fresh, all teams must rejoin")

init_db()
nuke_all_data()  # NUKE EVERYTHING on every server start
ensure_fixed_codes()  # Load fixed codes from codes.json

# ============= SETTINGS HELPERS =============

def get_setting(key, default=None):
    """Get a setting value from database, return default if not found"""
    try:
        with db_connect() as conn:
            result = conn.execute(
                "SELECT value FROM settings WHERE key = ?", 
                (key,)
            ).fetchone()
            value = result['value'] if result else default
            logger.debug(f"[SETTINGS] get_setting('{key}') = '{value}'")
            return value
    except Exception as e:
        logger.warning(f"[SETTINGS] Failed to get setting '{key}': {e}")
        return default

def set_setting(key, value, description=''):
    """Save a setting to database"""
    try:
        with db_connect() as conn:
            conn.execute("""
                INSERT OR REPLACE INTO settings (key, value, description, updated_at)
                VALUES (?, ?, ?, CURRENT_TIMESTAMP)
            """, (key, value, description))
            conn.commit()
            logger.debug(f"[SETTINGS] Setting updated: {key} = {value}")
            return True
    except Exception as e:
        logger.error(f"[SETTINGS] Failed to set setting '{key}': {e}")
        return False

def get_current_ai_model():
    """Get the current AI model to use.
    Priority: database setting > AI_MODEL env var > hardcoded default.
    """
    db_value = get_setting('ai_model', '')
    if db_value:
        valid_ids = [m['id'] for m in AI_MODEL_CHOICES]
        if db_value in valid_ids:
            return db_value
        else:
            logger.warning(f"[AI] Unknown model in database: '{db_value}', falling back to default")
    return AI_MODEL_DEFAULT

def build_claude_api_kwargs(max_tokens_default):
    """Build keyword arguments for client.messages.create() based on current settings.

    When extended thinking is enabled, removes temperature and adds thinking parameter.
    When disabled, uses temperature=0.
    """
    thinking_enabled = get_setting('extended_thinking_enabled', 'false') == 'true'

    if thinking_enabled:
        budget = int(get_setting('thinking_budget_tokens', '10000'))
        budget = max(budget, 1024)
        effective_max_tokens = budget + max_tokens_default
        return {
            'max_tokens': effective_max_tokens,
            'thinking': {
                'type': 'enabled',
                'budget_tokens': budget,
            },
        }
    else:
        return {
            'max_tokens': max_tokens_default,
            'temperature': 0,
        }

def extract_response_text(message):
    """Extract the text content from a Claude API response.

    When extended thinking is enabled, message.content contains a thinking block
    followed by a text block. This finds the text block regardless.
    """
    for block in message.content:
        if block.type == 'text':
            return block.text
    return message.content[0].text

# Anthropic SDK requires streaming when max_tokens exceeds this threshold
# to avoid HTTP timeouts on long-running extended thinking requests.
STREAMING_THRESHOLD = 21333

def call_claude_api(client, model, messages, api_kwargs):
    """Call Claude API, using streaming when max_tokens exceeds SDK threshold.

    When extended thinking is enabled and max_tokens > 21333, the Anthropic SDK
    requires streaming to avoid HTTP timeouts. This helper automatically switches
    to streaming in that case, returning the same Message object either way.
    """
    use_streaming = (
        'thinking' in api_kwargs
        and api_kwargs.get('max_tokens', 0) > STREAMING_THRESHOLD
    )

    if use_streaming:
        logger.debug(f"[AI] Using streaming (max_tokens={api_kwargs['max_tokens']} > {STREAMING_THRESHOLD})")
        with client.messages.stream(
            model=model,
            messages=messages,
            **api_kwargs
        ) as stream:
            return stream.get_final_message()
    else:
        return client.messages.create(
            model=model,
            messages=messages,
            **api_kwargs
        )

# ============= HELPERS =============

def similar(a, b):
    """Check if answers are similar (for auto-checking)"""
    if not a or not b:
        return False
    a_clean = a.lower().strip()
    b_clean = b.lower().strip()
    if a_clean == b_clean:
        logger.debug(f"[SCORING] similar() exact match: '{a}' == '{b}'")
        return True
    ratio = SequenceMatcher(None, a_clean, b_clean).ratio()
    if ratio > 0.9:
        logger.debug(f"[SCORING] similar() fuzzy match: '{a}' ~ '{b}' (ratio={ratio:.3f})")
        return True
    return False

PHOTO_SCAN_PROMPT = """You are extracting handwritten answers from a Family Feud paper answer sheet.

The page contains up to 4 team answer blocks arranged in a 2x2 grid. Each block has this layout:

LAYOUT OF EACH BLOCK:
- "Team Name:" label on the left, followed by a handwritten team name on the line
- A 4-LETTER CODE (like "ABAR", "HJNK", "XMPR") is written separately in the TOP RIGHT CORNER of the block, AWAY from the team name. The code is NOT part of the team name — it is a separate identifier. It is always exactly 4 uppercase letters with no numbers.
- "Answer 1:" through "Answer 6:" — handwritten answers on labeled lines
- "Tie Breaker #" — a number (typically 0-100)

CRITICAL: The 4-letter code and the team name are TWO SEPARATE THINGS. The code is in the top-right corner of the block. The team name is on the "Team Name:" line. Do NOT combine them. For example, if you see "Tina" written after "Team Name:" and "ABAR" written in the corner, the team_name is "Tina" and the code is "ABAR".

Extract ALL team blocks visible on the page that have at least a team name filled in. Skip completely blank blocks.

Rules:
- The code is ALWAYS exactly 4 uppercase letters (A-Z). No numbers, no spaces.
- The code uses only these letters: A B E F H J K M N P R S T W X Y Z
- Read handwriting as accurately as possible, even if messy
- If a field is blank/empty, use an empty string ""
- The tiebreaker should be an integer. If unclear or blank, use 0
- Team names may be creative/unusual — transcribe exactly what is written
- Answers may contain multiple words, abbreviations, or slang — transcribe as-is
- If you cannot find the 4-letter code, use "" but look carefully in the top-right area first
- List any fields where you are NOT confident in the "low_confidence_fields" array

Respond with ONLY valid JSON in this exact format (no markdown, no explanation):
{
  "teams": [
    {
      "code": "ABAR",
      "team_name": "Tina",
      "answers": ["chicken", "pizza", "broccoli", "", "", ""],
      "tiebreaker": 42,
      "low_confidence_fields": ["answers.2"]
    }
  ]
}

Always return exactly 6 entries in the answers array per team (use "" for blank ones).
For low_confidence_fields, use: "code", "team_name", "tiebreaker", or "answers.0" through "answers.5"."""


def extract_answers_from_photo(image_b64):
    """
    Use Claude Vision API to extract handwritten answers from a photo of a paper answer sheet.

    Args:
        image_b64: Base64-encoded JPEG image string (no data URI prefix)

    Returns:
        List of dicts with keys: code, team_name, answers (list of 6 strings), tiebreaker (int), low_confidence_fields (list)
    """
    if not ANTHROPIC_AVAILABLE or not ANTHROPIC_API_KEY:
        logger.error("[PHOTO-SCAN] extract_answers_from_photo() called but AI not available")
        return []

    try:
        current_model = get_current_ai_model()
        logger.info(f"[PHOTO-SCAN] Calling Claude Vision API (model: {current_model}, image size: {len(image_b64)} chars base64)")

        client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

        api_kwargs = build_claude_api_kwargs(max_tokens_default=2048)
        logger.info(f"[PHOTO-SCAN] Extended thinking: {'ON' if 'thinking' in api_kwargs else 'OFF'}")

        message = call_claude_api(
            client=client,
            model=current_model,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": "image/jpeg",
                            "data": image_b64
                        }
                    },
                    {
                        "type": "text",
                        "text": PHOTO_SCAN_PROMPT
                    }
                ]
            }],
            api_kwargs=api_kwargs
        )

        response_text = extract_response_text(message)
        logger.info(f"[PHOTO-SCAN] Claude Vision response: {response_text[:500]}")

        # Parse JSON response - same fallback pattern as score_with_ai()
        response_json = None
        try:
            response_json = json.loads(response_text)
        except json.JSONDecodeError:
            brace_start = response_text.find('{')
            brace_end = response_text.rfind('}')
            if brace_start != -1 and brace_end != -1:
                try:
                    response_json = json.loads(response_text[brace_start:brace_end + 1])
                except json.JSONDecodeError:
                    pass

        if response_json and 'teams' in response_json:
            teams = response_json['teams']
            # Validate and normalize each team
            for team in teams:
                team.setdefault('code', '')
                team.setdefault('team_name', '')
                team.setdefault('tiebreaker', 0)
                team.setdefault('low_confidence_fields', [])
                # Ensure exactly 6 answers
                answers = team.get('answers', [])
                while len(answers) < 6:
                    answers.append('')
                team['answers'] = answers[:6]
                # Ensure tiebreaker is int
                try:
                    team['tiebreaker'] = int(team['tiebreaker'])
                except (ValueError, TypeError):
                    team['tiebreaker'] = 0

            logger.info(f"[PHOTO-SCAN] Extracted {len(teams)} teams from photo")
            return teams
        else:
            logger.warning(f"[PHOTO-SCAN] Could not parse teams from response")
            return []

    except Exception as e:
        logger.error(f"[PHOTO-SCAN] Claude Vision API call failed: {e}", exc_info=True)
        raise


def score_with_ai(question, survey_answers, team_answers):
    """
    Use Claude AI to determine semantic matches between team answers and survey answers.

    Args:
        question: The Family Feud question text
        survey_answers: List of dicts with 'number', 'text', 'points' keys
        team_answers: List of strings (team's submitted answers)

    Returns:
        Dict with 'matches' (list of ints) and 'reasoning' (list of dicts)
    """
    if not ANTHROPIC_AVAILABLE or not ANTHROPIC_API_KEY:
        logger.error("[AI-SCORING] score_with_ai() called but AI not available")
        return {'matches': [], 'reasoning': []}

    # Build the prompt
    prompt = f"""You are scoring a Family Feud game. Determine which survey answers semantically match the team's submitted answers.

Question: "{question}"

Survey Answers (the correct answers from the survey):
"""
    for ans in survey_answers:
        prompt += f"{ans['number']}. {ans['text']} ({ans['points']} points)\n"

    prompt += "\nTeam's Submitted Answers:\n"
    for ans in team_answers:
        prompt += f"- {ans}\n"

    # === Fetch past corrections for long-term training ===
    # Load from persistent JSON file (survives deploys and DB resets)
    all_corrections = load_corrections_history()

    # Prioritize: same question first, then all others (most recent last)
    same_q = [c for c in all_corrections if c.get('question') == question]
    other_q = [c for c in all_corrections if c.get('question') != question]
    # Take up to 10 same-question + fill remaining with others, max 30 total
    recent_corrections = same_q[-10:] + other_q[-20:]
    if recent_corrections:
        logger.debug(f"[AI-SCORING] Loaded {len(recent_corrections)} corrections for training ({len(same_q)} same-question)")

    if recent_corrections:
        prompt += "\nPast Corrections (learn from these host overrides — apply similar logic to current answers):\n"
        for idx, corr in enumerate(recent_corrections, 1):
            if corr['correction_type'] == 'host_added':
                prompt += f'{idx}. SHOULD match: "{corr["team_answer"]}" matches "{corr["survey_answer"]}"'
            else:
                prompt += f'{idx}. Should NOT match: "{corr["team_answer"]}" does NOT match "{corr["survey_answer"]}"'
            # Prioritize host's explanation over AI's original reasoning
            if corr.get('host_reason'):
                prompt += f' — Host says: "{corr["host_reason"]}"'
            elif corr.get('ai_reasoning'):
                prompt += f' (you previously thought: {corr["ai_reasoning"]})'
            prompt += '\n'
        prompt += '\n'

    prompt += """
Matching Rules:
- Exact matches count (e.g., "car" matches "car")
- Synonyms count (e.g., "automobile" matches "car")
- Common abbreviations count (e.g., "bike" matches "bicycle")
- Specific types count (e.g., "minivan" matches "van")
- Plurals/singulars are the same (e.g., "dogs" matches "dog")
- Minor misspellings count if intent is clear
- DO NOT match if the meaning is different
- DO NOT match partial words that change meaning

Respond with ONLY a JSON object in this exact format:
{
  "matches": [1, 3, 5],
  "reasoning": [
    {"team_answer": "car", "matched_to": 1, "survey_answer": "Automobile", "why": "Car is a common synonym for automobile"},
    {"team_answer": "food", "matched_to": null, "survey_answer": null, "why": "Too vague, no survey answer about food"}
  ]
}

"matches" = list of survey answer numbers that have a semantic match in the team's answers.
"reasoning" = one entry per team answer, in the order they were submitted:
  - "team_answer" = the team's submitted text
  - "matched_to" = the survey answer number (integer) it matches, or null if no match
  - "survey_answer" = the text of the matched survey answer, or null if no match
  - "why" = one short sentence explaining the decision

If no matches at all, return: {"matches": [], "reasoning": [...]}"""

    try:
        current_model = get_current_ai_model()
        logger.debug(f"[AI-SCORING] Calling Claude API (model: {current_model}, prompt length: {len(prompt)} chars)")

        client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

        api_kwargs = build_claude_api_kwargs(max_tokens_default=1024)
        logger.debug(f"[AI-SCORING] Extended thinking: {'ON' if 'thinking' in api_kwargs else 'OFF'}")

        message = call_claude_api(
            client=client,
            model=current_model,
            messages=[
                {"role": "user", "content": prompt}
            ],
            api_kwargs=api_kwargs
        )

        response_text = extract_response_text(message)
        logger.debug(f"[AI-SCORING] Claude response: {response_text}")

        # Parse JSON response - try full parse first, then regex fallback
        response_json = None
        try:
            response_json = json.loads(response_text)
        except json.JSONDecodeError:
            # Try to extract JSON object from response text
            # Find the outermost { ... } block
            brace_start = response_text.find('{')
            brace_end = response_text.rfind('}')
            if brace_start != -1 and brace_end != -1:
                try:
                    response_json = json.loads(response_text[brace_start:brace_end + 1])
                except json.JSONDecodeError:
                    pass

        if response_json and 'matches' in response_json:
            matches = response_json.get('matches', [])
            reasoning = response_json.get('reasoning', [])

            # Validate matches are within valid range
            max_num = max(a['number'] for a in survey_answers) if survey_answers else 0
            valid_matches = [m for m in matches if isinstance(m, int) and 1 <= m <= max_num]

            logger.info(f"[AI-SCORING] Parsed {len(valid_matches)} valid matches: {valid_matches}, {len(reasoning)} reasoning entries")
            return {'matches': valid_matches, 'reasoning': reasoning}
        else:
            logger.warning(f"[AI-SCORING] Could not parse JSON from response: {response_text}")
            return {'matches': [], 'reasoning': []}

    except Exception as e:
        logger.error(f"[AI-SCORING] Claude API call failed: {e}", exc_info=True)
        raise

def time_ago(timestamp_str):
    """Convert timestamp to 'X minutes ago' format"""
    if not timestamp_str:
        return "just now"
    try:
        dt = datetime.strptime(timestamp_str, '%Y-%m-%d %H:%M:%S')
        diff = datetime.now() - dt
        minutes = int(diff.total_seconds() / 60)
        if minutes < 1:
            return "just now"
        elif minutes == 1:
            return "1 minute ago"
        elif minutes < 60:
            return f"{minutes} minutes ago"
        else:
            hours = minutes // 60
            return f"{hours} hour{'s' if hours > 1 else ''} ago"
    except:
        return "recently"

def format_timestamp(timestamp_str):
    """Format timestamp as '7:42:15 PM' for display"""
    if not timestamp_str:
        return ""
    try:
        dt = datetime.strptime(timestamp_str, '%Y-%m-%d %H:%M:%S')
        return dt.strftime('%I:%M:%S %p')  # e.g., "07:42:15 PM"
    except:
        return ""

# ============= HOST ROUTES =============

@app.route('/')
def index():
    return redirect(url_for('join'))

@app.route('/host/login', methods=['GET', 'POST'])
def host_login():
    """Host login page - password authentication"""
    if request.method == 'POST':
        password = request.form.get('password', '')
        if password == HOST_PASSWORD:
            session['host_authenticated'] = True
            logger.info("Host authenticated successfully")
            # On mobile, go straight to photo scan (if AI enabled)
            if AI_SCORING_ENABLED:
                ua = request.headers.get('User-Agent', '').lower()
                if any(m in ua for m in ['iphone', 'android', 'mobile']):
                    logger.info("[HOST] Mobile login — redirecting to photo scan")
                    return redirect(url_for('photo_scan'))
            return redirect(url_for('host_dashboard'))
        else:
            logger.warning("Failed host login attempt")
            return render_template('host_login.html', error=True)
    return render_template('host_login.html', error=False)

@app.route('/host/logout')
def host_logout():
    """Logout from host panel"""
    session.pop('host_authenticated', None)
    logger.info("Host logged out")
    return redirect(url_for('host_login'))

@app.route('/host')
@host_required
def host_dashboard():
    """Main host dashboard"""
    logger.debug("[HOST] host_dashboard() - loading dashboard")
    with db_connect() as conn:
        codes_raw = conn.execute("""
            SELECT code, used, team_name, reconnected, last_heartbeat
            FROM team_codes
            ORDER BY id ASC
        """).fetchall()
        
        # Process codes to add active status
        codes = []
        for code in codes_raw:
            code_dict = dict(code)
            # Calculate if team is active (heartbeat within last 30 seconds)
            if code['last_heartbeat']:
                from datetime import datetime
                try:
                    # Parse timestamp
                    last_hb = datetime.fromisoformat(code['last_heartbeat'])
                    now = datetime.now()
                    time_diff = (now - last_hb).total_seconds()
                    code_dict['is_active'] = time_diff < 30
                except:
                    code_dict['is_active'] = False
            else:
                code_dict['is_active'] = False
            codes.append(code_dict)
        
        rounds = conn.execute("SELECT * FROM rounds ORDER BY round_number ASC").fetchall()
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()
        
        # Count unscored submissions for active round
        unscored_count = 0
        submission_count = 0
        if active_round:
            unscored_count = conn.execute("""
                SELECT COUNT(*) as cnt FROM submissions 
                WHERE round_id = ? AND scored = 0
            """, (active_round['id'],)).fetchone()['cnt']
            
            # Total submissions for active round
            submission_count = conn.execute("""
                SELECT COUNT(*) as cnt FROM submissions 
                WHERE round_id = ?
            """, (active_round['id'],)).fetchone()['cnt']
    
    logger.debug(f"[HOST] host_dashboard() - {len(codes)} codes, {len(rounds)} rounds, "
                 f"active_round={'R'+str(active_round['round_number']) if active_round else 'None'}, "
                 f"submissions={submission_count}, unscored={unscored_count}")
    return render_template('host.html',
                         codes=codes,
                         rounds=[dict(r) for r in rounds],
                         active_round=dict(active_round) if active_round else None,
                         unscored_count=unscored_count,
                         submission_count=submission_count,
                         rounds_config=ROUNDS_CONFIG,
                         ai_scoring_available=AI_SCORING_ENABLED)

@app.route('/host/codes-status')
@host_required
def codes_status():
    """API endpoint - returns code statuses as JSON for auto-refresh"""
    logger.debug("[CODES] codes_status() called")
    with db_connect() as conn:
        codes = conn.execute("""
            SELECT code, used, team_name
            FROM team_codes
            ORDER BY id ASC
        """).fetchall()
        
        codes_data = []
        for code in codes:
            codes_data.append({
                'code': code['code'],
                'used': bool(code['used']),
                'team_name': code['team_name'] if code['team_name'] else None
            })
        
        used_count = sum(1 for c in codes_data if c['used'])
        
        logger.debug(f"[CODES] codes_status() returning {len(codes_data)} total, {used_count} used")
        
        return jsonify({
            'codes': codes_data,
            'total': len(codes_data),
            'used': used_count
        })

@app.route('/host/generate-codes', methods=['POST'])
@host_required
def generate_codes():
    """Reload fixed team codes from codes.json"""
    logger.info("[CODES] generate_codes() - reloading fixed codes from codes.json")
    ensure_fixed_codes()
    codes = load_fixed_codes()
    logger.info(f"[CODES] generate_codes() - {len(codes)} fixed codes loaded")
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
        <title>Codes Generated</title>
        <style>
            body {{ 
                font-family: Arial; 
                background: linear-gradient(135deg, #1e3c72 0%, #2a5298 100%);
                color: #ffd700;
                text-align: center;
                padding: 50px;
            }}
            .box {{
                background: #000;
                border: 3px solid #ffd700;
                padding: 40px;
                border-radius: 15px;
                max-width: 500px;
                margin: 0 auto;
            }}
            h1 {{ font-size: 3em; margin-bottom: 20px; }}
            button {{
                background: #ffd700;
                color: #000;
                border: none;
                padding: 20px 40px;
                font-size: 1.5em;
                font-weight: bold;
                border-radius: 10px;
                cursor: pointer;
                margin-top: 20px;
            }}
            button:hover {{ background: #ffed4e; }}
        </style>
    </head>
    <body>
        <div class="box">
            <h1>✅ Success!</h1>
            <p style="font-size: 1.5em;">{len(codes)} fixed team codes loaded!</p>
            <button onclick="window.location.href='/host'">Back to Dashboard</button>
        </div>
    </body>
    </html>
    """

@app.route('/host/reclaim-code/<code>', methods=['POST'])
@host_required
def reclaim_code(code):
    """Reclaim a used code - removes team and frees code for reuse"""
    code = code.upper()
    logger.debug(f"[CODES] reclaim_code() - attempting to reclaim code={code}")

    with db_connect() as conn:
        code_row = conn.execute("SELECT * FROM team_codes WHERE code = ?", (code,)).fetchone()

        if not code_row:
            logger.warning(f"[CODES] reclaim_code() - code={code} not found")
            return jsonify({"success": False, "message": "Code not found"}), 404

        if not code_row['used']:
            logger.warning(f"[CODES] reclaim_code() - code={code} is not in use")
            return jsonify({"success": False, "message": "Code is not in use"}), 400
        
        team_name = code_row['team_name']
        
        # Delete all submissions for this team
        conn.execute("DELETE FROM submissions WHERE code = ?", (code,))
        
        # Reset the code (mark as unused, clear team name and reconnect flag)
        conn.execute("""
            UPDATE team_codes 
            SET used = 0, team_name = NULL, reconnected = 0, last_heartbeat = NULL 
            WHERE code = ?
        """, (code,))
        
        conn.commit()
        
        logger.info(f"[CODES] reclaim_code() - code={code} reclaimed from team='{team_name}', submissions deleted")
        
        return jsonify({
            "success": True, 
            "message": f"Code {code} reclaimed. Team '{team_name}' removed."
        })

@app.route('/host/print-codes')
@host_required
def print_codes():
    """Generate landscape HTML page with QR codes for mobile play (replaces paper)"""
    logger.debug("[CODES] print_codes() - generating mobile play QR code page")
    codes = load_fixed_codes()
    if not codes:
        return "No codes available. Generate codes first!", 400
    qr_base_url = get_qr_base_url()
    return render_template('print_qr_codes.html', codes=codes, qr_base_url=qr_base_url,
                           mode='play', title='Mobile Play — Scan to Play on Your Phone')

@app.route('/host/print-codes-landscape')
@host_required
def print_codes_landscape():
    """Generate landscape HTML page with QR codes for view-only status (companion to paper)"""
    logger.debug("[CODES] print_codes_landscape() - generating view-only QR code page")
    codes = load_fixed_codes()
    if not codes:
        return "No codes available. Generate codes first!", 400
    qr_base_url = get_qr_base_url()
    return render_template('print_qr_codes.html', codes=codes, qr_base_url=qr_base_url,
                           mode='view', title='View Only — See Your Submitted Answers')

@app.route('/host/print-answer-sheets')
@host_required
def print_answer_sheets():
    """Generate printable answer sheets with pre-printed codes.
    Accepts ?group=1 (codes 1-30) or ?group=2 (codes 31-60).
    """
    all_codes = load_fixed_codes()
    group = request.args.get('group', '1')
    if group == '2':
        codes = all_codes[30:60]
        group_label = 'Group 2 (31-60)'
    else:
        codes = all_codes[0:30]
        group_label = 'Group 1 (1-30)'
    logger.info(f"[CODES] print_answer_sheets() - generating {group_label} ({len(codes)} codes)")
    qr_base_url = get_qr_base_url()

    return render_template('print_answer_sheets.html', codes=codes, group_label=group_label, rounds_config=ROUNDS_CONFIG, qr_base_url=qr_base_url)

def parse_pptx(filepath):
    """Parse PowerPoint file and extract questions/answers

    IMPROVED VERSION - Handles text boxes with answer/count pairs
    Correctly distinguishes rank indicators (1,2,3) from answer counts (10,20,43)
    """
    from pptx import Presentation
    
    prs = Presentation(filepath)
    slides = list(prs.slides)
    
    rounds_data = []
    
    # Strategy: Find question slides, then parse the next slide as answers
    i = 0
    while i < len(slides):
        slide = slides[i]
        
        # Extract all text from current slide
        all_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    all_text.append(text)
        
        # Check if this is a question slide
        is_question_slide = False
        question_text = ""
        
        for text in all_text:
            # Look for question markers
            if 'Survey Has' in text and 'Responses' in text:
                is_question_slide = True
            # Extract actual question (not the metadata)
            elif len(text) > 10 and 'Round #' not in text and 'Survey Has' not in text:
                if not question_text:  # Take first substantial text
                    question_text = text
        
        if is_question_slide and i + 1 < len(slides):
            # Next slide should have answers
            answer_slide = slides[i + 1]
            answers = []
            
            # Extract all text from answer slide
            answer_text_elements = []
            for shape in answer_slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        answer_text_elements.append(text)
            
            # Parse answer/count pairs from text
            j = 0
            while j < len(answer_text_elements):
                text = answer_text_elements[j]
                
                # Skip UI elements
                skip_keywords = ['Round:', 'ROUND', 'Score Multiplier:', 'BACK TO SCORES', 
                               'NEXT ROUND', 'And The Survey Says', 'X', '«', '»',
                               'type only', '(type', 'Click', 'Press']
                
                if any(keyword in text for keyword in skip_keywords):
                    j += 1
                    continue
                
                # Skip rank indicators (single-digit numbers 1-8)
                if text.isdigit() and len(text) <= 2 and int(text) <= 8:
                    j += 1
                    continue
                
                # If it's text (potential answer), look for count
                if not text.isdigit() and len(text) > 1:
                    answer_text = text
                    count = 0
                    
                    # Look ahead for count
                    if j + 1 < len(answer_text_elements):
                        next_text = answer_text_elements[j + 1]
                        if next_text.isdigit():
                            try:
                                count_value = int(next_text)
                                # Count numbers are typically 5+ (answer counts, not ranks)
                                # But also accept low counts (some answers might have count of 1-4)
                                # The key is: if we just read answer text, next number IS the count
                                if count_value > 0:
                                    answers.append({'answer': answer_text, 'count': count_value})
                                    j += 1  # Skip the count
                            except ValueError:
                                pass
                
                j += 1
            
            # Add round if we found answers
            if answers:
                rounds_data.append({
                    'question': question_text,
                    'answers': answers
                })
            
            # Skip the answer slide
            i += 2
        else:
            i += 1
    
    return rounds_data

@app.route('/host/upload-answers', methods=['POST'])
@host_required
def upload_answers():
    """Upload DOCX or PPTX answer sheet and auto-create all rounds"""
    logger.info("[UPLOAD] upload_answers() - file upload started")
    try:
        if 'file' not in request.files:
            logger.warning("[UPLOAD] No file in request")
            flash('No file uploaded!', 'error')
            return redirect(url_for('host_dashboard'))

        file = request.files['file']
        if file.filename == '':
            logger.warning("[UPLOAD] No file selected")
            flash('No file selected!', 'error')
            return redirect(url_for('host_dashboard'))

        # Accept .docx, .pptx, and .pptm files
        file_ext = os.path.splitext(file.filename)[1].lower()
        logger.info(f"[UPLOAD] File received: '{file.filename}', type={file_ext}")
        if file_ext not in ['.docx', '.pptx', '.pptm']:
            logger.warning(f"[UPLOAD] Invalid file type: {file_ext}")
            flash('Please upload a .docx, .pptx, or .pptm file!', 'error')
            return redirect(url_for('host_dashboard'))
        
        # Save temp file
        temp_path = os.path.join(BASE_DIR, f'temp_answers_{int(time.time())}{file_ext}')
        file.save(temp_path)
        
        # Parse based on file type
        rounds_data = []
        
        if file_ext == '.docx':
            # DOCX parsing (existing code)
            from docx import Document
            doc = Document(temp_path)
            
            # ROBUST: Extract questions with flexible matching (handles both - and – dashes)
            questions = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text and len(text) > 0 and text[0].isdigit():
                    # Match both regular dash (-) and em-dash (–)
                    if '-' in text or '–' in text:
                        # Split on either dash type
                        separator = '–' if '–' in text else '-'
                        parts = text.split(separator, 1)
                        if len(parts) > 1:
                            question = parts[1].strip()
                            questions.append(question)
            
            # ROBUST: Parse ALL 8 tables regardless of question count
            for table_idx, table in enumerate(doc.tables):
                if table_idx >= 8:
                    break
                    
                answers = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if len(cells) >= 3:
                        # Skip header rows - only process if first cell is a number (rank)
                        if not cells[0] or not cells[0].strip().isdigit():
                            continue
                        
                        answer = cells[1]
                        points_count = cells[2]
                        
                        # ROBUST: Flexible count parsing (handles various spacing)
                        count = 0
                        if points_count:
                            # Try both dash types
                            for separator in ['-', '–']:
                                if separator in points_count:
                                    parts = points_count.split(separator)
                                    if len(parts) > 1:
                                        try:
                                            # Extract just the digits from the second part
                                            count_str = ''.join(filter(str.isdigit, parts[1]))
                                            if count_str:
                                                count = int(count_str)
                                            break
                                        except ValueError:
                                            count = 0
                        
                        answers.append({'answer': answer, 'count': count})
                
                # Use question by index, or empty string if not found
                question = questions[table_idx] if table_idx < len(questions) else ''
                
                rounds_data.append({
                    'question': question,
                    'answers': answers
                })
        
        elif file_ext in ['.pptx', '.pptm']:
            # PowerPoint parsing (new code)
            rounds_data = parse_pptx(temp_path)
        
        # Always create all rounds found (should be 8)
        with db_connect() as conn:
            conn.execute("DELETE FROM rounds")
            conn.execute("DELETE FROM submissions")
            
            for idx, round_data in enumerate(rounds_data):
                round_num = idx + 1
                config = ROUNDS_CONFIG[idx]
                num_answers = config['answers']
                
                fields = ['round_number', 'question', 'num_answers', 'is_active']
                values = [round_num, round_data['question'], num_answers, 0]
                
                for i in range(1, num_answers + 1):
                    if i <= len(round_data['answers']):
                        fields.append(f'answer{i}')
                        fields.append(f'answer{i}_count')
                        values.append(round_data['answers'][i-1]['answer'])
                        values.append(round_data['answers'][i-1]['count'])
                
                placeholders = ','.join(['?'] * len(values))
                conn.execute(f"INSERT INTO rounds ({','.join(fields)}) VALUES ({placeholders})", values)
            
            conn.commit()
        
        if os.path.exists(temp_path):
            os.remove(temp_path)
        
        rounds_created = len(rounds_data)
        logger.info(f"[UPLOAD] Complete: {rounds_created} rounds created from '{file.filename}'")
        for idx, rd in enumerate(rounds_data):
            logger.debug(f"[UPLOAD]   Round {idx+1}: Q='{rd['question'][:60]}', {len(rd['answers'])} answers")
        flash(f'✅ Success! {rounds_created} rounds created!', 'success')
        return redirect(url_for('host_dashboard'))
        
    except FileNotFoundError as e:
        logger.error(f"[UPLOAD] FileNotFoundError: {e}")
        try:
            if 'temp_path' in locals() and os.path.exists(temp_path):
                os.remove(temp_path)
        except:
            pass
        flash(f'❌ File error: Could not read the uploaded file. Please try again.', 'error')
        return redirect(url_for('host_dashboard'))
    except ImportError as e:
        logger.error(f"[UPLOAD] ImportError (missing library): {e}")
        try:
            if 'temp_path' in locals() and os.path.exists(temp_path):
                os.remove(temp_path)
        except:
            pass
        flash(f'❌ Missing library: {str(e)}. Please install required dependencies.', 'error')
        return redirect(url_for('host_dashboard'))
    except Exception as e:
        logger.error(f"[UPLOAD] Unexpected error: {type(e).__name__}: {e}")
        import traceback
        logger.error(f"[UPLOAD] Traceback:\n{traceback.format_exc()}")
        try:
            if 'temp_path' in locals() and os.path.exists(temp_path):
                os.remove(temp_path)
        except:
            pass
        
        # Provide helpful error messages based on the error type
        error_msg = str(e)
        if 'pptx' in error_msg.lower() or 'presentation' in error_msg.lower():
            flash(f'❌ PowerPoint parsing error: The file format may be corrupted or unsupported. Details: {error_msg}', 'error')
        elif 'docx' in error_msg.lower() or 'document' in error_msg.lower():
            flash(f'❌ Word document parsing error: The file format may be corrupted. Details: {error_msg}', 'error')
        elif 'table' in error_msg.lower():
            flash(f'❌ Table parsing error: Could not read answer tables. Make sure your file has the correct format. Details: {error_msg}', 'error')
        else:
            flash(f'❌ Upload failed: {error_msg}', 'error')
        
        return redirect(url_for('host_dashboard'))

@app.route('/host/round/create', methods=['POST'])
@host_required
def create_round():
    """Create a round manually"""
    round_num = int(request.form.get('round_number'))
    question = request.form.get('question', '').strip()
    logger.info(f"[ROUND] create_round() - round_num={round_num}, question='{question[:50]}'")

    config = next((r for r in ROUNDS_CONFIG if r['round'] == round_num), None)
    if not config:
        logger.warning(f"[ROUND] create_round() - invalid round number: {round_num}")
        return "Invalid round number", 400

    with db_connect() as conn:
        conn.execute("UPDATE rounds SET is_active = 0")
        conn.execute("""
            INSERT INTO rounds (round_number, question, num_answers, is_active)
            VALUES (?, ?, ?, 1)
        """, (round_num, question, config['answers']))
        conn.commit()
    logger.info(f"[ROUND] create_round() - round {round_num} created and activated")
    return redirect(url_for('host_dashboard'))

@app.route('/host/round/<int:round_id>/activate', methods=['POST'])
@host_required
def activate_round(round_id):
    """Activate a specific round"""
    logger.info(f"[ROUND] activate_round() - requesting activation of round_id={round_id}")
    with db_connect() as conn:
        # CRITICAL FIX: Validate that round has answers before activating
        round_data = conn.execute(
            "SELECT answer1, question FROM rounds WHERE id = ?",
            (round_id,)
        ).fetchone()

        if not round_data:
            logger.warning(f"[ROUND] activate_round() - round_id={round_id} not found")
            flash('❌ Round not found!', 'error')
            return redirect(url_for('host_dashboard'))

        if not round_data['answer1']:
            logger.warning(f"[ROUND] activate_round() - round_id={round_id} has no answers, blocking activation")
            flash('❌ Cannot activate round without answers! Please set answers first.', 'error')
            return redirect(url_for('host_dashboard'))
        
        # CRITICAL FIX: Use transaction to prevent race conditions
        # Deactivate ALL rounds, then activate the selected one atomically
        conn.execute("BEGIN IMMEDIATE")  # Lock database to prevent race conditions
        try:
            conn.execute("UPDATE rounds SET is_active = 0")
            conn.execute("UPDATE rounds SET is_active = 1 WHERE id = ?", (round_id,))
            conn.commit()
            logger.info(f"[ROUND] activate_round() - round_id={round_id} now active (deactivated all others)")
            flash(f'✅ Round activated: {round_data["question"]}', 'success')
        except Exception as e:
            conn.rollback()
            logger.error(f"[ROUND] activate_round() - error: {e}")
            flash(f'❌ Error activating round: {str(e)}', 'error')
    
    return redirect(url_for('host_dashboard'))

@app.route('/host/round/<int:round_id>/answers', methods=['POST'])
@host_required
def set_answers(round_id):
    """Set correct answers for a round"""
    logger.debug(f"[ROUND] set_answers() - round_id={round_id}")
    with db_connect() as conn:
        round_info = conn.execute("SELECT * FROM rounds WHERE id = ?", (round_id,)).fetchone()
        num_answers = round_info['num_answers']
        logger.debug(f"[ROUND] Setting {num_answers} answers for round {round_info['round_number']}")
        
        fields = []
        values = []
        for i in range(1, 7):
            if i <= num_answers:
                fields.append(f'answer{i} = ?')
                if i == 1:
                    fields.append(f'answer{i}_count = ?')
                    values.append(request.form.get(f'answer{i}', '').strip())
                    values.append(int(request.form.get(f'answer{i}_count', 0) or 0))
                else:
                    values.append(request.form.get(f'answer{i}', '').strip())
        
        values.append(round_id)
        conn.execute(f"UPDATE rounds SET {', '.join(fields)} WHERE id = ?", values)
        conn.commit()
    logger.info(f"[ROUND] set_answers() - answers saved for round_id={round_id}")
    return redirect(url_for('host_dashboard'))

@app.route('/host/scoring-queue')
@host_required
def scoring_queue():
    """Manual scoring page - shows unscored submissions"""
    logger.debug("[SCORING] scoring_queue() - loading scoring queue")
    with db_connect() as conn:
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()

        if not active_round:
            logger.warning("[SCORING] scoring_queue() - no active round")
            flash('No active round!', 'error'); return redirect(url_for('host_dashboard'))
        
        # Get unscored submissions
        submissions = conn.execute("""
            SELECT s.*, tc.team_name
            FROM submissions s
            JOIN team_codes tc ON s.code = tc.code
            WHERE s.round_id = ? AND s.scored = 0
            ORDER BY s.submitted_at ASC
        """, (active_round['id'],)).fetchall()
        
        submissions_data = []
        for sub in submissions:
            sub_dict = dict(sub)
            sub_dict['time_ago'] = time_ago(sub['submitted_at'])
            sub_dict['submitted_time'] = format_timestamp(sub['submitted_at'])
            
            # All boxes unchecked by default — host reviews manually
            auto_checks = {i: False for i in range(1, active_round['num_answers'] + 1)}
            
            sub_dict['auto_checks'] = auto_checks
            submissions_data.append(sub_dict)
    logger.debug(f"[SCORING] scoring_queue() - {len(submissions_data)} unscored submissions for round {active_round['round_number']}")
    ai_enabled = AI_SCORING_ENABLED and get_setting('ai_scoring_enabled', 'true') == 'true'
    return render_template('scoring_queue.html',
                         round=dict(active_round),
                         submissions=submissions_data,
                         ai_scoring_enabled=ai_enabled)

@app.route('/host/check-active-round')
@host_required
def check_active_round():
    """API endpoint to check if there's an active round (for AJAX polling)"""
    with db_connect() as conn:
        active_round = conn.execute("SELECT id FROM rounds WHERE is_active = 1").fetchone()
        has_active = active_round is not None
        logger.debug(f"[API] check_active_round() = {has_active}")
        return jsonify({'has_active_round': has_active})

@app.route('/host/count-unscored')
@host_required
def count_unscored():
    """API endpoint to get count of unscored submissions"""
    with db_connect() as conn:
        active_round = conn.execute("SELECT id FROM rounds WHERE is_active = 1").fetchone()
        
        if not active_round:
            return jsonify({'count': 0})
        
        count = conn.execute("""
            SELECT COUNT(*) as cnt FROM submissions
            WHERE round_id = ? AND scored = 0
        """, (active_round['id'],)).fetchone()['cnt']
        logger.debug(f"[API] count_unscored() = {count}")
        return jsonify({'count': count})

@app.route('/host/score-team/<int:submission_id>', methods=['POST'])
@host_required
def score_team(submission_id):
    """Submit score for a single team"""
    logger.debug(f"[SCORING] score_team() - submission_id={submission_id}")
    checked_answers = []
    for key in request.form:
        if key.startswith('answer_'):
            checked_answers.append(int(key.split('_')[1]))
    logger.debug(f"[SCORING] Checked answers: {sorted(checked_answers)}")

    with db_connect() as conn:
        submission = conn.execute("SELECT * FROM submissions WHERE id = ?", (submission_id,)).fetchone()
        round_info = conn.execute("SELECT * FROM rounds WHERE id = ?", (submission['round_id'],)).fetchone()

        # Get team name
        team_info = conn.execute("SELECT team_name FROM team_codes WHERE code = ?", (submission['code'],)).fetchone()
        team_name = team_info['team_name'] if team_info else 'Unknown Team'
        logger.debug(f"[SCORING] Team: {team_name} (code={submission['code']})")
        
        # Calculate score based on checked boxes
        score = 0
        for ans_num in checked_answers:
            points = round_info['num_answers'] - ans_num + 1
            score += points
        
        # Store which answers were checked (e.g., "1,3,5")
        checked_answers_str = ','.join(map(str, sorted(checked_answers))) if checked_answers else ''
        
        # Store current score as previous_score before updating (for undo functionality)
        current_score = submission['score']
        
        # Update submission with new score and save previous
        logger.debug(f"[SCORING] Score calculated: {score} points (checked: {checked_answers_str}), previous_score: {current_score}")
        logger.info(f"[SCORING] Scored: {team_name} ({submission['code']}) = {score}pts (answers: {checked_answers_str})")

        conn.execute("""
            UPDATE submissions
            SET score = ?, scored = 1, scored_at = CURRENT_TIMESTAMP, checked_answers = ?, previous_score = ?
            WHERE id = ?
        """, (score, checked_answers_str, current_score, submission_id))

        # === AI CORRECTIONS: Detect and store host overrides ===
        ai_matches_str = request.form.get('ai_matches', '').strip()
        ai_reasoning_str = request.form.get('ai_reasoning', '').strip()
        # Per-answer override notes: ai_note_1, ai_note_2, etc.
        ai_notes = {}
        for key, val in request.form.items():
            if key.startswith('ai_note_') and val.strip():
                try:
                    answer_num = int(key.split('_')[2])
                    ai_notes[answer_num] = val.strip()[:200]
                except (ValueError, IndexError):
                    pass

        if ai_matches_str:
            logger.debug(f"[AI-CORRECTIONS] Processing corrections for submission_id={submission_id}")
            ai_matches = set(int(x) for x in ai_matches_str.split(',') if x.strip())
            host_matches = set(checked_answers)

            # Parse AI reasoning for context
            ai_reasoning_list = []
            if ai_reasoning_str:
                try:
                    ai_reasoning_list = json.loads(ai_reasoning_str)
                except Exception:
                    logger.warning("[AI-CORRECTIONS] Failed to parse ai_reasoning JSON")

            host_added = host_matches - ai_matches    # Host checked, AI didn't
            host_removed = ai_matches - host_matches  # AI checked, host unchecked

            logger.debug(f"[AI-CORRECTIONS] AI={sorted(ai_matches)}, Host={sorted(host_matches)}, added={host_added}, removed={host_removed}")

            corrections_count = 0

            for survey_num in host_added:
                survey_answer = round_info[f'answer{survey_num}']
                # Find the team answer from reasoning that relates to this survey answer
                team_answer = None
                ai_reason = None
                for entry in ai_reasoning_list:
                    # Check unmatched entries — AI didn't match them, but host says they match this survey answer
                    if entry.get('matched_to') is None and entry.get('team_answer'):
                        team_answer = entry.get('team_answer', '')
                        ai_reason = entry.get('why', '')
                        break
                if not team_answer:
                    # Fallback: use any team answer from the submission
                    for j in range(1, round_info['num_answers'] + 1):
                        ans = submission[f'answer{j}']
                        if ans and ans.strip():
                            team_answer = ans.strip()
                            break
                if team_answer:
                    host_note = ai_notes.get(survey_num, None)
                    conn.execute("""
                        INSERT INTO ai_corrections (round_id, submission_id, question, team_answer, survey_answer, survey_num, correction_type, ai_reasoning, host_reason)
                        VALUES (?, ?, ?, ?, ?, ?, 'host_added', ?, ?)
                    """, (submission['round_id'], submission_id, round_info['question'], team_answer, survey_answer, survey_num, ai_reason, host_note))
                    save_correction_to_history({
                        'team_answer': team_answer, 'survey_answer': survey_answer,
                        'correction_type': 'host_added', 'ai_reasoning': ai_reason,
                        'host_reason': host_note, 'question': round_info['question']
                    })
                    corrections_count += 1

            for survey_num in host_removed:
                survey_answer = round_info[f'answer{survey_num}']
                # Find the team answer AI matched to this survey answer
                team_answer = None
                ai_reason = None
                for entry in ai_reasoning_list:
                    if entry.get('matched_to') == survey_num:
                        team_answer = entry.get('team_answer', '')
                        ai_reason = entry.get('why', '')
                        break
                if team_answer:
                    host_note = ai_notes.get(survey_num, None)
                    conn.execute("""
                        INSERT INTO ai_corrections (round_id, submission_id, question, team_answer, survey_answer, survey_num, correction_type, ai_reasoning, host_reason)
                        VALUES (?, ?, ?, ?, ?, ?, 'host_removed', ?, ?)
                    """, (submission['round_id'], submission_id, round_info['question'], team_answer, survey_answer, survey_num, ai_reason, host_note))
                    save_correction_to_history({
                        'team_answer': team_answer, 'survey_answer': survey_answer,
                        'correction_type': 'host_removed', 'ai_reasoning': ai_reason,
                        'host_reason': host_note, 'question': round_info['question']
                    })
                    corrections_count += 1

            if corrections_count > 0:
                logger.info(f"[AI-CORRECTIONS] Stored {corrections_count} correction(s) for submission_id={submission_id}")
                if ai_notes:
                    logger.info(f"[AI-CORRECTIONS] Host notes: {ai_notes}")

        conn.commit()

        # Check if all submissions for this round are scored
        total_subs = conn.execute("SELECT COUNT(*) as cnt FROM submissions WHERE round_id = ?", 
                                   (submission['round_id'],)).fetchone()['cnt']
        scored_subs = conn.execute("SELECT COUNT(*) as cnt FROM submissions WHERE round_id = ? AND scored = 1", 
                                    (submission['round_id'],)).fetchone()['cnt']
        
        logger.info(f"[SCORING] Round progress: {scored_subs}/{total_subs} teams scored")

        # If all scored, find winner and update round
        if total_subs > 0 and scored_subs == total_subs:
            logger.info("[SCORING] ALL TEAMS SCORED - determining winner")
            winner = conn.execute("""
                SELECT code, score FROM submissions 
                WHERE round_id = ? 
                ORDER BY score DESC, tiebreaker DESC 
                LIMIT 1
            """, (submission['round_id'],)).fetchone()
            
            if winner:
                conn.execute("UPDATE rounds SET winner_code = ? WHERE id = ?", 
                           (winner['code'], submission['round_id']))
                conn.commit()
                logger.info(f"[SCORING] WINNER: code={winner['code']}, score={winner['score']} for round_id={submission['round_id']}")
    
    # Check if AJAX request
    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        # Return JSON for AJAX
        return jsonify({
            'success': True,
            'score': score,
            'team_name': team_name
        })
    else:
        # Traditional form submit (fallback)
        flash(f'{team_name} scored {score} points!', 'success')
        return redirect(url_for('scoring_queue'))

@app.route('/host/ai-score/<int:submission_id>', methods=['POST'])
@host_required
def ai_score_submission(submission_id):
    """Use Claude AI to suggest scoring for a submission"""
    logger.debug(f"[AI-SCORING] ai_score_submission() - submission_id={submission_id}")

    if not AI_SCORING_ENABLED:
        logger.error("[AI-SCORING] AI scoring not enabled at server level")
        return jsonify({'error': 'AI scoring not enabled'}), 500

    if get_setting('ai_scoring_enabled', 'true') != 'true':
        logger.error("[AI-SCORING] AI scoring disabled in settings")
        return jsonify({'error': 'AI scoring is turned off in settings'}), 500

    try:
        with db_connect() as conn:
            submission = conn.execute(
                "SELECT * FROM submissions WHERE id = ?", (submission_id,)
            ).fetchone()

            if not submission:
                logger.error(f"[AI-SCORING] Submission {submission_id} not found")
                return jsonify({'error': 'Submission not found'}), 404

            round_info = conn.execute(
                "SELECT * FROM rounds WHERE id = ?", (submission['round_id'],)
            ).fetchone()

            if not round_info:
                logger.error(f"[AI-SCORING] Round {submission['round_id']} not found")
                return jsonify({'error': 'Round not found'}), 404

            # Build survey answers list
            survey_answers = []
            for i in range(1, round_info['num_answers'] + 1):
                answer = round_info[f'answer{i}']
                if answer:
                    survey_answers.append({
                        'number': i,
                        'text': answer,
                        'points': round_info['num_answers'] - i + 1
                    })

            # Build team answers list (only non-blank)
            team_answers = []
            for i in range(1, round_info['num_answers'] + 1):
                answer = submission[f'answer{i}']
                if answer and answer.strip():
                    team_answers.append(answer.strip())

            if not team_answers:
                logger.info("[AI-SCORING] No team answers to score")
                return jsonify({'success': True, 'matches': [], 'reasoning': []})

            logger.debug(f"[AI-SCORING] Scoring {len(team_answers)} team answers against {len(survey_answers)} survey answers")

            ai_result = score_with_ai(
                question=round_info['question'],
                survey_answers=survey_answers,
                team_answers=team_answers
            )

            logger.info(f"[AI-SCORING] Result: matches={ai_result['matches']}, reasoning_count={len(ai_result.get('reasoning', []))}")

            return jsonify({
                'success': True,
                'matches': ai_result['matches'],
                'reasoning': ai_result.get('reasoning', [])
            })

    except Exception as e:
        logger.error(f"[AI-SCORING] Error: {e}", exc_info=True)
        return jsonify({'error': f'AI scoring failed: {str(e)}'}), 500

@app.route('/host/undo-score/<int:submission_id>', methods=['POST'])
@host_required
def undo_score(submission_id):
    """Undo the last score for a submission"""
    logger.info(f"[SCORING] undo_score() - submission_id={submission_id}")
    with db_connect() as conn:
        submission = conn.execute("SELECT * FROM submissions WHERE id = ?", (submission_id,)).fetchone()

        if not submission:
            logger.warning(f"[SCORING] undo_score() - submission {submission_id} not found")
            return jsonify({"success": False, "message": "Submission not found"}), 404

        if submission['previous_score'] is None:
            logger.warning(f"[SCORING] undo_score() - no previous score for submission {submission_id}")
            return jsonify({"success": False, "message": "No previous score to restore"}), 400
        
        # Get team name
        team_info = conn.execute("SELECT team_name FROM team_codes WHERE code = ?", (submission['code'],)).fetchone()
        team_name = team_info['team_name'] if team_info else 'Unknown Team'
        
        # Restore previous score
        previous_score = submission['previous_score']
        conn.execute("""
            UPDATE submissions 
            SET score = ?, previous_score = NULL
            WHERE id = ?
        """, (previous_score, submission_id))
        conn.commit()
        
        logger.info(f"[SCORING] undo_score() - {team_name} reverted from {submission['score']} to {previous_score}")
        
        return jsonify({
            "success": True,
            "message": f"{team_name}'s score restored to {previous_score}",
            "new_score": previous_score
        })

@app.route('/host/round-summary')
@host_required
def round_summary():
    """Show round summary after all teams scored"""
    try:
        with db_connect() as conn:
            active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()

            if not active_round:
                logger.warning("[SCORING] No active round found")
                flash('No active round!', 'error'); return redirect(url_for('host_dashboard'))

            logger.debug(f"[SCORING] Round summary: R{active_round['round_number']}, Q='{active_round['question'][:50]}'")

            # Get all scored submissions for this round
            submissions = conn.execute("""
                SELECT s.*, tc.team_name
                FROM submissions s
                JOIN team_codes tc ON s.code = tc.code
                WHERE s.round_id = ? AND s.scored = 1
                ORDER BY s.score DESC,
                         ABS(COALESCE(s.tiebreaker, 0) - ?) ASC,
                         s.submitted_at ASC
            """, (active_round['id'], active_round['answer1_count'] or 0)).fetchall()

            if not submissions:
                logger.warning("[SCORING] No scored teams found")
                flash('No scored teams yet!', 'warning'); return redirect(url_for('scoring_queue'))

            logger.debug(f"[SCORING] {len(submissions)} scored teams")
            for i, sub in enumerate(submissions):
                logger.debug(f"[SCORING]   {i+1}. {sub['team_name']} ({sub['code']}) - Score: {sub['score']}, TB: {sub['tiebreaker']}")

            # Get winner (first in sorted list)
            winner = dict(submissions[0])

            # Check if there was a tie
            tied = False
            tiebreaker_info = None
            ultimate_tie = False

            if len(submissions) > 1:
                second = submissions[1]
                if winner['score'] == second['score']:
                    tied = True
                    actual_count = active_round['answer1_count'] or 0
                    winner_diff = abs((winner['tiebreaker'] or 0) - actual_count)
                    second_diff = abs((second['tiebreaker'] or 0) - actual_count)

                    logger.debug(f"[SCORING] Tie details: winner_tb={winner['tiebreaker']}, second_tb={second['tiebreaker']}, actual={actual_count}")

                    # Check if tiebreaker guesses are also the same
                    if winner_diff == second_diff:
                        ultimate_tie = True

                        # Calculate time difference
                        try:
                            winner_time = datetime.strptime(winner['submitted_at'], '%Y-%m-%d %H:%M:%S')
                            second_time = datetime.strptime(second['submitted_at'], '%Y-%m-%d %H:%M:%S')
                            time_diff_seconds = abs((winner_time - second_time).total_seconds())
                        except:
                            time_diff_seconds = 0

                        tiebreaker_info = {
                            'winner_guess': winner['tiebreaker'] or 0,
                            'second_guess': second['tiebreaker'] or 0,
                            'actual_count': actual_count,
                            'difference': (winner['tiebreaker'] or 0) - actual_count,
                            'tied_score': winner['score'],
                            'ultimate_tie': True,
                            'winner_time': format_timestamp(winner['submitted_at']),
                            'second_time': format_timestamp(second['submitted_at']),
                            'second_name': second['team_name'],
                            'time_diff_seconds': int(time_diff_seconds)
                        }
                        logger.info(f"[ROUND-SUMMARY] R{active_round['round_number']}: {len(submissions)} teams | ULTIMATE TIE: {winner['team_name']} wins by submission time ({winner['score']}pts)")
                    else:
                        # Regular tiebreaker (different guesses)
                        tiebreaker_info = {
                            'winner_guess': winner['tiebreaker'] or 0,
                            'actual_count': actual_count,
                            'difference': (winner['tiebreaker'] or 0) - actual_count,
                            'tied_score': winner['score'],
                            'ultimate_tie': False
                        }
                        logger.info(f"[ROUND-SUMMARY] R{active_round['round_number']}: {len(submissions)} teams | TIEBREAKER: {winner['team_name']} wins ({winner['score']}pts)")
                else:
                    logger.info(f"[ROUND-SUMMARY] R{active_round['round_number']}: {len(submissions)} teams | Winner: {winner['team_name']} ({winner['score']}pts)")
            else:
                logger.info(f"[ROUND-SUMMARY] R{active_round['round_number']}: 1 team | Winner: {winner['team_name']} ({winner['score']}pts)")
        return render_template('round_summary.html',
                             round=dict(active_round),
                             winner=winner,
                             tied=tied,
                             tiebreaker_info=tiebreaker_info,
                             total_teams=len(submissions))
    except Exception as e:
        logger.error(f"[SCORING] round_summary() failed: {type(e).__name__}: {e}", exc_info=True)
        flash(f'Error loading summary: {str(e)}', 'error'); return redirect(url_for('host_dashboard'))

@app.route('/host/start-next-round', methods=['POST'])
@host_required
def start_next_round():
    """Move to next round"""
    logger.debug("[ROUND] start_next_round() called")
    with db_connect() as conn:
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()

        if active_round:
            current_num = active_round['round_number']
            logger.info(f"[ROUND] Current active round: {current_num}, advancing to {current_num + 1}")
            # Deactivate current
            conn.execute("UPDATE rounds SET is_active = 0 WHERE id = ?", (active_round['id'],))

            # Activate next round
            next_round = conn.execute("""
                SELECT * FROM rounds WHERE round_number = ?
            """, (current_num + 1,)).fetchone()

            if next_round:
                conn.execute("UPDATE rounds SET is_active = 1 WHERE id = ?", (next_round['id'],))
                conn.commit()
                logger.info(f"[ROUND] Activated round {current_num + 1} (id={next_round['id']})")
            else:
                # No more rounds - game over
                conn.commit()
                logger.info(f"[ROUND] No round {current_num + 1} found - game complete!")
                flash('All rounds complete!', 'info'); return redirect(url_for('host_dashboard'))
        
        conn.commit()
    
    return redirect(url_for('host_dashboard'))

@app.route('/host/scored-teams')
@host_required
def scored_teams():
    """View all scored teams"""
    logger.debug("[SCORING] scored_teams() - loading scored teams list")
    with db_connect() as conn:
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()

        if not active_round:
            logger.warning("[SCORING] scored_teams() - no active round")
            flash('No active round!', 'error'); return redirect(url_for('host_dashboard'))
        
        submissions = conn.execute("""
            SELECT s.*, tc.team_name
            FROM submissions s
            JOIN team_codes tc ON s.code = tc.code
            WHERE s.round_id = ? AND s.scored = 1
            ORDER BY s.score DESC, 
                     ABS(COALESCE(s.tiebreaker, 0) - ?) ASC
        """, (active_round['id'], active_round['answer1_count'] or 0)).fetchall()
        
        # Add formatted timestamps
        submissions_data = []
        for sub in submissions:
            sub_dict = dict(sub)
            sub_dict['submitted_time'] = format_timestamp(sub['submitted_at'])
            submissions_data.append(sub_dict)
    
    logger.debug(f"[SCORING] scored_teams() - {len(submissions_data)} scored teams for round {active_round['round_number']}")
    return render_template('scored_teams.html',
                         round=dict(active_round),
                         submissions=submissions_data)

@app.route('/host/edit-score/<int:submission_id>')
@host_required
def edit_score(submission_id):
    """Edit an already-scored submission"""
    logger.debug(f"[SCORING] edit_score() - loading edit form for submission_id={submission_id}")
    with db_connect() as conn:
        submission = conn.execute("""
            SELECT s.*, tc.team_name
            FROM submissions s
            JOIN team_codes tc ON s.code = tc.code
            WHERE s.id = ?
        """, (submission_id,)).fetchone()
        
        round_info = conn.execute("SELECT * FROM rounds WHERE id = ?", (submission['round_id'],)).fetchone()
        
        # Use stored checked_answers if available
        checked_set = set()
        if submission['checked_answers']:
            # Parse "1,3,5" into set {1, 3, 5}
            checked_set = set(map(int, submission['checked_answers'].split(',')))
        else:
            checked_set = set()  # No auto-matching — start unchecked
        
        # Convert to dict for template
        auto_checks = {i: (i in checked_set) for i in range(1, round_info['num_answers'] + 1)}
    logger.debug(f"[SCORING] edit_score() - team={submission['team_name']}, current_score={submission['score']}, checked={checked_set}")
    return render_template('edit_score.html',
                         round=dict(round_info),
                         submission=dict(submission),
                         auto_checks=auto_checks)

@app.route('/host/update-score/<int:submission_id>', methods=['POST'])
@host_required
def update_score(submission_id):
    """Update score for edited submission"""
    logger.debug(f"[SCORING] update_score() - submission_id={submission_id}")
    checked_answers = []
    for key in request.form:
        if key.startswith('answer_'):
            checked_answers.append(int(key.split('_')[1]))
    
    # Get tiebreaker if provided
    tiebreaker = request.form.get('tiebreaker', type=int)
    
    with db_connect() as conn:
        submission = conn.execute("SELECT * FROM submissions WHERE id = ?", (submission_id,)).fetchone()
        round_info = conn.execute("SELECT * FROM rounds WHERE id = ?", (submission['round_id'],)).fetchone()
        
        # Store previous score before updating
        previous_score = submission['score'] if submission['score'] is not None else 0
        
        score = 0
        for ans_num in checked_answers:
            points = round_info['num_answers'] - ans_num + 1
            score += points
        
        # Store which answers were checked
        checked_answers_str = ','.join(map(str, sorted(checked_answers))) if checked_answers else ''
        
        # Update score, tiebreaker, checked_answers, and previous_score
        if tiebreaker is not None:
            conn.execute("UPDATE submissions SET score = ?, tiebreaker = ?, checked_answers = ?, previous_score = ? WHERE id = ?", 
                        (score, tiebreaker, checked_answers_str, previous_score, submission_id))
        else:
            conn.execute("UPDATE submissions SET score = ?, checked_answers = ?, previous_score = ? WHERE id = ?",
                        (score, checked_answers_str, previous_score, submission_id))
        conn.commit()
    logger.info(f"[SCORING] update_score() - old_score={previous_score}, new_score={score}, checked={checked_answers_str}, tiebreaker={tiebreaker}")
    return redirect(url_for('scored_teams'))

@app.route('/host/edit-submission/<int:submission_id>')
@host_required
def edit_submission(submission_id):
    """Edit a team's submitted answers (answer1-6 + tiebreaker) before scoring"""
    logger.info(f"[SCORING] edit_submission() - loading edit form for submission_id={submission_id}")
    with db_connect() as conn:
        submission = conn.execute("""
            SELECT s.*, tc.team_name
            FROM submissions s
            JOIN team_codes tc ON s.code = tc.code
            WHERE s.id = ?
        """, (submission_id,)).fetchone()

        if not submission:
            flash('Submission not found!', 'error')
            return redirect(url_for('scoring_queue'))

        round_info = conn.execute("SELECT * FROM rounds WHERE id = ?", (submission['round_id'],)).fetchone()

        if not round_info:
            flash('Round not found!', 'error')
            return redirect(url_for('scoring_queue'))

    return render_template('edit_submission.html',
                         round=dict(round_info),
                         submission=dict(submission))

@app.route('/host/update-submission/<int:submission_id>', methods=['POST'])
@host_required
def update_submission(submission_id):
    """Save edited team answers (answer1-6 + tiebreaker)"""
    logger.info(f"[SCORING] update_submission() - submission_id={submission_id}")

    with db_connect() as conn:
        submission = conn.execute("SELECT * FROM submissions WHERE id = ?", (submission_id,)).fetchone()

        if not submission:
            flash('Submission not found!', 'error')
            return redirect(url_for('scoring_queue'))

        round_info = conn.execute("SELECT * FROM rounds WHERE id = ?", (submission['round_id'],)).fetchone()
        num_answers = round_info['num_answers']

        # Collect edited answers
        updates = []
        values = []
        for i in range(1, num_answers + 1):
            answer_val = request.form.get(f'answer{i}', '').strip()
            updates.append(f'answer{i} = ?')
            values.append(answer_val)

        # Collect edited tiebreaker
        tiebreaker = request.form.get('tiebreaker', type=int)
        if tiebreaker is None or tiebreaker < 0 or tiebreaker > 100:
            tiebreaker = 0
        updates.append('tiebreaker = ?')
        values.append(tiebreaker)

        values.append(submission_id)

        conn.execute(
            f"UPDATE submissions SET {', '.join(updates)} WHERE id = ?",
            values
        )
        conn.commit()

    logger.info(f"[SCORING] update_submission() - answers updated for submission_id={submission_id}")

    # Return JSON for AJAX (inline edit) requests
    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        answers = {}
        for i in range(1, num_answers + 1):
            answers[f'answer{i}'] = request.form.get(f'answer{i}', '').strip()
        return jsonify(success=True, answers=answers, tiebreaker=tiebreaker)

    flash('Submission answers updated!', 'success')
    return redirect(url_for('scoring_queue'))

@app.route('/host/revert-score/<int:submission_id>')
@host_required
def revert_score(submission_id):
    """Revert score to previous value"""
    logger.info(f"[SCORING] revert_score() - submission_id={submission_id}")
    with db_connect() as conn:
        submission = conn.execute("SELECT previous_score FROM submissions WHERE id = ?", (submission_id,)).fetchone()

        if submission and submission['previous_score'] is not None:
            logger.info(f"[SCORING] revert_score() - reverting to previous_score={submission['previous_score']}")
            # Swap current and previous scores
            current_previous = submission['previous_score']
            conn.execute("UPDATE submissions SET score = ?, previous_score = score WHERE id = ?", 
                        (current_previous, submission_id))
            conn.commit()
    
    return redirect(url_for('scored_teams'))

@app.route('/host/manual-entry')
@host_required
def manual_entry():
    """Manual entry form for paper submissions"""
    logger.debug("[SCORING] manual_entry() - loading manual entry form")
    with db_connect() as conn:
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()

        if not active_round:
            flash('No active round! Please activate a round first.', 'error')
            return redirect(url_for('host_dashboard'))
        
        # Get ALL codes (both used and unused)
        all_codes = conn.execute("""
            SELECT code, team_name, used FROM team_codes 
            ORDER BY code ASC
        """).fetchall()
    
    return render_template('manual_entry.html',
                         round=dict(active_round),
                         codes=all_codes)

@app.route('/host/manual-entry/submit', methods=['POST'])
@host_required
def manual_entry_submit():
    """Process manual paper submission"""
    code = request.form.get('code')
    team_name = request.form.get('team_name', '').strip()
    round_id = request.form.get('round_id')
    tiebreaker = int(request.form.get('tiebreaker', 0) or 0)
    logger.info(f"[SCORING] manual_entry_submit() - code={code}, team_name='{team_name}', round_id={round_id}")
    
    if not code or not team_name:
        # Check if AJAX request
        if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
            return jsonify({'success': False, 'error': 'Please fill in all required fields!'}), 400
        else:
            flash('Please fill in all required fields!', 'error')
            return redirect(url_for('manual_entry'))
    
    with db_connect() as conn:
        # Mark code as used with team name
        conn.execute("UPDATE team_codes SET used = 1, team_name = ? WHERE code = ?", (team_name, code))
        
        # Get round info
        round_info = conn.execute("SELECT num_answers FROM rounds WHERE id = ?", (round_id,)).fetchone()
        num_answers = round_info['num_answers']
        
        # Collect answers
        answers = {f'answer{i}': request.form.get(f'answer{i}', '').strip() for i in range(1, num_answers + 1)}
        
        # Insert submission
        fields = ['code', 'round_id', 'tiebreaker'] + [f'answer{i}' for i in range(1, num_answers + 1)]
        placeholders = ', '.join(['?'] * len(fields))
        values = [code, round_id, tiebreaker] + [answers[f'answer{i}'] for i in range(1, num_answers + 1)]
        
        try:
            conn.execute(f"INSERT INTO submissions ({', '.join(fields)}) VALUES ({placeholders})", values)
            conn.commit()
            logger.info(f"[SCORING] manual_entry_submit() - submission created for team '{team_name}' (code={code})")
        except sqlite3.IntegrityError:
            logger.warning(f"[SCORING] manual_entry_submit() - duplicate submission for code={code}")
            # Check if AJAX request
            if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
                return jsonify({'success': False, 'error': 'This code has already submitted for this round!'}), 400
            else:
                flash('This code has already submitted for this round!', 'error')
                return redirect(url_for('manual_entry'))
    
    # Check if AJAX request
    if request.headers.get('X-Requested-With') == 'XMLHttpRequest':
        # Return JSON for AJAX
        return jsonify({
            'success': True,
            'team_name': team_name
        })
    else:
        # Traditional form submit (fallback)
        flash('✅ Manual entry submitted successfully!', 'success')
        return redirect(url_for('host_dashboard'))

@app.route('/host/photo-scan')
@app.route('/host/scan')
@host_required
def photo_scan():
    """Photo scan page — mobile camera UI for scanning paper answer sheets"""
    logger.debug("[PHOTO-SCAN] photo_scan() - loading photo scan page")

    if not AI_SCORING_ENABLED:
        flash('AI features are required for Photo Scan. Enable AI scoring in Settings.', 'error')
        return redirect(url_for('host_dashboard'))

    with db_connect() as conn:
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()

        if not active_round:
            flash('No active round! Please activate a round first.', 'error')
            return redirect(url_for('host_dashboard'))

    return render_template('photo_scan.html',
                         round=dict(active_round))


@app.route('/host/photo-scan/upload', methods=['POST'])
@host_required
def photo_scan_upload():
    """Receive photo, extract answers via Claude Vision, insert into submissions"""
    logger.info("[PHOTO-SCAN] photo_scan_upload() - processing image")

    if not AI_SCORING_ENABLED:
        return jsonify({'success': False, 'error': 'AI features not available'}), 503

    data = request.get_json()
    if not data or 'image' not in data:
        return jsonify({'success': False, 'error': 'No image provided'}), 400

    image_b64 = data['image']
    round_id = data.get('round_id')

    with db_connect() as conn:
        round_info = conn.execute("SELECT * FROM rounds WHERE id = ?", (round_id,)).fetchone()
        if not round_info:
            return jsonify({'success': False, 'error': 'Round not found'}), 404
        num_answers = round_info['num_answers']

        # Get valid codes for matching
        valid_codes = {row['code'].upper(): row['code'] for row in
                       conn.execute("SELECT code FROM team_codes").fetchall()}

        # Extract answers from photo
        try:
            teams = extract_answers_from_photo(image_b64)
        except Exception as e:
            logger.error(f"[PHOTO-SCAN] Extraction failed: {e}")
            return jsonify({'success': False, 'error': 'Failed to read photo. Try again with better lighting.'}), 500

        if not teams:
            return jsonify({'success': False, 'error': 'No teams found in photo. Make sure the answer sheet is clearly visible.'}), 400

        # Insert each team into submissions
        results = []
        for team in teams:
            code_raw = team.get('code', '').strip()
            team_name = team.get('team_name', '').strip()
            tiebreaker = team.get('tiebreaker', 0)
            answers = team.get('answers', [''] * 6)

            # Match code: exact first, fuzzy fallback
            code = valid_codes.get(code_raw.upper(), '')

            if not code and code_raw:
                # Fuzzy fallback — find closest code (3 of 4 letters must match)
                best_ratio = 0
                best_code = ''
                code_upper = code_raw.upper()
                for valid_upper, valid_original in valid_codes.items():
                    ratio = SequenceMatcher(None, code_upper, valid_upper).ratio()
                    if ratio > best_ratio and ratio >= 0.75:
                        best_ratio = ratio
                        best_code = valid_original
                if best_code:
                    code = best_code
                    logger.info(f"[PHOTO-SCAN] Fuzzy code match: '{code_raw}' → {code} (ratio={best_ratio:.2f})")

            if not code:
                results.append({
                    'team_name': team_name or '(blank)',
                    'code': code_raw,
                    'success': False,
                    'error': f'Code "{code_raw}" not found'
                })
                continue

            # Answer sheet is authoritative for team names
            existing = conn.execute("SELECT team_name, used FROM team_codes WHERE code = ?", (code,)).fetchone()
            old_name = existing['team_name'] if existing else None

            if team_name:
                # Sheet has a name — use it (first registration OR rename)
                pending_name_update = team_name
            elif old_name:
                # No name on sheet but code already registered — keep existing name
                team_name = old_name
                pending_name_update = None
                logger.info(f"[PHOTO-SCAN] No name on sheet for code={code}, keeping existing: '{team_name}'")
            else:
                # No name on sheet AND code not registered — assign placeholder
                suffix = ''.join(secrets.choice(string.digits) for _ in range(4))
                team_name = f"NO_NAME_{suffix}"
                pending_name_update = team_name
                logger.info(f"[PHOTO-SCAN] No name on sheet for unregistered code={code}, assigned: '{team_name}'")

            # Build and insert submission (same logic as manual_entry_submit)
            fields = ['code', 'round_id', 'tiebreaker'] + [f'answer{i}' for i in range(1, num_answers + 1)]
            placeholders = ', '.join(['?'] * len(fields))
            values = [code, round_id, tiebreaker] + [answers[i] if i < len(answers) else '' for i in range(num_answers)]

            try:
                conn.execute(f"INSERT INTO submissions ({', '.join(fields)}) VALUES ({placeholders})", values)
                # Only update team name after submission succeeds to avoid
                # corrupting the canonical name on duplicate/failed inserts
                if pending_name_update:
                    if old_name and old_name != pending_name_update:
                        logger.info(f"[PHOTO-SCAN] Team name changed: code={code} '{old_name}' -> '{pending_name_update}'")
                    conn.execute("UPDATE team_codes SET used = 1, team_name = ? WHERE code = ?",
                                (pending_name_update, code))
                result_entry = {
                    'team_name': team_name,
                    'code': code,
                    'success': True
                }
                if old_name and old_name != team_name:
                    result_entry['name_changed_from'] = old_name
                results.append(result_entry)
                logger.info(f"[PHOTO-SCAN] Submitted: team='{team_name}' code={code}")
            except sqlite3.IntegrityError:
                results.append({
                    'team_name': old_name or team_name,
                    'code': code,
                    'success': False,
                    'error': 'Already submitted for this round'
                })
                logger.warning(f"[PHOTO-SCAN] Duplicate: code={code}")

        conn.commit()

    succeeded = sum(1 for r in results if r['success'])
    failed = sum(1 for r in results if not r['success'])
    logger.info(f"[PHOTO-SCAN] Done: {succeeded} succeeded, {failed} failed")

    return jsonify({
        'success': True,
        'results': results,
        'summary': {
            'total': len(results),
            'succeeded': succeeded,
            'failed': failed
        }
    })


@app.route('/host/round/<int:round_id>/edit-answer/<int:answer_num>')
@host_required
def edit_single_answer(round_id, answer_num):
    """Edit a single answer"""
    logger.info(f"[ROUND] edit_single_answer() - round_id={round_id}, answer_num={answer_num}")
    with db_connect() as conn:
        round_info = conn.execute("SELECT * FROM rounds WHERE id = ?", (round_id,)).fetchone()

        if not round_info:
            logger.warning(f"[ROUND] edit_single_answer() - round_id={round_id} not found")
            flash('Round not found!', 'error'); return redirect(url_for('host_dashboard'))
        
        current_answer = round_info[f'answer{answer_num}']
        current_count = round_info['answer1_count'] if answer_num == 1 else None
    
    return render_template('edit_answer.html',
                         round=dict(round_info),
                         answer_num=answer_num,
                         current_answer=current_answer,
                         current_count=current_count)

@app.route('/host/round/<int:round_id>/update-answer/<int:answer_num>', methods=['POST'])
@host_required
def update_single_answer(round_id, answer_num):
    """Update a single answer"""
    new_answer = request.form.get('answer', '').strip()
    logger.info(f"[ROUND] update_single_answer() - round_id={round_id}, answer_num={answer_num}, new_answer='{new_answer}'")
    
    with db_connect() as conn:
        if answer_num == 1:
            new_count = int(request.form.get('count', 0) or 0)
            conn.execute("""
                UPDATE rounds 
                SET answer1 = ?, answer1_count = ?
                WHERE id = ?
            """, (new_answer, new_count, round_id))
        else:
            conn.execute(f"""
                UPDATE rounds 
                SET answer{answer_num} = ?
                WHERE id = ?
            """, (new_answer, round_id))
        
        conn.commit()
    
    flash('✅ Answer updated!', 'success'); return redirect(url_for('host_dashboard'))

@app.route('/host/create-round-manual')
@host_required
def create_round_manual_form():
    """Show manual round creation form"""
    return render_template('create_round_manual.html',
                         rounds_config=ROUNDS_CONFIG,
                         prebuilt_surveys=PREBUILT_SURVEYS)

@app.route('/host/create-round-manual/submit', methods=['POST'])
@host_required
def create_round_manual_submit():
    """Process manual round creation for ALL 8 rounds"""
    logger.info("[ROUND] create_round_manual_submit() - creating all 8 rounds manually")
    try:
        with db_connect() as conn:
            # Delete any existing rounds and submissions
            conn.execute("DELETE FROM rounds")
            conn.execute("DELETE FROM submissions")
            
            # Create all 8 rounds
            for config in ROUNDS_CONFIG:
                round_num = config['round']
                num_answers = config['answers']
                
                # Get question for this round
                question = request.form.get(f'question{round_num}', '').strip()
                
                # Build insert for this round
                fields = ['round_number', 'question', 'num_answers', 'is_active']
                is_active = 0
                values = [round_num, question, num_answers, is_active]
                
                # Get answers for this round
                for i in range(1, num_answers + 1):
                    answer = request.form.get(f'round{round_num}_answer{i}', '').strip()
                    fields.append(f'answer{i}')
                    values.append(answer)
                    
                    # Get count only for answer #1
                    if i == 1:
                        count = int(request.form.get(f'round{round_num}_answer1_count', 0) or 0)
                        fields.append(f'answer{i}_count')
                        values.append(count)
                
                # Insert this round
                placeholders = ','.join(['?'] * len(values))
                conn.execute(f"INSERT INTO rounds ({','.join(fields)}) VALUES ({placeholders})", values)
            
            conn.commit()
        
        logger.info("[ROUND] create_round_manual_submit() - all 8 rounds created successfully")
        flash('✅ All 8 rounds created!', 'success'); return redirect(url_for('host_dashboard'))

    except Exception as e:
        logger.error(f"[ROUND] create_round_manual_submit() error: {e}")
        flash(f'Error creating rounds: {str(e)}', 'error'); return redirect(url_for('host_dashboard'))


@app.route('/host/reset', methods=['POST'])
@host_required
def reset_game():
    """Reset game but keep codes and team names - for setup fixes"""
    logger.info("[HOST] reset_game() - resetting game (keeping teams)")
    with db_connect() as conn:
        conn.execute("DELETE FROM submissions")
        conn.execute("DELETE FROM rounds")
        # DO NOT touch team_codes table - keep teams joined!
        conn.commit()
    logger.info("[HOST] reset_game() - submissions and rounds deleted, team_codes untouched")
    flash('Game reset! Teams are still joined. Upload new questions to start fresh.', 'success')
    return redirect(url_for('host_dashboard'))

@app.route('/host/reset-all', methods=['POST'])
@host_required
def reset_all():
    """Reset everything - clear teams but keep code values"""
    global RESET_COUNTER
    
    with db_connect() as conn:
        conn.execute("DELETE FROM submissions")
        conn.execute("DELETE FROM rounds")
        # Reset codes to unused but keep the code values (HNCL, LZLX, etc)
        conn.execute("UPDATE team_codes SET used = 0, team_name = NULL")
        conn.commit()
    
    # Increment reset counter to invalidate all team sessions
    RESET_COUNTER += 1
    logger.info(f"[HOST] reset_all() - RESET_COUNTER incremented to {RESET_COUNTER}")
    logger.info("[HOST] All team sessions are now invalid - teams will see Game Over page")
    
    flash('Everything reset! All codes are now unused and ready for new teams.', 'success')
    return redirect(url_for('host_dashboard'))

@app.route('/host/settings', methods=['GET', 'POST'])
@host_required
def settings():
    """Settings page for configuring game options"""
    if request.method == 'POST':
        # Get form data
        qr_base_url = request.form.get('qr_base_url', '').strip()
        logger.info(f"[SETTINGS] settings() POST - qr_base_url='{qr_base_url}'")
        
        # Basic validation
        if not qr_base_url:
            flash('QR Base URL cannot be empty!', 'error')
        elif ' ' in qr_base_url:
            flash('URL cannot contain spaces!', 'error')
        else:
            # Save setting
            if set_setting('qr_base_url', qr_base_url, 'Base URL for QR codes on printed sheets'):
                flash('Settings saved successfully!', 'success')
            else:
                flash('Failed to save settings. Please try again.', 'error')
        
        return redirect(url_for('settings'))
    
   # GET - show form with current settings
    current_qr_url = get_qr_base_url()
    allow_team_registration = get_setting('allow_team_registration', 'true') == 'true'
    system_paused = get_setting('system_paused', 'false') == 'true'
    broadcast_message = get_setting('broadcast_message', '')
    ai_scoring_enabled = get_setting('ai_scoring_enabled', 'true') == 'true'
    extended_thinking_enabled = get_setting('extended_thinking_enabled', 'false') == 'true'
    thinking_budget_tokens = int(get_setting('thinking_budget_tokens', '10000'))

    # Count corrections in current session
    corrections_count = len(load_corrections_history())

    return render_template('settings.html',
                         qr_base_url=current_qr_url,
                         allow_team_registration=allow_team_registration,
                         system_paused=system_paused,
                         broadcast_message=broadcast_message,
                         ai_scoring_available=AI_SCORING_ENABLED,
                         ai_scoring_enabled=ai_scoring_enabled,
                         corrections_count=corrections_count,
                         ai_model_choices=AI_MODEL_CHOICES,
                         current_ai_model=get_current_ai_model(),
                         extended_thinking_enabled=extended_thinking_enabled,
                         thinking_budget_tokens=thinking_budget_tokens)

@app.route('/host/save-training', methods=['POST'])
@host_required
def save_training():
    """Save AI corrections to GitHub repo for long-term persistence."""
    if not GITHUB_TOKEN:
        return jsonify({'success': False, 'error': 'GITHUB_TOKEN not configured. Set it in Render environment variables.'}), 400

    corrections = load_corrections_history()
    if not corrections:
        return jsonify({'success': False, 'error': 'No corrections to save.'}), 400

    try:
        file_path = 'corrections_history.json'
        api_url = f'https://api.github.com/repos/{GITHUB_REPO}/contents/{file_path}'

        # First, get the current file SHA (needed for updates)
        get_req = urllib.request.Request(api_url, headers={
            'Authorization': f'token {GITHUB_TOKEN}',
            'Accept': 'application/vnd.github.v3+json'
        })

        existing_sha = None
        existing_data = []
        try:
            with urllib.request.urlopen(get_req) as resp:
                file_info = json.loads(resp.read().decode())
                existing_sha = file_info.get('sha')
                # Decode existing content and merge
                existing_content = base64.b64decode(file_info.get('content', '')).decode('utf-8')
                existing_data = json.loads(existing_content) if existing_content.strip() else []
        except urllib.error.HTTPError as e:
            if e.code == 404:
                existing_data = []  # File doesn't exist yet
            else:
                raise

        # Merge: add new corrections that aren't already in the file
        # Use a simple dedup by converting to comparable tuples
        existing_set = set()
        for c in existing_data:
            key = (c.get('team_answer', ''), c.get('survey_answer', ''), c.get('correction_type', ''), c.get('question', ''))
            existing_set.add(key)

        new_corrections = []
        for c in corrections:
            key = (c.get('team_answer', ''), c.get('survey_answer', ''), c.get('correction_type', ''), c.get('question', ''))
            if key not in existing_set:
                new_corrections.append(c)

        if not new_corrections:
            return jsonify({'success': True, 'message': f'All {len(corrections)} corrections already saved. No new data.'})

        merged = existing_data + new_corrections
        content_b64 = base64.b64encode(json.dumps(merged, indent=2).encode('utf-8')).decode('utf-8')

        # Commit to GitHub
        payload = json.dumps({
            'message': f'Update AI training data (+{len(new_corrections)} corrections, {len(merged)} total)',
            'content': content_b64,
            'sha': existing_sha  # None if new file
        }).encode('utf-8')

        # Remove sha key if None (new file)
        payload_dict = json.loads(payload)
        if payload_dict.get('sha') is None:
            del payload_dict['sha']
        payload = json.dumps(payload_dict).encode('utf-8')

        put_req = urllib.request.Request(api_url, data=payload, method='PUT', headers={
            'Authorization': f'token {GITHUB_TOKEN}',
            'Accept': 'application/vnd.github.v3+json',
            'Content-Type': 'application/json'
        })

        with urllib.request.urlopen(put_req) as resp:
            if resp.status in (200, 201):
                logger.info(f"[AI-CORRECTIONS] Saved {len(new_corrections)} new corrections to GitHub ({len(merged)} total)")
                return jsonify({'success': True, 'message': f'Saved {len(new_corrections)} new corrections to GitHub ({len(merged)} total)'})

        return jsonify({'success': False, 'error': 'Unexpected response from GitHub'}), 500

    except urllib.error.HTTPError as e:
        error_body = e.read().decode() if e.fp else str(e)
        logger.error(f"[AI-CORRECTIONS] GitHub API error: {e.code} - {error_body}")
        return jsonify({'success': False, 'error': f'GitHub API error ({e.code}). Check your token permissions.'}), 500
    except Exception as e:
        logger.error(f"[AI-CORRECTIONS] Failed to save to GitHub: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/host/clear-training', methods=['POST'])
@host_required
def clear_training():
    """Clear all AI training corrections from local file and database."""
    try:
        # Clear the local JSON file
        with open(CORRECTIONS_FILE, 'w') as f:
            json.dump([], f)
        logger.info("[AI-CORRECTIONS] Cleared corrections_history.json")

        # Clear the database table
        with db_connect() as conn:
            conn.execute("DELETE FROM ai_corrections")
            conn.commit()
        logger.info("[AI-CORRECTIONS] Cleared ai_corrections table")

        return jsonify({'success': True, 'message': 'All training data cleared.'})
    except Exception as e:
        logger.error(f"[AI-CORRECTIONS] Failed to clear training data: {e}")
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/host/toggle-setting', methods=['POST'])
@host_required
def toggle_setting():
    """Toggle a boolean setting"""
    setting_key = request.form.get('setting_key')

    if setting_key in ['allow_team_registration', 'system_paused', 'ai_scoring_enabled', 'extended_thinking_enabled']:
        current_value = get_setting(setting_key, 'true' if setting_key == 'ai_scoring_enabled' else 'false')
        new_value = 'false' if current_value == 'true' else 'true'
        logger.info(f"[SETTINGS] toggle_setting() - {setting_key}: '{current_value}' -> '{new_value}'")

        set_setting(setting_key, new_value, '')

        # User-friendly messages
        if setting_key == 'allow_team_registration':
            if new_value == 'true':
                flash('✅ Team registration enabled - New teams can join!', 'success')
            else:
                flash('🚫 Team registration disabled - No new teams can join', 'success')
        elif setting_key == 'system_paused':
            if new_value == 'true':
                # Auto-disable registration when pausing
                set_setting('allow_team_registration', 'false', '')
                flash('⏸️ System PAUSED - Team registration also disabled', 'success')
            else:
                flash('▶️ System RESUMED - Remember to re-enable registration if needed', 'success')
        elif setting_key == 'ai_scoring_enabled':
            if new_value == 'true':
                flash('🤖 AI Scoring enabled - AI button will appear on scoring queue', 'success')
            else:
                flash('🤖 AI Scoring disabled - AI button hidden from scoring queue', 'success')
        elif setting_key == 'extended_thinking_enabled':
            if new_value == 'true':
                flash('🧠 Extended Thinking enabled - AI will think deeper (higher cost)', 'success')
            else:
                flash('🧠 Extended Thinking disabled - Using standard mode', 'success')
    
    return redirect(url_for('settings'))

@app.route('/host/set-ai-model', methods=['POST'])
@host_required
def set_ai_model():
    """Set the AI model for scoring and photo scanning"""
    model_id = request.form.get('ai_model', '').strip()

    valid_ids = [m['id'] for m in AI_MODEL_CHOICES]
    if model_id not in valid_ids:
        flash('Invalid model selection.', 'error')
        return redirect(url_for('settings'))

    set_setting('ai_model', model_id, 'AI model for scoring and photo scan')

    model_name = next((m['name'] for m in AI_MODEL_CHOICES if m['id'] == model_id), model_id)
    logger.info(f"[SETTINGS] AI model changed to: {model_id}")
    flash(f'AI Model set to {model_name}', 'success')

    return redirect(url_for('settings'))

@app.route('/host/set-thinking-budget', methods=['POST'])
@host_required
def set_thinking_budget():
    """Set the token budget for extended thinking"""
    budget_str = request.form.get('thinking_budget', '').strip()

    try:
        budget = int(budget_str)
    except (ValueError, TypeError):
        flash('Invalid budget value. Must be a number.', 'error')
        return redirect(url_for('settings'))

    if budget < 1024:
        flash('Thinking budget must be at least 1,024 tokens.', 'error')
        return redirect(url_for('settings'))

    if budget > 128000:
        flash('Thinking budget cannot exceed 128,000 tokens.', 'error')
        return redirect(url_for('settings'))

    set_setting('thinking_budget_tokens', str(budget), 'Token budget for extended thinking')

    logger.info(f"[SETTINGS] Thinking budget changed to: {budget}")
    flash(f'Thinking budget set to {budget:,} tokens', 'success')

    return redirect(url_for('settings'))

@app.route('/host/toggle-sleep', methods=['POST'])
@host_required
def toggle_sleep():
    """Toggle server sleep mode"""
    current_value = get_setting('server_sleep', 'false')
    new_value = 'false' if current_value == 'true' else 'true'
    
    set_setting('server_sleep', new_value, 'Server sleep mode - stops auto-refresh')
    
    if new_value == 'true':
        logger.info("[SETTINGS] Server sleep mode ENABLED - team auto-refresh will stop")
        flash('💤 Server sleep mode enabled - All auto-refresh stopped', 'success')
    else:
        logger.info("[SETTINGS] Server sleep mode DISABLED - team auto-refresh resumed")
        flash('⏰ Server awake - Auto-refresh resumed', 'success')
    
    return jsonify({'success': True, 'sleep_mode': new_value})

@app.route('/host/get-sleep-status')
@host_required
def get_sleep_status():
    """Get current sleep mode status"""
    sleep_mode = get_setting('server_sleep', 'false')
    logger.debug(f"[API] get_sleep_status() -> {sleep_mode}")
    return jsonify({'sleep_mode': sleep_mode})

@app.route('/host/send-broadcast', methods=['POST'])
@host_required
def send_broadcast():
    """Send broadcast message to all teams"""
    import html
    import json
    import time

    message = request.form.get('message', '').strip()
    logger.info(f"[HOST] send_broadcast() - message='{message[:50]}' (len={len(message)})")

    if not message:
        flash('⚠️ Message cannot be empty!', 'error')
        return redirect(url_for('settings'))
    
    # Security: Length limit (200 chars max)
    if len(message) > 200:
        flash('⚠️ Message too long! Maximum 200 characters.', 'error')
        return redirect(url_for('settings'))
    
    # Security: HTML escape to prevent XSS
    message = html.escape(message)
    
    # Store message with timestamp
    broadcast_data = {
        'message': message,
        'timestamp': time.time()
    }
    
    set_setting('broadcast_message', json.dumps(broadcast_data), 'Broadcast message to all teams')
    flash(f'📢 Message sent to all teams!', 'success')
    
    return redirect(url_for('settings'))

@app.route('/host/clear-broadcast', methods=['POST'])
@host_required
def clear_broadcast():
    """Clear broadcast message"""
    logger.info("[HOST] clear_broadcast() - broadcast message cleared")
    import json
    # Set empty broadcast with current timestamp
    broadcast_data = {
        'message': '',
        'timestamp': 0
    }
    set_setting('broadcast_message', json.dumps(broadcast_data), 'Broadcast message to all teams')
    flash('🗑️ Broadcast message cleared', 'success')
    return redirect(url_for('settings'))

@app.route('/host/close-round', methods=['POST'])
@host_required
def close_round():
    """Close submissions for the active round and move to scoring"""
    logger.info("[ROUND] close_round() called")
    with db_connect() as conn:
        # Get active round
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()

        if not active_round:
            logger.warning("[ROUND] close_round() - no active round to close")
            flash('⚠️ No active round to close', 'error')
            return redirect(url_for('host_dashboard'))
        
        # Mark round as closed
        conn.execute("UPDATE rounds SET submissions_closed = 1 WHERE id = ?", (active_round['id'],))
        conn.commit()
        
        # Count submissions
        sub_count = conn.execute("SELECT COUNT(*) as cnt FROM submissions WHERE round_id = ?", 
                                (active_round['id'],)).fetchone()['cnt']
        
        # Check if all submissions are already scored
        unscored_count = conn.execute(
            "SELECT COUNT(*) as cnt FROM submissions WHERE round_id = ? AND scored = 0",
            (active_round['id'],)
        ).fetchone()['cnt']
        
        logger.info(f"[ROUND] Round {active_round['round_number']} closed - {sub_count} submissions, {unscored_count} unscored")
        flash(f'🔒 Round {active_round["round_number"]} closed! {sub_count} teams submitted.', 'success')

        # If all teams are already scored, skip scoring queue and go straight to winner announcement
        if unscored_count == 0 and sub_count > 0:
            logger.info(f"[ROUND] All {sub_count} teams already scored - redirecting to round_summary")
            return redirect(url_for('round_summary'))
        
        return redirect(url_for('scoring_queue'))

# ============= TEAM ROUTES =============

@app.route('/join')
def join():
    """Team join page - step 1"""
    paused = get_setting('system_paused', 'false') == 'true'
    reg_closed = get_setting('allow_team_registration', 'true') == 'false'
    logger.debug(f"[TEAM] join() page loaded | paused={paused}, registration_closed={reg_closed}")
    # Check if system is paused
    if paused:
        return render_template('join.html', error="⏸️ System is currently paused. Please wait for the host to resume.")

    # Check if team registration is allowed
    if reg_closed:
        return render_template('join.html', error="🚫 Team registration is currently closed.")

    # Support ?code= query param to pre-fill code from QR scan
    prefill_code = request.args.get('code', '').strip().upper()
    return render_template('join.html', prefill_code=prefill_code)

@app.route('/terms')
def terms():
    """Terms and conditions page"""
    return render_template('terms.html')

@app.route('/join/validate-code', methods=['POST'])
def validate_code():
    """Step 1: Validate team code"""
    logger.debug(f"[TEAM] validate_code() - code='{request.form.get('code', '').strip().upper()}'")
    # Check if system is paused
    if get_setting('system_paused', 'false') == 'true':
        return render_template('join.html', error="⏸️ System is currently paused. Please wait for the host to resume.")
    
    # Check if team registration is allowed
    if get_setting('allow_team_registration', 'true') == 'false':
        return render_template('join.html', error="🚫 Team registration is currently closed.")
    
    code = request.form.get('code', '').strip().upper()
    
    if not code:
        return render_template('join.html', error="Please enter a code")
    
    with db_connect() as conn:
        code_row = conn.execute("SELECT * FROM team_codes WHERE code = ?", (code,)).fetchone()
        
        if not code_row:
            logger.warning(f"[TEAM] validate_code() - code '{code}' not found in database")
            return render_template('join.html', error="Invalid code. Check your code and try again.")

        if code_row['used']:
            logger.info(f"[TEAM] validate_code() - code '{code}' already used by '{code_row['team_name']}', showing reconnect form")
            # Code is in use - show reconnection form
            return render_template('join.html', code=code, show_reconnect_form=True, existing_team=code_row['team_name'])

    logger.info(f"[TEAM] validate_code() - code '{code}' is valid and unused, showing team name form")
    return render_template('join.html', code=code, show_team_form=True)

def _rejoin_team(conn, code, code_row, source="REJOIN"):
    """Shared rejoin logic for both join_submit and join_reconnect routes.

    Handles: DB update (reconnected flag, heartbeat), session creation, redirect.
    Called AFTER validation has already confirmed the code is used and team name matches.

    Args:
        conn: Active database connection
        code: Team code (uppercase)
        code_row: Database row for the team code
        source: Log label ("REJOIN" or "RECONNECT") for distinguishing routes in logs

    Returns:
        Flask redirect to team_play
    """
    original_name = code_row['team_name']  # Always use DB capitalization

    logger.debug(f"[TEAM] {source}: Team '{original_name}' rejoining with code {code}")

    # Mark as reconnected + refresh heartbeat
    conn.execute(
        "UPDATE team_codes SET reconnected = 1, last_heartbeat = CURRENT_TIMESTAMP WHERE code = ?",
        (code,)
    )
    conn.commit()

    # Create session with current server state
    session['code'] = code
    session['team_name'] = original_name
    session['startup_id'] = STARTUP_ID
    session['reset_counter'] = RESET_COUNTER

    logger.info(f"[TEAM] {source}: Session created for '{original_name}' (Code: {code}), redirecting to team_play")
    return redirect(url_for('team_play'))


@app.route('/join/reconnect', methods=['POST'])
def join_reconnect():
    """Reconnect with existing team code"""
    logger.debug("[TEAM] RECONNECT: Attempt started")

    # Check if system is paused
    if get_setting('system_paused', 'false') == 'true':
        logger.info(f"[TEAM] RECONNECT: Blocked - system paused")
        return render_template('join.html', error="⏸️ System is currently paused. Please wait for the host to resume.")

    code = request.form.get('code', '').strip().upper()
    team_name = request.form.get('team_name', '').strip()

    logger.debug(f"[TEAM] RECONNECT: code='{code}', team_name='{team_name}'")

    if not code or not team_name:
        logger.warning(f"[TEAM] RECONNECT: Missing code or team_name")
        return render_template('join.html', error="Please enter both code and team name")

    with db_connect() as conn:
        code_row = conn.execute("SELECT * FROM team_codes WHERE code = ?", (code,)).fetchone()

        if not code_row:
            logger.warning(f"[TEAM] RECONNECT: Code '{code}' not found in database")
            return render_template('join.html', error="Invalid code")

        logger.debug(f"[TEAM] RECONNECT: Code found - used={code_row['used']}, team_name='{code_row['team_name']}'")

        if not code_row['used']:
            logger.warning(f"[TEAM] RECONNECT: Code '{code}' not yet used, rejecting reconnect")
            return render_template('join.html', error="This code hasn't been used yet. Use regular join.")

        # Case-insensitive team name comparison
        if code_row['team_name'].lower() != team_name.lower():
            logger.warning(f"[TEAM] RECONNECT: Name mismatch - DB='{code_row['team_name']}', submitted='{team_name}'")
            return render_template('join.html',
                code=code,
                show_reconnect_form=True,
                existing_team=code_row['team_name'],
                error="❌ Team name doesn't match. This code belongs to another team. Get a new code from the host.")

        # Validation passed - use shared rejoin logic
        return _rejoin_team(conn, code, code_row, source="RECONNECT")

@app.route('/join/submit', methods=['POST'])
def join_submit():
    """Step 2: Submit team name"""
    code_val = request.form.get('code', '').strip().upper()
    team_val = request.form.get('team_name', '').strip()
    logger.debug(f"[TEAM] join_submit() - code='{code_val}', team_name='{team_val}'")
    # Check if system is paused
    if get_setting('system_paused', 'false') == 'true':
        return render_template('join.html', error="⏸️ System is currently paused. Please wait for the host to resume.")
    
    # Check if team registration is allowed
    if get_setting('allow_team_registration', 'true') == 'false':
        return render_template('join.html', error="🚫 Team registration is currently closed.")
    
    code = request.form.get('code', '').strip().upper()
    team_name = request.form.get('team_name', '').strip()
    
    # Validation: Check for empty code or team name
    if not code or not team_name:
        return render_template('join.html', code=code, error="Please enter both code and team name")
    
    # Validation: Check for whitespace-only team name
    if len(team_name.strip()) == 0:
        return render_template('join.html', code=code, error="Team name cannot be empty or just spaces")
    
    # Validation: Team name character limit (30 chars max)
    if len(team_name) > 30:
        return render_template('join.html', code=code, error="Team name too long! Maximum 30 characters.")
    
    with db_connect() as conn:
        # Validation: Check for duplicate team names (case-insensitive)
        existing_team = conn.execute(
            "SELECT team_name FROM team_codes WHERE LOWER(team_name) = LOWER(?) AND used = 1 AND code != ?",
            (team_name, code)
        ).fetchone()
        if existing_team:
            logger.warning(f"[TEAM] join_submit() - team name '{team_name}' already taken")
            # Suggest alternative names
            base_name = team_name if len(team_name) <= 27 else team_name[:27]  # Leave room for " 2"
            counter = 2
            suggested_name = f"{base_name} {counter}"
            while conn.execute(
                "SELECT team_name FROM team_codes WHERE LOWER(team_name) = LOWER(?) AND used = 1 AND code != ?",
                (suggested_name, code)
            ).fetchone():
                counter += 1
                suggested_name = f"{base_name} {counter}"
            
            return render_template('join.html', code=code, error=f'Team name "{team_name}" already taken! Try: "{suggested_name}"')
        
        code_row = conn.execute("SELECT * FROM team_codes WHERE code = ?", (code,)).fetchone()
        
        if not code_row:
            return render_template('join.html', error="Invalid code")
        
        if code_row['used']:
            # Code is already used - check if it's the same team trying to rejoin
            if code_row['team_name'] and code_row['team_name'].lower() == team_name.lower():
                # Same team rejoining - use shared rejoin logic
                return _rejoin_team(conn, code, code_row, source="REJOIN")
            else:
                # Different team trying to use an already-used code
                logger.warning(f"[TEAM] REJOIN BLOCKED: Code {code} used by '{code_row['team_name']}', attempted by '{team_name}'")
                return render_template('join.html', error="Code already used by another team")
        
        # Code is unused - claim it
        conn.execute("UPDATE team_codes SET used = 1, team_name = ? WHERE code = ?", (team_name, code))
        conn.commit()
        logger.info(f"[TEAM] join_submit() - code '{code}' claimed by team '{team_name}', session created")

        # Store current startup_id and reset_counter in session
        # If server restarts or game resets, session becomes invalid
        session['code'] = code
        session['team_name'] = team_name
        session['startup_id'] = STARTUP_ID
        session['reset_counter'] = RESET_COUNTER
        
        return redirect(url_for('team_play'))

@app.route('/api/heartbeat', methods=['POST'])
@team_session_valid
def heartbeat():
    """Update last heartbeat timestamp for active tab detection"""
    code = session.get('code')
    logger.debug(f"[API] heartbeat() - code={code}")

    if not code:
        return jsonify({"success": False}), 401
    
    with db_connect() as conn:
        conn.execute("""
            UPDATE team_codes 
            SET last_heartbeat = CURRENT_TIMESTAMP 
            WHERE code = ?
        """, (code,))
        conn.commit()
    
    return jsonify({"success": True})

@app.route('/host/team-status')
@host_required
def get_team_status():
    """Get status of all teams (online/offline) for host dashboard"""
    logger.debug("[API] get_team_status() called")
    with db_connect() as conn:
        teams = conn.execute("""
            SELECT code, team_name, used, last_heartbeat,
                   CASE 
                       WHEN last_heartbeat IS NULL THEN 0
                       WHEN (julianday('now') - julianday(last_heartbeat)) * 86400 <= 15 THEN 1
                       ELSE 0
                   END as is_online
            FROM team_codes
            ORDER BY code
        """).fetchall()
        
        result = []
        for team in teams:
            team_dict = dict(team)
            # Calculate last seen time
            if team['last_heartbeat']:
                from datetime import datetime
                try:
                    last_seen = datetime.fromisoformat(team['last_heartbeat'].replace('Z', '+00:00'))
                    now = datetime.now(last_seen.tzinfo) if last_seen.tzinfo else datetime.now()
                    seconds_ago = int((now - last_seen).total_seconds())
                    
                    if seconds_ago < 60:
                        team_dict['last_seen_text'] = f"{seconds_ago} seconds ago"
                    elif seconds_ago < 3600:
                        minutes = seconds_ago // 60
                        team_dict['last_seen_text'] = f"{minutes} minute{'s' if minutes != 1 else ''} ago"
                    else:
                        hours = seconds_ago // 3600
                        team_dict['last_seen_text'] = f"{hours} hour{'s' if hours != 1 else ''} ago"
                except:
                    team_dict['last_seen_text'] = "Unknown"
            else:
                team_dict['last_seen_text'] = "Never"
            
            result.append(team_dict)

        online_count = sum(1 for t in result if t.get('is_online'))
        logger.debug(f"[API] get_team_status() -> {len(result)} teams, {online_count} online")
        return jsonify(result)

@app.route('/api/check-round-status')
def check_round_status():
    """API endpoint to check if there's an active round (for AJAX polling)"""
    code = session.get('code')
    logger.debug(f"[API] check_round_status() - code={code}")

    if not code:
        return jsonify({'error': 'No code in session'}), 401
    
    # Check if server was restarted (startup_id mismatch)
    session_startup_id = session.get('startup_id')
    if session_startup_id != STARTUP_ID:
        return jsonify({'error': 'Server restarted', 'reload': True}), 401
    
    # Check if game was reset (reset_counter mismatch)
    session_reset_counter = session.get('reset_counter', 0)
    if session_reset_counter != RESET_COUNTER:
        return jsonify({'error': 'Game was reset', 'reload': True}), 401
    
    # Check if server is in sleep mode
    server_sleep = get_setting('server_sleep', 'false')
    if server_sleep == 'true':
        return jsonify({'sleep_mode': True, 'message': 'Server in sleep mode'}), 200
    
    with db_connect() as conn:
        # Check if there's an active round
        active_round = conn.execute("SELECT id, round_number, submissions_closed FROM rounds WHERE is_active = 1").fetchone()
        
        if active_round:
            # Check if this team already submitted for this round
            submission = conn.execute(
                "SELECT id FROM submissions WHERE code = ? AND round_id = ?",
                (code, active_round['id'])
            ).fetchone()

            result = {
                'has_active_round': True,
                'round_id': active_round['id'],
                'round_number': active_round['round_number'],
                'submissions_closed': bool(active_round['submissions_closed']),
                'already_submitted': submission is not None
            }

            # Include previous round's winner (for winner interstitial on round transition)
            prev_round = conn.execute("""
                SELECT r.round_number, r.winner_code, tc.team_name, s.score
                FROM rounds r
                LEFT JOIN team_codes tc ON r.winner_code = tc.code
                LEFT JOIN submissions s ON r.winner_code = s.code AND r.id = s.round_id
                WHERE r.round_number = ? - 1
            """, (active_round['round_number'],)).fetchone()

            if prev_round and prev_round['winner_code']:
                result['prev_winner_team'] = prev_round['team_name']
                result['prev_winner_score'] = prev_round['score']
                result['prev_round_number'] = prev_round['round_number']

            logger.debug(f"[API] check_round_status() -> round={active_round['round_number']}, closed={bool(active_round['submissions_closed'])}, submitted={submission is not None}")
            return jsonify(result)
        else:
            logger.debug("[API] check_round_status() -> no active round")
            return jsonify({
                'has_active_round': False
            })

@app.route('/view/<code>')
def team_view(code):
    """View-only page for manually-entered teams. Auth is the code in the URL."""
    code = code.strip().upper()
    logger.debug(f"[VIEW] team_view() - code={code}")

    with db_connect() as conn:
        team = conn.execute(
            "SELECT code, team_name, used FROM team_codes WHERE code = ?",
            (code,)
        ).fetchone()

        if not team:
            logger.warning(f"[VIEW] team_view() - code not found: {code}")
            return render_template('view.html',
                team_name=f"Code: {code}",
                code=code,
                state='code_not_found',
                round_num=0,
                question='')

        if not team['used'] or not team['team_name']:
            logger.debug(f"[VIEW] team_view() - code exists but not yet registered: {code}")
            return render_template('view.html',
                team_name=f"Code: {code}",
                code=code,
                state='waiting_for_registration',
                round_num=0,
                question='')

        team_name = team['team_name']

        active_round = conn.execute(
            "SELECT * FROM rounds WHERE is_active = 1"
        ).fetchone()

        if not active_round:
            return render_template('view.html',
                team_name=team_name,
                code=code,
                state='waiting_for_round',
                round_num=0,
                question='')

        submission = conn.execute(
            "SELECT * FROM submissions WHERE code = ? AND round_id = ?",
            (code, active_round['id'])
        ).fetchone()

        if not submission:
            return render_template('view.html',
                team_name=team_name,
                code=code,
                state='waiting_for_entry',
                round_num=active_round['round_number'],
                question=active_round['question'])

        if not submission['scored']:
            return render_template('view.html',
                team_name=team_name,
                code=code,
                state='waiting_for_scoring',
                round_num=active_round['round_number'],
                question=active_round['question'],
                num_answers=active_round['num_answers'],
                submission=dict(submission))

        return render_template('view.html',
            team_name=team_name,
            code=code,
            state='scored',
            round_num=active_round['round_number'],
            question=active_round['question'],
            num_answers=active_round['num_answers'],
            submission=dict(submission))

@app.route('/play')
@team_session_valid
def team_play():
    """Team answer submission page"""
    code = session.get('code')
    team_name = session.get('team_name')
    logger.debug(f"[TEAM] team_play() - code={code}, team={team_name}")

    if not code:
        logger.warning("[TEAM] team_play() - no code in session, redirecting to join")
        return redirect(url_for('join'))
    
    with db_connect() as conn:
        # DEFENSIVE: Verify team still exists in database
        team = conn.execute("SELECT * FROM team_codes WHERE code = ?", (code,)).fetchone()
        
        if not team:
            # Team doesn't exist anymore - session is stale
            logger.error(f"[TEAM] team_play() - team {code} not found in database, clearing session")
            session.clear()
            return redirect(url_for('join'))
        
        # DEFENSIVE: Initialize last_heartbeat if NULL (for rejoining teams)
        if team['last_heartbeat'] is None:
            logger.debug(f"[TEAM] team_play() - initializing heartbeat for team {code} ({team_name})")
            conn.execute(
                "UPDATE team_codes SET last_heartbeat = CURRENT_TIMESTAMP WHERE code = ?",
                (code,)
            )
            conn.commit()
        
        active_round = conn.execute("SELECT * FROM rounds WHERE is_active = 1").fetchone()
        
        if not active_round:
            logger.debug(f"[TEAM] team_play() - no active round, showing waiting screen")
            return render_template('play.html',
                                 team_name=team_name,
                                 code=code,
                                 no_active_round=True)
        
        submission = conn.execute("""
            SELECT * FROM submissions 
            WHERE code = ? AND round_id = ?
        """, (code, active_round['id'])).fetchone()
        
        if submission:
            logger.debug(f"[TEAM] team_play() - round {active_round['round_number']}, already submitted")
            # Get last_submission from session (for answer preview)
            last_submission = session.pop('last_submission', None)
            
            return render_template('play.html',
                                 team_name=team_name,
                                 code=code,
                                 round_num=active_round['round_number'],
                                 question=active_round['question'],
                                 num_answers=active_round['num_answers'],
                                 already_submitted=True,
                                 submissions_closed=active_round['submissions_closed'],
                                 submission=dict(submission),
                                 last_submission=last_submission)
    
    logger.debug(f"[TEAM] team_play() - round {active_round['round_number']}, showing answer form ({active_round['num_answers']} answers)")
    return render_template('play.html',
                         team_name=team_name,
                         code=code,
                         round_num=active_round['round_number'],
                         question=active_round['question'],
                         num_answers=active_round['num_answers'],
                         round_id=active_round['id'],
                         submissions_closed=active_round['submissions_closed'])

@app.route('/play/submit', methods=['POST'])
@team_session_valid
def submit_answers():
    """Submit team answers"""
    code = session.get('code')
    round_id = request.form.get('round_id')
    logger.info(f"[TEAM] submit_answers() - code={code}, round_id={round_id}")

    # Check if system is paused
    if get_setting('system_paused', 'false') == 'true':
        logger.warning(f"[TEAM] submit_answers() - system paused, rejecting from code={code}")
        flash('⏸️ System is currently paused. Submissions are disabled.', 'error')
        return redirect(url_for('team_play'))

    if not code:
        return redirect(url_for('join'))
    
    # Validation: Tiebreaker must be 0-100
    try:
        tiebreaker = int(request.form.get('tiebreaker', 0) or 0)
        if tiebreaker < 0 or tiebreaker > 100:
            flash('⚠️ Tiebreaker must be between 0 and 100', 'error')
            return redirect(url_for('team_play'))
    except ValueError:
        tiebreaker = 0
    
    with db_connect() as conn:
        # Validate that this is still the active round (prevent stale submissions)
        active_round = conn.execute("SELECT id, submissions_closed FROM rounds WHERE is_active = 1").fetchone()
        if not active_round or str(active_round['id']) != str(round_id):
            logger.warning(f"[TEAM] submit_answers() - stale round_id={round_id}, active={active_round['id'] if active_round else 'None'}")
            # Round has changed - redirect to play page to get current round
            return redirect(url_for('team_play'))

        # Check if round is closed
        if active_round['submissions_closed']:
            logger.warning(f"[TEAM] submit_answers() - round closed, rejecting submission from code={code}")
            flash('⏰ Round has ended. Submissions are closed.', 'error')
            return redirect(url_for('team_play'))
        
        # CRITICAL FIX: Check for duplicate submission BEFORE attempting insert
        existing_submission = conn.execute(
            "SELECT id FROM submissions WHERE code = ? AND round_id = ?",
            (code, round_id)
        ).fetchone()
        
        if existing_submission:
            logger.warning(f"[TEAM] submit_answers() - duplicate submission from code={code} for round_id={round_id}")
            flash('✅ You have already submitted for this round!', 'warning')
            return redirect(url_for('team_play'))
        
        round_info = conn.execute("SELECT num_answers FROM rounds WHERE id = ?", (round_id,)).fetchone()
        num_answers = round_info['num_answers']
        
        answers = {f'answer{i}': request.form.get(f'answer{i}', '').strip() for i in range(1, num_answers + 1)}
        
        try:
            fields = ['code', 'round_id', 'tiebreaker'] + [f'answer{i}' for i in range(1, num_answers + 1)]
            placeholders = ', '.join(['?'] * len(fields))
            values = [code, round_id, tiebreaker] + [answers[f'answer{i}'] for i in range(1, num_answers + 1)]
            
            conn.execute(f"INSERT INTO submissions ({', '.join(fields)}) VALUES ({placeholders})", values)
            conn.commit()
            logger.info(f"[TEAM] submit_answers() - submission saved for code={code}, round_id={round_id}, tiebreaker={tiebreaker}, answers={answers}")

            # Store submission for answer preview
            session['last_submission'] = {
                'round_id': round_id,
                'answers': answers,
                'tiebreaker': tiebreaker
            }
        except sqlite3.IntegrityError:
            # Fallback: UNIQUE constraint caught it
            logger.warning(f"[TEAM] submit_answers() - UNIQUE constraint caught duplicate from code={code}")
            flash('✅ You have already submitted for this round!', 'warning')
    
    return redirect(url_for('team_play'))

@app.route('/api/broadcast-message')
def api_broadcast_message():
    """API endpoint for teams to get current broadcast message"""
    logger.debug("[API] api_broadcast_message() called")
    import json
    
    broadcast_json = get_setting('broadcast_message', '')
    
    # Handle legacy format (plain string) or new format (JSON)
    try:
        if broadcast_json:
            broadcast_data = json.loads(broadcast_json)
            return jsonify({
                'message': broadcast_data.get('message', ''),
                'timestamp': broadcast_data.get('timestamp', 0)
            })
        else:
            return jsonify({'message': '', 'timestamp': 0})
    except (json.JSONDecodeError, TypeError):
        # Legacy format - just a plain string
        return jsonify({'message': broadcast_json, 'timestamp': 0})

@app.route('/api/view-status/<code>')
def api_view_status(code):
    """API endpoint for view-only page polling. Returns round + scoring state."""
    code = code.strip().upper()
    logger.debug(f"[API] api_view_status() - code={code}")

    server_sleep = get_setting('server_sleep', 'false')
    if server_sleep == 'true':
        return jsonify({'sleep_mode': True}), 200

    with db_connect() as conn:
        team = conn.execute(
            "SELECT code, team_name, used FROM team_codes WHERE code = ?",
            (code,)
        ).fetchone()

        if not team:
            return jsonify({'state': 'code_not_found', 'has_active_round': False})

        if not team['used'] or not team['team_name']:
            return jsonify({'state': 'waiting_for_registration', 'has_active_round': False})

        active_round = conn.execute(
            "SELECT * FROM rounds WHERE is_active = 1"
        ).fetchone()

        if not active_round:
            result = {
                'has_active_round': False,
                'state': 'waiting_for_round'
            }

            last_round = conn.execute("""
                SELECT r.round_number, r.winner_code, tc.team_name, s.score
                FROM rounds r
                LEFT JOIN team_codes tc ON r.winner_code = tc.code
                LEFT JOIN submissions s ON r.winner_code = s.code AND r.id = s.round_id
                ORDER BY r.round_number DESC LIMIT 1
            """).fetchone()

            if last_round and last_round['winner_code']:
                result['prev_winner_team'] = last_round['team_name']
                result['prev_winner_score'] = last_round['score']
                result['prev_round_number'] = last_round['round_number']

            return jsonify(result)

        submission = conn.execute(
            "SELECT * FROM submissions WHERE code = ? AND round_id = ?",
            (code, active_round['id'])
        ).fetchone()

        result = {
            'has_active_round': True,
            'round_id': active_round['id'],
            'round_number': active_round['round_number'],
            'question': active_round['question'],
            'num_answers': active_round['num_answers'],
            'submissions_closed': bool(active_round['submissions_closed'])
        }

        if not submission:
            result['state'] = 'waiting_for_entry'
        elif not submission['scored']:
            result['state'] = 'waiting_for_scoring'
            result['answers'] = {
                f'answer{i}': submission[f'answer{i}']
                for i in range(1, active_round['num_answers'] + 1)
            }
            result['tiebreaker'] = submission['tiebreaker']
        else:
            # Players only see their submitted answers - no scoring data
            result['state'] = 'scored'
            result['answers'] = {
                f'answer{i}': submission[f'answer{i}']
                for i in range(1, active_round['num_answers'] + 1)
            }
            result['tiebreaker'] = submission['tiebreaker']

        prev_round = conn.execute("""
            SELECT r.round_number, r.winner_code, tc.team_name, s.score
            FROM rounds r
            LEFT JOIN team_codes tc ON r.winner_code = tc.code
            LEFT JOIN submissions s ON r.winner_code = s.code AND r.id = s.round_id
            WHERE r.round_number = ? - 1
        """, (active_round['round_number'],)).fetchone()

        if prev_round and prev_round['winner_code']:
            result['prev_winner_team'] = prev_round['team_name']
            result['prev_winner_score'] = prev_round['score']
            result['prev_round_number'] = prev_round['round_number']

        return jsonify(result)

if __name__ == '__main__':
    import socket
    hostname = socket.gethostname()
    local_ip = socket.gethostbyname(hostname)
    
    print("\n" + "="*60)
    print("🎮 FAMILY FEUD - PRODUCTION SERVER")
    print("="*60)
    print(f"\n📱 Team Join: http://{local_ip}:5000/join")
    print(f"🖥️  Host Dashboard: http://localhost:5000/host")
    print(f"🏆 Scoring Queue: http://localhost:5000/host/scoring-queue")
    print(f"\n💡 Upload answer sheet, generate codes, start playing!")
    print("="*60 + "\n")
    
    app.run(host='0.0.0.0', port=5000, debug=False)
