"""
File parsing utilities for Survey Says.

Parses PowerPoint (.pptx/.pptm) and Word (.docx) answer files into
a common rounds_data format for database insertion.
"""

from config import logger


def parse_pptx(filepath):
    """Parse PowerPoint file and extract questions/answers

    IMPROVED VERSION - Handles text boxes with answer/count pairs
    Correctly distinguishes rank indicators (1,2,3) from answer counts (10,20,43)
    """
    from pptx import Presentation

    prs = Presentation(filepath)
    slides = list(prs.slides)

    rounds_data = []

    # Strategy: Find question slides, then parse the next slide as answers
    i = 0
    while i < len(slides):
        slide = slides[i]

        # Extract all text from current slide
        all_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    all_text.append(text)

        # Check if this is a question slide
        is_question_slide = False
        question_text = ""

        for text in all_text:
            # Look for question markers
            if 'Survey Has' in text and 'Responses' in text:
                is_question_slide = True
            # Extract actual question (not the metadata)
            elif len(text) > 10 and 'Round #' not in text and 'Survey Has' not in text:
                if not question_text:  # Take first substantial text
                    question_text = text

        if is_question_slide and i + 1 < len(slides):
            # Next slide should have answers
            answer_slide = slides[i + 1]
            answers = []

            # Extract all text from answer slide
            answer_text_elements = []
            for shape in answer_slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        answer_text_elements.append(text)

            # Parse answer/count pairs from text
            j = 0
            while j < len(answer_text_elements):
                text = answer_text_elements[j]

                # Skip UI elements
                skip_keywords = ['Round:', 'ROUND', 'Score Multiplier:', 'BACK TO SCORES',
                               'NEXT ROUND', 'And The Survey Says', 'X', '\u00ab', '\u00bb',
                               'type only', '(type', 'Click', 'Press']

                if any(keyword in text for keyword in skip_keywords):
                    j += 1
                    continue

                # Skip rank indicators (single-digit numbers 1-8)
                if text.isdigit() and len(text) <= 2 and int(text) <= 8:
                    j += 1
                    continue

                # If it's text (potential answer), look for count
                if not text.isdigit() and len(text) > 1:
                    answer_text = text
                    count = 0

                    # Look ahead for count
                    if j + 1 < len(answer_text_elements):
                        next_text = answer_text_elements[j + 1]
                        if next_text.isdigit():
                            try:
                                count_value = int(next_text)
                                # Count numbers are typically 5+ (answer counts, not ranks)
                                # But also accept low counts (some answers might have count of 1-4)
                                # The key is: if we just read answer text, next number IS the count
                                if count_value > 0:
                                    answers.append({'answer': answer_text, 'count': count_value})
                                    j += 1  # Skip the count
                            except ValueError:
                                pass

                j += 1

            # Add round if we found answers
            if answers:
                rounds_data.append({
                    'question': question_text,
                    'answers': answers
                })

            # Skip the answer slide
            i += 2
        else:
            i += 1

    return rounds_data


def parse_docx(filepath):
    """Parse Word document and extract questions/answers.

    Extracts questions from numbered paragraphs (with dash separators)
    and answers from up to 8 tables. Handles both '-' and '\u2013' (em-dash) separators.
    """
    from docx import Document

    doc = Document(filepath)
    rounds_data = []

    # ROBUST: Extract questions with flexible matching (handles both - and \u2013 dashes)
    questions = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text and len(text) > 0 and text[0].isdigit():
            # Match both regular dash (-) and em-dash (\u2013)
            if '-' in text or '\u2013' in text:
                # Split on either dash type
                separator = '\u2013' if '\u2013' in text else '-'
                parts = text.split(separator, 1)
                if len(parts) > 1:
                    question = parts[1].strip()
                    questions.append(question)

    # ROBUST: Parse up to 12 tables regardless of question count
    for table_idx, table in enumerate(doc.tables):
        if table_idx >= 12:
            break

        answers = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if len(cells) >= 3:
                # Skip header rows - only process if first cell is a number (rank)
                if not cells[0] or not cells[0].strip().isdigit():
                    continue

                answer = cells[1]
                points_count = cells[2]

                # ROBUST: Flexible count parsing (handles various spacing)
                count = 0
                if points_count:
                    # Try both dash types
                    for separator in ['-', '\u2013']:
                        if separator in points_count:
                            parts = points_count.split(separator)
                            if len(parts) > 1:
                                try:
                                    # Extract just the digits from the second part
                                    count_str = ''.join(filter(str.isdigit, parts[1]))
                                    if count_str:
                                        count = int(count_str)
                                    break
                                except ValueError:
                                    count = 0

                answers.append({'answer': answer, 'count': count})

        # Use question by index, or empty string if not found
        question = questions[table_idx] if table_idx < len(questions) else ''

        rounds_data.append({
            'question': question,
            'answers': answers
        })

    return rounds_data
